
# 상품 식별 원칙:
# - 동일 코드 또는 동일 모델/본품/용량·수량 fingerprint -> 동일 product master로 과거 이력 연결 가능
# - 단, 증정품/추가구성 차이는 _product_variant_key로 별도 보존
# - 성과 집계/재편성 추천에서 variant가 실질적으로 다른 판매구성이면 별도 행 유지

# VERIFIED BASE: app_v4_2_8_gender_target_filter.py + promotion columns
# VERIFIED BUILD: V4.2.8-20260719-GENDER-TARGET-FILTER

from __future__ import annotations

import io
import math
import os
from pathlib import Path
import re
from datetime import datetime
from typing import Iterable

import pandas as pd
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import requests
import urllib3
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt


# === V4.4.38 DAILY INSIGHT ENGINE POLICY ===
# 성과등급 기준(주문금액):
# 핵심 상품: 5,000,000원 이상
# 우수 상품: 3,000,000원 이상 ~ 5,000,000원 미만
# 안정 상품: 2,000,000원 이상 ~ 3,000,000원 미만
# 관찰 상품: 1,000,000원 이상 ~ 2,000,000원 미만
# 부진 상품: 1,000,000원 미만
DAILY_GRADE_RULES = [
    (5_000_000, "핵심 상품"),
    (3_000_000, "우수 상품"),
    (2_000_000, "안정 상품"),
    (1_000_000, "관찰 상품"),
    (0, "부진 상품"),
]

def _daily_grade(amount):
    try:
        v = float(amount or 0)
    except Exception:
        v = 0.0
    for threshold, label in DAILY_GRADE_RULES:
        if v >= threshold:
            return label
    return "부진 상품"

def _daily_trend_label(values):
    """최근 최대 3회 추세: 상승/둔화만 명확할 때 반환, 혼조는 None."""
    vals = [float(x) for x in values if x is not None]
    if len(vals) < 3:
        return None
    vals = vals[-3:]
    if vals[0] < vals[1] < vals[2]:
        return "성과 성장"
    if vals[0] > vals[1] > vals[2]:
        return "성과 둔화"
    return None

def _daily_select_insights(candidates, max_core=3):
    """
    후보 인사이트를 중요도순으로 최대 3개 + 다음 운영 제안 1개로 압축.
    신규 첫 운영/신규 타겟 TEST는 최우선.
    금번 성과는 최대 1개.
    동일 근거의 성과/추세 중복을 제거.
    """
    if not candidates:
        return candidates

    rows = [str(x).strip() for x in candidates if str(x).strip()]
    actions = [x for x in rows if "다음 운영 제안" in x]
    warnings = [x for x in rows if x.startswith("주의:") or x.startswith("**주의")]
    rows = [x for x in rows if x not in actions and x not in warnings]

    # 완전 중복 제거
    uniq = []
    for x in rows:
        if x not in uniq:
            uniq.append(x)
    rows = uniq

    # 신규 첫 운영이 있으면 별도 금번 성과는 중복으로 제거
    has_new = any(("신규 첫 운영" in x or "신규 타겟 TEST" in x) for x in rows)
    if has_new:
        rows = [x for x in rows if "금번 성과" not in x or ("신규 첫 운영" in x or "신규 타겟 TEST" in x)]

    # 금번 성과 최대 1개
    seen_perf = False
    tmp = []
    for x in rows:
        if "금번 성과" in x:
            if seen_perf:
                continue
            seen_perf = True
        tmp.append(x)
    rows = tmp

    def score(x):
        if "신규 첫 운영" in x: return 100
        if "신규 타겟 TEST" in x: return 98
        if "금번 성과" in x: return 90
        if "성과 하락 요인" in x or "성과 상승 요인" in x: return 88
        if "가격 성공 조건" in x: return 84
        if "타겟 적합" in x: return 82
        if "프로모션" in x and ("의존" in x or "영향" in x): return 80
        if "운영 간격" in x or "피로도" in x: return 78
        if "가격 대비 성과" in x or "가격 경쟁력" in x or "가격 점검" in x: return 76
        if "시즌성" in x: return 65
        if "성과 성장" in x or "성과 둔화" in x: return 60
        return 50

    rows = sorted(enumerate(rows), key=lambda z: (-score(z[1]), z[0]))
    selected = [x for _, x in rows[:max_core]]

    if actions:
        selected.append(actions[0])

    # 주의 태그는 최대 2개 의미만 유지
    if warnings:
        w = warnings[0]
        if "·" in w:
            prefix, body = w.split(":", 1)
            tags = [t.strip() for t in body.split("·") if t.strip()][:2]
            w = prefix + ": " + " · ".join(tags)
        selected.append(w)
    return selected

def _daily_compare_priority():
    """원인 분석용 과거 비교 우선순위."""
    return (
        "동일 상품+동일 타겟/SEG",
        "동일 시즌",
        "최근 운영",
        "역대 최고",
    )

def _daily_test_principle():
    """재TEST는 원인 학습이 가능하도록 한 번에 한 변수 중심."""
    return "가격/구성/타겟/전시순서 중 핵심 변수 1개 우선 조정"

# 고급 분석 원칙:
# 1) 하락뿐 아니라 상승 시에도 과거 유사 회차 대비 성공 요인을 추정한다.
# 2) 역대 최고만 비교하지 않고 동일 타겟/SEG > 동일 시즌 > 최근 운영 > 역대 최고 순으로 비교한다.
# 3) 가격은 외부 최저가 차액뿐 아니라 과거 300만원 이상 달성 회차의 성공가격 조건을 참고한다.
#    단, 표본 1회면 '성공가격'으로 단정하지 않고 '과거 고성과 당시 가격'으로 표현한다.
# 4) 타겟 적합도는 평균매출 + 300만원 이상 달성률 + 운영횟수 + 최근 성과의 재현성을 함께 본다.
# 5) 프로모션 의존도는 프로모션/일반기간 표본이 충분할 때만 판단하고, 표본 부족 시 가능성으로 제한한다.
# 6) 피로도는 최근 하락 + 짧은 운영 간격이 함께 확인될 때 우선 추정한다.
# 7) 다음 운영 제안은 가격/구성/타겟/전시순서 중 한 번에 핵심 변수 1개를 우선 조정한다.
# 8) 출력은 핵심 인사이트 최대 3개 + 다음 운영 제안 1개를 기본으로 한다.
# === END V4.4.38 POLICY ===




st.set_page_config(
    page_title="MMS AI Dashboard",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown(
    """
    <style>
:root {
    --bg: #f5f7fb;
    --surface: #ffffff;
    --surface-soft: #f8fafc;
    --border: #e4e8ef;
    --text: #1f2937;
    --muted: #6b7280;
    --primary: #2f6fec;
    --primary-soft: #eef4ff;
    --success: #2e8b57;
    --warning: #d69e00;
    --danger: #d64545;
    --shadow: 0 6px 20px rgba(25, 42, 70, 0.06);
}

html, body, [class*="css"] {
    font-family: "Pretendard", "Noto Sans KR", "Apple SD Gothic Neo", sans-serif;
    color: var(--text);
}

.stApp {
    background: var(--bg);
}

.block-container {
    max-width: 100%;
    padding: 2.4rem 1.4rem 3rem;
}

[data-testid="stSidebar"] {
    background: #ffffff;
    border-right: 1px solid var(--border);
}

[data-testid="stSidebar"] .block-container {
    padding-top: 1.5rem;
}

.app-header {
    display: flex;
    align-items: center;
    justify-content: space-between;
    gap: 20px;
    background: linear-gradient(135deg, #ffffff 0%, #f7faff 100%);
    border: 1px solid var(--border);
    border-radius: 16px;
    padding: 22px 24px;
    box-shadow: var(--shadow);
    margin-bottom: 16px;
}

.app-title {
    font-size: 30px;
    font-weight: 800;
    letter-spacing: -0.7px;
    margin: 0;
    line-height: 1.35;
}

.app-subtitle {
    color: var(--muted);
    font-size: 14px;
    margin-top: 6px;
}

.status-badge {
    display: inline-flex;
    align-items: center;
    gap: 8px;
    background: var(--primary-soft);
    color: #2859bb;
    border: 1px solid #d8e5ff;
    border-radius: 999px;
    padding: 9px 14px;
    font-size: 13px;
    white-space: nowrap;
}

.section-title {
    font-size: 21px;
    font-weight: 800;
    letter-spacing: -0.4px;
    margin: 24px 0 12px;
    padding-left: 12px;
    border-left: 4px solid var(--primary);
}

.subsection-title {
    font-size: 17px;
    font-weight: 750;
    margin: 18px 0 10px;
}

.card-shell {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 14px;
    padding: 18px;
    box-shadow: var(--shadow);
    margin-bottom: 14px;
}

.metric-card {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 14px;
    padding: 16px 17px;
    min-height: 98px;
    box-shadow: 0 3px 12px rgba(25, 42, 70, 0.04);
    transition: transform .15s ease, box-shadow .15s ease;
}

.metric-card:hover {
    transform: translateY(-1px);
    box-shadow: 0 7px 18px rgba(25, 42, 70, 0.08);
}

.metric-label {
    font-size: 12px;
    color: var(--muted);
    font-weight: 650;
    margin-bottom: 8px;
}

.metric-value {
    font-size: 23px;
    font-weight: 800;
    letter-spacing: -0.5px;
    line-height: 1.2;
}

.metric-delta {
    font-size: 12px;
    margin-top: 8px;
    color: var(--muted);
}

.insight-box {
    background: linear-gradient(180deg, #f8fbff 0%, #f3f7fd 100%);
    border: 1px solid #dfe8f5;
    border-radius: 14px;
    padding: 18px 20px;
    line-height: 1.6;
    white-space: pre-wrap;
    box-shadow: 0 3px 12px rgba(25, 42, 70, 0.035);
}

.asset-card {
    border: 1px solid var(--border);
    border-radius: 14px;
    background: var(--surface);
    padding: 14px;
    box-sizing: border-box;
    box-shadow: 0 3px 12px rgba(25, 42, 70, 0.045);
}

/* 문구 카드가 좌우 전체 높이의 기준이 됨 */
.asset-message-card {
    height: auto;
    min-height: 320px;
    overflow: visible;
    white-space: pre-wrap;
    font-size: 15px;
    line-height: 1.75;
    background: #fbfcfe;
}

/* 이미지 자체 높이가 row를 늘리지 않도록 카드 안에 absolute 배치 */
.asset-image-card {
    position: relative;
    height: 100%;
    min-height: 320px;
    overflow: hidden;
    background: var(--surface);
}

.asset-image-card img {
    position: absolute;
    inset: 14px;
    width: calc(100% - 28px);
    height: calc(100% - 28px);
    min-height: 0;
    max-height: none;
    object-fit: contain;
    object-position: center center;
    border-radius: 10px;
}

.asset-empty {
    height: 398px;
    display: flex;
    align-items: center;
    justify-content: center;
    color: var(--muted);
    text-align: center;
}

[data-testid="stDataFrame"] {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 14px;
    overflow: hidden;
    box-shadow: 0 3px 12px rgba(25, 42, 70, 0.04);
}

[data-testid="stDataFrame"] [role="columnheader"] {
    background: #f5f8fc !important;
    font-weight: 750 !important;
}

[data-testid="stSelectbox"] > div,
[data-testid="stDateInput"] > div,
[data-testid="stMultiSelect"] > div {
    background: var(--surface);
    border-radius: 10px;
}

.stTabs [data-baseweb="tab-list"] {
    gap: 8px;
    background: #eef2f7;
    padding: 5px;
    border-radius: 12px;
}

.stTabs [data-baseweb="tab"] {
    border-radius: 9px;
    padding: 8px 14px;
    font-weight: 700;
}

.stTabs [aria-selected="true"] {
    background: #ffffff !important;
    color: var(--primary) !important;
    box-shadow: 0 2px 8px rgba(25, 42, 70, 0.08);
}

div[data-testid="stPlotlyChart"] {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 14px;
    padding: 8px;
    box-shadow: 0 3px 12px rgba(25, 42, 70, 0.04);
}

.stButton > button,
.stDownloadButton > button {
    border-radius: 10px;
    font-weight: 750;
    border: 1px solid #d5deed;
}

.stDownloadButton > button {
    background: var(--primary);
    color: #ffffff;
    border-color: var(--primary);
}

hr {
    border: 0;
    border-top: 1px solid var(--border);
}

/* 일일 상품 인사이트: 가독성 확보 */
.compact-insight {
    font-size: 16px !important;
    line-height: 1.72 !important;
    margin: 0;
    padding: 2px 0;
}
.compact-insight .insight-row {
    font-size: 16px !important;
    line-height: 1.72 !important;
    margin: 0 0 9px 0 !important;
    padding: 0;
}
.compact-insight .evidence {
    color: var(--muted);
    font-size: 13px !important;
    line-height: 1.55 !important;
}
[data-testid="stExpander"] details summary {
    padding-top: 0.55rem;
    padding-bottom: 0.55rem;
}
[data-testid="stExpander"] details > div {
    padding-top: 0.25rem;
}



/* V4.4.35 발송소재 실제 이미지 하단을 MMS 문구 [KT] 하단과 동기화 */
[data-testid="stHorizontalBlock"]:has(.asset-message-card) {
    align-items: stretch !important;
}
[data-testid="stHorizontalBlock"]:has(.asset-message-card) > [data-testid="stColumn"] {
    align-self: stretch !important;
    display: flex !important;
    flex-direction: column !important;
}
[data-testid="stHorizontalBlock"]:has(.asset-message-card) > [data-testid="stColumn"] > div,
[data-testid="stHorizontalBlock"]:has(.asset-message-card) > [data-testid="stColumn"] [data-testid="stVerticalBlock"],
[data-testid="stHorizontalBlock"]:has(.asset-message-card) > [data-testid="stColumn"] [data-testid="stVerticalBlockBorderWrapper"] {
    flex: 1 1 auto !important;
    height: 100% !important;
}
[data-testid="stHorizontalBlock"]:has(.asset-message-card) .asset-card {
    width: 100% !important;
    height: 100% !important;
    min-height: 430px !important;
    box-sizing: border-box !important;
}
[data-testid="stHorizontalBlock"]:has(.asset-message-card) .asset-message-card {
    height: 100% !important;
    min-height: 430px !important;
    overflow: visible !important;
}
[data-testid="stHorizontalBlock"]:has(.asset-message-card) .asset-image-card {
    position: relative !important;
    height: 100% !important;
    min-height: 430px !important;
    padding: 0 !important;
    overflow: hidden !important;
}
[data-testid="stHorizontalBlock"]:has(.asset-message-card) .asset-image-card img {
    position: absolute !important;
    inset: 0 !important;
    width: 100% !important;
    height: 100% !important;
    min-width: 0 !important;
    min-height: 0 !important;
    max-width: 100% !important;
    max-height: 100% !important;
    object-fit: contain !important;
    object-position: top center !important;
    transform: none !important;
    margin: 0 !important;
    padding: 0 !important;
    display: block !important;
}


/* V4.4.30 발송소재: 원본 이미지 비율은 유지하고 좌우 박스 하단만 정확히 맞춤 */
.daily-asset-pair {
    display: grid;
    grid-template-columns: minmax(0, 1fr) minmax(0, 1.35fr);
    gap: 1rem;
    align-items: stretch;
    width: 100%;
}
.daily-asset-pair .asset-card {
    min-height: 0 !important;
    margin: 0;
}
.daily-asset-pair .asset-image-card {
    /* 이미지 자체 높이가 Grid 행 높이를 키우지 않도록 카드 안에 절대배치 */
    position: relative;
    overflow: hidden;
    padding: 0;
    min-height: 0 !important;
    background: var(--surface);
}
.daily-asset-pair .asset-image-card img {
    /* 오른쪽 MMS 문구 카드가 정한 높이 안에서만 원본 전체 노출 */
    position: absolute !important;
    top: 20px !important;
    right: 20px !important;
    bottom: 20px !important;
    left: 20px !important;
    width: calc(100% - 40px) !important;
    height: calc(100% - 40px) !important;
    min-width: 0 !important;
    min-height: 0 !important;
    max-width: calc(100% - 40px) !important;
    max-height: calc(100% - 40px) !important;
    object-fit: contain !important;
    object-position: center center !important;
    margin: 0 !important;
    padding: 0 !important;
    border-radius: 10px;
    box-sizing: border-box !important;
}
.daily-asset-pair .asset-message-card {
    /* 문구 내용 높이가 전체 Grid 행 높이의 기준 */
    min-height: 0 !important;
    height: auto !important;
    overflow: visible;
    padding: 18px 14px 14px;
}

@media (max-width: 900px) {

    [data-testid="stHorizontalBlock"]:has(.asset-message-card) > [data-testid="stColumn"],
    [data-testid="stHorizontalBlock"]:has(.asset-message-card) > [data-testid="stColumn"] > div,
    [data-testid="stHorizontalBlock"]:has(.asset-message-card) > [data-testid="stColumn"] [data-testid="stVerticalBlock"] {
        height: auto !important;
        min-height: 0 !important;
    }

    .block-container {
        padding: 1.6rem 0.75rem 2rem;
    }

    .app-header {
        flex-direction: column;
        align-items: flex-start;
    }

    .app-title {
        font-size: 25px;
    }

    .status-badge {
        white-space: normal;
    }

    .metric-value {
        font-size: 20px;
    }

    .asset-card,
    .asset-message-card {
        height: auto;
        min-height: 320px;
    }

    .asset-image-card {
        height: auto;
        min-height: 320px;
    }

    .asset-image-card img {
        height: auto;
        min-height: 0;
        max-height: 380px;
    }
}

@media (max-width: 900px) {
    .daily-asset-pair {
        grid-template-columns: 1fr;
    }
    .daily-asset-pair .asset-image-card,
    .daily-asset-pair .asset-message-card {
        height: auto;
    }
    .daily-asset-pair .asset-image-card img {
        height: auto;
        max-height: none;
    }
}


/* V4.4.30 DESKTOP FINAL OVERRIDE: legacy selector 충돌 방지 */
@media (min-width: 901px) {
    .daily-asset-pair .asset-image-card {
        position: relative !important;
        overflow: hidden !important;
        padding: 0 !important;
        box-sizing: border-box !important;
    }
    .daily-asset-pair .asset-image-card > img {
        position: absolute !important;
        top: 20px !important;
        right: 20px !important;
        bottom: 20px !important;
        left: 20px !important;
        inset: 20px !important;
        width: calc(100% - 40px) !important;
        height: calc(100% - 40px) !important;
        min-width: 0 !important;
        min-height: 0 !important;
        max-width: calc(100% - 40px) !important;
        max-height: calc(100% - 40px) !important;
        object-fit: contain !important;
        object-position: center center !important;
        display: block !important;
        margin: 0 !important;
        padding: 0 !important;
        box-sizing: border-box !important;
    }
}

</style>
    """,
    unsafe_allow_html=True,
)

GRADE_ORDER = ["핵심 상품", "우수 상품", "안정 상품", "관찰 상품", "부진 상품"]
CASE_ORDER = [
    "가격 경쟁력 부족 사례",
    "기네스 갱신 사례",
    "타겟 확대 운영 사례",
    "시즌 상품 사례",
    "운영 피로도 사례",
    "보답프로그램 영향 사례",
]


def first_col(df: pd.DataFrame, names: Iterable[str]) -> str | None:
    for name in names:
        if name in df.columns:
            return name
    return None


def num(series: pd.Series) -> pd.Series:
    if series.dtype == object:
        series = (
            series.astype(str)
            .str.replace(",", "", regex=False)
            .str.replace("%", "", regex=False)
            .replace({"-": "0", "nan": "0", "None": "0", "": "0"})
        )
    return pd.to_numeric(series, errors="coerce").fillna(0)


def safe_div(a, b):
    return a / b.replace(0, pd.NA)


def fmt_num(v) -> str:
    return f"{int(round(float(v))):,}"


def fmt_pct(v) -> str:
    return f"{float(v) * 100:.1f}%"


def compact_money(v: float) -> str:
    v = float(v)
    if abs(v) >= 10_000_000:
        return f"{v/10_000_000:.1f}천만원"
    if abs(v) >= 1_000_000:
        return f"{v/1_000_000:.1f}백만원"
    if abs(v) >= 10_000:
        return f"{v/10_000:.1f}만원"
    return f"{int(v):,}원"




def stable_variant(key: str, options: list[str]) -> str:
    """동일 상품·조건에서는 같은 문장을 유지하면서 표현 반복을 줄입니다."""
    if not options:
        return ""
    score = sum((idx + 1) * ord(ch) for idx, ch in enumerate(str(key)))
    return options[score % len(options)]

def product_grade(amount: float) -> str:
    if amount < 1_000_000:
        return "부진 상품"
    if amount < 2_000_000:
        return "관찰 상품"
    if amount < 3_000_000:
        return "안정 상품"
    if amount < 5_000_000:
        return "우수 상품"
    return "핵심 상품"


def parse_yyyymmdd_date(series: pd.Series) -> pd.Series:
    """20260716 같은 숫자·문자 날짜와 일반 날짜를 안전하게 변환합니다."""
    text = series.astype(str).str.strip().str.replace(r"\.0$", "", regex=True)
    compact_mask = text.str.fullmatch(r"\d{8}")
    parsed = pd.Series(pd.NaT, index=series.index, dtype="datetime64[ns]")
    parsed.loc[compact_mask] = pd.to_datetime(
        text.loc[compact_mask], format="%Y%m%d", errors="coerce"
    )
    parsed.loc[~compact_mask] = pd.to_datetime(
        series.loc[~compact_mask], errors="coerce"
    )
    return parsed


def normalize_product(df: pd.DataFrame) -> pd.DataFrame:
    d = df.copy()
    date_col = first_col(d, ["발송일", "날짜", "일자"])
    if date_col is None:
        raise ValueError("상품 시트에서 발송일/날짜/일자 열을 찾을 수 없습니다.")
    d["_date"] = parse_yyyymmdd_date(d[date_col])
    d = d[d["_date"].notna()].copy()
    for c in ["정상가", "멤버십혜택가", "할인율", "전시순서", "주문건수", "주문수량", "주문금액", "발송일 최저가"]:
        if c in d.columns:
            d[c] = num(d[c])
    d["_year"] = d["_date"].dt.year.astype(int)
    d["_month"] = d["_date"].dt.month.astype(int)
    return d


def normalize_send(df: pd.DataFrame) -> pd.DataFrame:
    d = df.copy()
    date_col = first_col(d, ["발송일시2", "발송일", "날짜", "일자"])
    if date_col is None:
        raise ValueError("소재 시트에서 발송일시2/발송일/날짜/일자 열을 찾을 수 없습니다.")
    d["_date"] = parse_yyyymmdd_date(d[date_col])
    d = d[d["_date"].notna()].copy()
    for c in [
        "상품수", "URL", "총 발송 건수", "발송 성공 건수",
        "클릭 수", "클릭 수(uniq)", "반응율", "반응율(uniq)",
        "주문건수", "주문수량", "주문금액", "객단가",
        "클릭>구매 전환율", "발송>구매 전환율", "SPM",
    ]:
        if c in d.columns:
            d[c] = num(d[c])
    d["_year"] = d["_date"].dt.year.astype(int)
    d["_month"] = d["_date"].dt.month.astype(int)
    return d


def normalize_message(df: pd.DataFrame | None) -> pd.DataFrame:
    """선택적인 '문구' 시트를 정규화합니다."""
    if df is None or df.empty:
        return pd.DataFrame(columns=["캠페인명", "MMS문구"])

    d = df.copy()
    campaign_col = first_col(d, ["캠페인명", "캠페인", "Campaign", "campaign"])
    message_col = first_col(d, ["MMS문구", "MMS 문구", "발송문구", "문구"])

    if campaign_col is None or message_col is None:
        return pd.DataFrame(columns=["캠페인명", "MMS문구"])

    d = d[[campaign_col, message_col]].copy()
    d.columns = ["캠페인명", "MMS문구"]
    d["캠페인명"] = d["캠페인명"].fillna("").astype(str).str.strip()
    d["MMS문구"] = d["MMS문구"].apply(clean_mms_message if "clean_mms_message" in globals() else lambda x: str(x))
    d = d[d["캠페인명"].ne("")].drop_duplicates("캠페인명", keep="last")
    return d.reset_index(drop=True)


def normalize_lowest(df: pd.DataFrame | None) -> pd.DataFrame:
    """선택적인 '최저가' 시트를 정규화합니다."""
    if df is None or df.empty:
        return pd.DataFrame()

    d = df.copy()
    for c in ["쇼라코드", "알파코드"]:
        if c in d.columns:
            d[c] = d[c].astype(str).str.replace(".0", "", regex=False).str.strip()

    date_col = first_col(d, ["발송일", "날짜", "일자"])
    if date_col:
        d["_date"] = parse_yyyymmdd_date(d[date_col])

    for c in [
        "네이버 최저가", "최저가", "비교가", "네이버가",
        "멤버십혜택가", "가격차이"
    ]:
        if c in d.columns:
            d[c] = num(d[c])

    price_col = first_col(d, ["네이버 최저가", "최저가", "네이버가", "비교가"])
    if price_col and price_col != "최저가":
        d = d.rename(columns={price_col: "최저가"})

    return d


def normalize_promotion(df: pd.DataFrame | None) -> pd.DataFrame:
    """로우 프로모션 시트의 프로모션명·시작일·종료일을 정규화합니다."""
    if df is None or df.empty:
        return pd.DataFrame(columns=["프로모션명", "_start_date", "_end_date", "스킴"])
    d = df.copy()
    name_col = first_col(d, ["프로모션명", "프로모션", "행사명", "기획전명", "보답프로그램"])
    start_col = first_col(d, ["시작일", "시작일자", "발송일", "날짜", "일자"])
    end_col = first_col(d, ["종료일", "종료일자"])
    if name_col is None or start_col is None:
        return pd.DataFrame(columns=["프로모션명", "_start_date", "_end_date", "스킴"])
    out = pd.DataFrame()
    out["프로모션명"] = d[name_col].fillna("").astype(str).str.strip()
    out["_start_date"] = parse_yyyymmdd_date(d[start_col]).dt.normalize()
    out["_end_date"] = parse_yyyymmdd_date(d[end_col]).dt.normalize() if end_col else out["_start_date"]
    out["_end_date"] = out["_end_date"].fillna(out["_start_date"])
    scheme_col = first_col(d, ["스킴", "혜택", "프로모션 스킴"])
    out["스킴"] = d[scheme_col].fillna("").astype(str).str.strip() if scheme_col else ""
    return out[out["프로모션명"].ne("") & out["_start_date"].notna()].reset_index(drop=True)


def promotion_name_for_date(date_value, promotions: pd.DataFrame) -> str:
    date_value = pd.to_datetime(date_value, errors="coerce")
    if promotions is None or promotions.empty or pd.isna(date_value):
        return "-"
    day = date_value.normalize()
    matched = promotions[(promotions["_start_date"] <= day) & (promotions["_end_date"] >= day)]
    if matched.empty:
        return "-"
    return str(matched.iloc[-1]["프로모션명"]).strip() or "-"


def apply_promotion_periods(products: pd.DataFrame, promotions: pd.DataFrame) -> pd.DataFrame:
    out = products.copy()
    out["프로모션명"] = out["_date"].map(lambda value: promotion_name_for_date(value, promotions))
    return out


@st.cache_data(show_spinner=False)
def load_excel_bytes(file_bytes: bytes):
    workbook = pd.ExcelFile(io.BytesIO(file_bytes))
    product = pd.read_excel(workbook, sheet_name="상품")
    send = pd.read_excel(workbook, sheet_name="소재")
    lowest = (
        pd.read_excel(workbook, sheet_name="최저가")
        if "최저가" in workbook.sheet_names
        else pd.DataFrame()
    )
    messages = (
        pd.read_excel(workbook, sheet_name="문구")
        if "문구" in workbook.sheet_names
        else pd.DataFrame()
    )
    promotion_sheet = "로우 프로모션" if "로우 프로모션" in workbook.sheet_names else ("프로모션" if "프로모션" in workbook.sheet_names else None)
    promotions = pd.read_excel(workbook, sheet_name=promotion_sheet) if promotion_sheet else pd.DataFrame()
    normalized_promotions = normalize_promotion(promotions)
    normalized_products = apply_promotion_periods(normalize_product(product), normalized_promotions)
    return (
        normalized_products,
        normalize_send(send),
        normalize_lowest(lowest),
        normalize_message(messages),
        normalized_promotions,
    )


def extract_google_sheet_id(url: str) -> str:
    match = re.search(r"/spreadsheets/d/([a-zA-Z0-9-_]+)", url)
    if not match:
        raise ValueError("구글시트 주소에서 문서 ID를 찾을 수 없습니다.")
    return match.group(1)



def google_requests_get(url: str, **kwargs):
    """
    Google 요청은 먼저 정상 SSL 검증으로 시도합니다.
    회사망 자체 서명 인증서 오류가 발생할 때만 검증 없이 한 번 재시도합니다.
    """
    try:
        return requests.get(url, **kwargs)
    except requests.exceptions.SSLError:
        urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
        retry_kwargs = dict(kwargs)
        retry_kwargs["verify"] = False
        return requests.get(url, **retry_kwargs)

def read_google_csv(sheet_id: str, sheet_name: str) -> pd.DataFrame:
    csv_url = (
        f"https://docs.google.com/spreadsheets/d/{sheet_id}/gviz/tq"
        f"?tqx=out:csv&sheet={requests.utils.quote(sheet_name)}"
    )
    response = google_requests_get(csv_url, timeout=30)
    response.raise_for_status()
    text = response.text
    if text.lstrip().lower().startswith("<!doctype html") or "<html" in text[:300].lower():
        raise PermissionError(
            f"'{sheet_name}' 탭을 불러오지 못했습니다. "
            "구글시트 공유 권한을 '링크가 있는 모든 사용자/뷰어'로 설정해주세요."
        )
    return pd.read_csv(io.StringIO(text))


@st.cache_data(ttl=300, show_spinner=False)
def load_google_sheet(url: str):
    """구글시트를 불러오고 5분간 캐시합니다."""
    sheet_id = extract_google_sheet_id(url)
    errors = []

    # 1차: 전체 문서를 XLSX로 내보내기
    try:
        export_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=xlsx"
        response = google_requests_get(
            export_url,
            timeout=45,
            headers={"User-Agent": "Mozilla/5.0"},
        )
        response.raise_for_status()
        content_type = response.headers.get("content-type", "").lower()
        if "text/html" not in content_type and response.content[:2] == b"PK":
            return load_excel_bytes(response.content)
        errors.append("전체 XLSX 내보내기 응답이 엑셀 파일이 아니었습니다.")
    except Exception as exc:
        errors.append(
            "전체 XLSX 내보내기 실패: "
            f"{exc}"
        )

    # 2차: 상품/소재 탭을 각각 CSV로 불러오기
    try:
        product = read_google_csv(sheet_id, "상품")
        send = read_google_csv(sheet_id, "소재")
        try:
            lowest = read_google_csv(sheet_id, "최저가")
        except Exception:
            lowest = pd.DataFrame()
        try:
            messages = read_google_csv(sheet_id, "문구")
        except Exception:
            messages = pd.DataFrame()
        try:
            promotions_raw = read_google_csv(sheet_id, "로우 프로모션")
        except Exception:
            try:
                promotions_raw = read_google_csv(sheet_id, "프로모션")
            except Exception:
                promotions_raw = pd.DataFrame()
        normalized_promotions = normalize_promotion(promotions_raw)
        normalized_products = apply_promotion_periods(normalize_product(product), normalized_promotions)
        return (
            normalized_products,
            normalize_send(send),
            normalize_lowest(lowest),
            normalize_message(messages),
            normalized_promotions,
        )
    except Exception as exc:
        errors.append(f"상품·소재 CSV 불러오기 실패: {exc}")

    raise RuntimeError(" / ".join(errors))


def sync_google_sheet(url: str, force: bool = False):
    """자동 또는 수동으로 구글시트를 세션 데이터에 반영합니다."""
    if force:
        load_google_sheet.clear()

    products, sends, lowest, messages, promotions = load_google_sheet(url)
    st.session_state.products = products
    st.session_state.sends = sends
    st.session_state.lowest = lowest
    st.session_state.messages = messages
    st.session_state.promotions = promotions
    st.session_state.source_name = "구글시트 자동연동"
    st.session_state.synced_at = datetime.now()
    st.session_state.google_sync_error = None


def aggregate_send(data: pd.DataFrame, mode: str) -> pd.DataFrame:
    d = data.copy()
    if d.empty:
        return pd.DataFrame()

    if mode == "Monthly":
        d["_label"] = d["_date"].dt.strftime("%Y-%m")
        d["_sort1"] = d["_date"].dt.year
        d["_sort2"] = d["_date"].dt.month
    elif mode == "Weekly":
        d["_label"] = d["주차"].astype(str)
        d["_sort1"] = d["_date"].dt.year
        d["_sort2"] = pd.to_numeric(
            d["_label"].str.extract(r"(\d+)")[0], errors="coerce"
        ).fillna(0)
    else:
        d["_label"] = d["_date"].dt.strftime("%m%d")
        d["_sort1"] = d["_date"].dt.year
        d["_sort2"] = d["_date"].dt.dayofyear

    send_col = first_col(d, ["발송 성공 건수", "총 발송 건수"])
    click_col = first_col(d, ["클릭 수(uniq)", "클릭 수"])

    g = d.groupby("_label", as_index=False).agg(
        연도=("_sort1", "max"),
        _sort1=("_sort1", "max"),
        _sort2=("_sort2", "max"),
        발송횟수=("_label", "size"),
        상품수=("상품수", "sum"),
        URL=("URL", "sum"),
        발송건수=(send_col, "sum"),
        클릭수=(click_col, "sum"),
        주문건수=("주문건수", "sum"),
        주문수량=("주문수량", "sum"),
        주문금액=("주문금액", "sum"),
    ).sort_values(["_sort1", "_sort2"])

    if mode == "Monthly":
        g["월"] = g["_label"].str[-2:].astype(int)
    elif mode == "Daily":
        month_map = d.groupby("_label")["_month"].max()
        g["월"] = g["_label"].map(month_map).astype(int)

    g["반응율(Uniq CTR)"] = safe_div(g["클릭수"], g["발송건수"])
    g["객단가"] = safe_div(g["주문금액"], g["주문건수"])
    g["클릭 CVR"] = safe_div(g["주문건수"], g["클릭수"])
    g["발송 CVR"] = safe_div(g["주문건수"], g["발송건수"])
    g["클릭당매출(RPC)"] = safe_div(g["주문금액"], g["클릭수"])
    g["발송대비매출(SPM)"] = safe_div(g["주문금액"], g["발송건수"])
    g["발송당매출(발송횟수)"] = safe_div(g["주문금액"], g["발송횟수"])

    return g.fillna(0).reset_index(drop=True)


def add_changes(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    out["주문금액 증감"] = out["주문금액"].pct_change()
    out["CTR 증감"] = out["반응율(Uniq CTR)"].diff()
    out["SPM 증감"] = out["발송대비매출(SPM)"].pct_change()
    out["발송당매출 증감"] = out["발송당매출(발송횟수)"].pct_change()
    return out.replace([float("inf"), float("-inf")], pd.NA)


def change_label(value, pp: bool = False) -> str:
    if pd.isna(value):
        return "-"
    arrow = "▲" if value > 0 else ("▼" if value < 0 else "-")
    if pp:
        return f"{arrow}{abs(value)*100:.1f}%p"
    return f"{arrow}{abs(value)*100:.1f}%"


def format_home_table(df: pd.DataFrame, mode: str) -> pd.DataFrame:
    out = add_changes(df).copy()
    if mode == "Monthly":
        out = out.rename(columns={"_label": "기간"})
        order = [
            "연도", "월", "발송횟수", "상품수", "URL", "발송건수", "클릭수",
            "반응율(Uniq CTR)", "주문건수", "주문수량", "주문금액", "객단가",
            "클릭 CVR", "발송 CVR", "클릭당매출(RPC)",
            "발송대비매출(SPM)", "발송당매출(발송횟수)",
            "주문금액 증감", "CTR 증감", "SPM 증감", "발송당매출 증감",
        ]
    elif mode == "Weekly":
        out = out.rename(columns={"_label": "주차"})
        order = [
            "연도", "주차", "발송횟수", "상품수", "URL", "발송건수", "클릭수",
            "반응율(Uniq CTR)", "주문건수", "주문수량", "주문금액", "객단가",
            "클릭 CVR", "발송 CVR", "클릭당매출(RPC)",
            "발송대비매출(SPM)", "발송당매출(발송횟수)",
            "주문금액 증감", "CTR 증감", "SPM 증감", "발송당매출 증감",
        ]
    else:
        out = out.rename(columns={"_label": "일자"})
        order = [
            "연도", "월", "일자", "발송횟수", "상품수", "URL", "발송건수", "클릭수",
            "반응율(Uniq CTR)", "주문건수", "주문수량", "주문금액", "객단가",
            "클릭 CVR", "발송 CVR", "클릭당매출(RPC)",
            "발송대비매출(SPM)", "발송당매출(발송횟수)",
            "주문금액 증감", "CTR 증감", "SPM 증감", "발송당매출 증감",
        ]

    out = out[[c for c in order if c in out.columns]].copy()

    for c in ["반응율(Uniq CTR)", "클릭 CVR", "발송 CVR"]:
        if c in out.columns:
            out[c] = out[c].map(fmt_pct)

    for c in ["주문금액 증감", "SPM 증감", "발송당매출 증감"]:
        if c in out.columns:
            out[c] = out[c].map(change_label)

    if "CTR 증감" in out.columns:
        out["CTR 증감"] = out["CTR 증감"].map(lambda x: change_label(x, pp=True))

    for c in [
        "발송횟수", "상품수", "URL", "발송건수", "클릭수",
        "주문건수", "주문수량", "주문금액", "객단가",
        "클릭당매출(RPC)", "발송당매출(발송횟수)",
    ]:
        if c in out.columns:
            out[c] = out[c].map(fmt_num)

    if "발송대비매출(SPM)" in out.columns:
        out["발송대비매출(SPM)"] = out["발송대비매출(SPM)"].map(lambda x: f"{x:.1f}")

    return out


def filter_monthly_period(df: pd.DataFrame, option: str, start_key: str | None, end_key: str | None):
    if df.empty or option == "전체":
        return df
    if option.startswith("최근"):
        count = int(re.search(r"\d+", option).group())
        return df.tail(count)
    if option == "직접 선택" and start_key and end_key:
        labels = df["_label"].astype(str)
        return df[(labels >= start_key) & (labels <= end_key)]
    return df


def filter_weekly_period(df: pd.DataFrame, start_week: str, end_week: str):
    if df.empty:
        return df
    labels = df["_label"].astype(str).tolist()
    if start_week not in labels or end_week not in labels:
        return df
    start_idx, end_idx = labels.index(start_week), labels.index(end_week)
    lo, hi = sorted([start_idx, end_idx])
    return df.iloc[lo : hi + 1]


def trend_chart(df: pd.DataFrame, title: str, color: str) -> go.Figure:
    labels = df["_label"].astype(str).tolist()
    vals = df["발송대비매출(SPM)"].astype(float).tolist()
    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=labels,
            y=vals,
            marker_color=color,
            name="SPM",
            text=[f"{v:.1f}" for v in vals],
            textposition="outside",
            cliponaxis=False,
        )
    )
    if len(vals) >= 2:
        trend = pd.Series(vals).rolling(3, min_periods=1).mean()
        fig.add_trace(
            go.Scatter(
                x=labels,
                y=trend,
                mode="lines",
                name="추세",
                line=dict(color="#333", width=2, dash="dot"),
            )
        )
    ymax = max(vals) if vals else 0
    fig.update_layout(
        title=dict(text=title, x=.5, font=dict(size=23)),
        height=480,
        margin=dict(l=60, r=30, t=72, b=100),
        plot_bgcolor="#ffffff",
        barmode="overlay",
        yaxis=dict(
            tickformat=",.1f",
            gridcolor="#ddd",
            range=[0, max(ymax * 1.22, 10)],
        ),
        xaxis=dict(
            tickangle=-35 if len(labels) > 12 else 0,
            automargin=True,
        ),
        legend=dict(orientation="h", y=-.23),
    )
    return fig


def delta_for_latest(df: pd.DataFrame, metric: str, pp: bool = False) -> str:
    if len(df) < 2:
        return "-"
    cur, prev = df.iloc[-1][metric], df.iloc[-2][metric]
    if pp:
        return change_label(cur - prev, pp=True)
    if prev == 0:
        return "-"
    return change_label((cur - prev) / abs(prev))


def classify_cases(row: pd.Series, history: pd.DataFrame) -> list[str]:
    cases = []
    name = str(row["상품명"])
    amount = float(row["주문금액"])
    prior = history[(history["상품명"] == name) & (history["_date"] < row["_date"])].sort_values("_date")

    if len(prior) and amount > prior["주문금액"].max() and amount >= 3_000_000:
        cases.append("기네스 갱신 사례")

    if len(prior):
        last = prior.iloc[-1]
        if (
            str(last.get("성별", "")) != str(row.get("성별", ""))
            or str(last.get("연령", "")) != str(row.get("연령", ""))
        ):
            cases.append("타겟 확대 운영 사례")
        gap = (row["_date"] - last["_date"]).days
        if gap <= 21 and amount < last["주문금액"] * 0.75:
            cases.append("운영 피로도 사례")
        if row.get("멤버십혜택가", 0) < last.get("멤버십혜택가", 0) and amount < last["주문금액"]:
            cases.append("가격 경쟁력 부족 사례")

    season_words = ["선풍기", "에어컨", "우양산", "래쉬가드", "삼계탕", "장어", "아이스크림", "드라이기"]
    if any(word in name for word in season_words):
        cases.append("시즌 상품 사례")

    if amount < 1_000_000 and float(row.get("할인율", 0)) >= 0.5:
        cases.append("가격 경쟁력 부족 사례")

    return list(dict.fromkeys(cases))


def product_history_rows(row: pd.Series, history: pd.DataFrame) -> pd.DataFrame:
    """현재 행보다 이전의 동일 상품 이력을 코드 우선순위로 찾습니다."""
    prior = history[history["_date"] < row["_date"]].copy()

    for key in ["쇼라코드", "알파코드"]:
        current = clean_identifier_value(row.get(key, ""))
        if current and key in prior.columns:
            candidates = prior[prior[key].map(clean_identifier_value).eq(current)]
            if not candidates.empty:
                return candidates.sort_values("_date")

    name = str(row.get("상품명", "")).strip()
    if name and "상품명" in prior.columns:
        return prior[prior["상품명"].astype(str).str.strip().eq(name)].sort_values("_date")

    return prior.iloc[0:0].copy()


def target_label(row: pd.Series, include_seg: bool = True) -> str:
    values = []
    for col in ["성별", "연령"]:
        value = clean_identifier_value(row.get(col, ""))
        if value:
            values.append(value)
    if include_seg:
        seg = clean_identifier_value(row.get("SEG", ""))
        if seg:
            values.append(f"SEG{seg}")
    return " ".join(values).strip()


def base_target_label(row: pd.Series) -> str:
    """SEG를 제외한 성별·연령 기준 타겟 라벨."""
    return target_label(row, include_seg=False)


def product_history_summary(row: pd.Series, history: pd.DataFrame) -> dict:
    prior = product_history_rows(row, history)
    if prior.empty:
        return {
            "운영횟수": 0,
            "평균매출": 0,
            "최고매출": 0,
            "최고타겟": "",
            "최근이력": None,
            "동일타겟이력": pd.DataFrame(),
            "과거이력": prior,
        }

    prior = prior.copy()
    prior["_target"] = prior.apply(target_label, axis=1)
    target_stats = (
        prior.groupby("_target", dropna=False)["주문금액"]
        .agg(["mean", "max", "count"])
        .sort_values(["mean", "max", "count"], ascending=False)
    )
    best_target = ""
    if not target_stats.empty:
        best_target = str(target_stats.index[0]).strip()

    current_target = target_label(row)
    same_target = prior[prior["_target"].eq(current_target)].copy()

    return {
        "운영횟수": int(len(prior)),
        "평균매출": float(prior["주문금액"].mean()),
        "최고매출": float(prior["주문금액"].max()),
        "최고타겟": best_target,
        "최근이력": prior.iloc[-1],
        "동일타겟이력": same_target,
        "과거이력": prior,
    }


def product_history_including_current(row: pd.Series, history: pd.DataFrame) -> pd.DataFrame:
    """현재 발송 건을 포함해 동일 상품의 누적 이력을 찾습니다."""
    cumulative = history[history["_date"] <= row["_date"]].copy()

    for key in ["쇼라코드", "알파코드"]:
        current = clean_identifier_value(row.get(key, ""))
        if current and key in cumulative.columns:
            candidates = cumulative[cumulative[key].map(clean_identifier_value).eq(current)]
            if not candidates.empty:
                return candidates.sort_values(["_date", "주문금액"])

    name = str(row.get("상품명", "")).strip()
    if name and "상품명" in cumulative.columns:
        return cumulative[
            cumulative["상품명"].astype(str).str.strip().eq(name)
        ].sort_values(["_date", "주문금액"])

    return cumulative.iloc[0:0].copy()


def add_history_columns(current_df: pd.DataFrame, history: pd.DataFrame) -> pd.DataFrame:
    """일일 상품표용 최고 실적 정보는 현재 발송일을 포함해 계산합니다."""
    out = current_df.copy()

    highest_amounts = []
    highest_dates = []
    highest_targets = []

    for _, row in out.iterrows():
        cumulative = product_history_including_current(row, history)

        if cumulative.empty:
            highest_amounts.append(float(row.get("주문금액", 0)))
            highest_dates.append(row.get("_date"))
            highest_targets.append(target_label(row))
            continue

        max_amount = float(cumulative["주문금액"].max())
        # 동일 최고매출이 여러 건이면 가장 최근 발송 건을 사용
        best_rows = cumulative[cumulative["주문금액"].eq(max_amount)].sort_values("_date")
        best_row = best_rows.iloc[-1]

        highest_amounts.append(max_amount)
        highest_dates.append(best_row.get("_date"))
        highest_targets.append(target_label(best_row))

    out["최고매출"] = highest_amounts
    out["최고일자"] = [
        pd.to_datetime(value).strftime("%Y-%m-%d")
        if pd.notna(pd.to_datetime(value, errors="coerce"))
        else ""
        for value in highest_dates
    ]
    out["최고타겟"] = highest_targets
    return out


def promotion_label(row: pd.Series) -> str:
    """프로모션 관련 컬럼의 실제 운영값만 표시합니다."""
    non_promo = {"", "-", "0", "x", "n", "no", "미진행", "일반", "일반기간", "해당없음", "없음", "nan", "none"}
    for col in ["프로모션명", "프로모션", "보답프로그램", "행사명", "기획전명"]:
        value = str(row.get(col, "")).strip()
        if value.lower() not in non_promo:
            return value
    return "-"


def is_promotional(row: pd.Series) -> bool:
    return promotion_label(row) != "-"


def coefficient_of_variation(values: pd.Series) -> float:
    values = pd.to_numeric(values, errors="coerce").dropna()
    if len(values) < 2 or values.mean() == 0:
        return 0.0
    return float(values.std(ddof=0) / abs(values.mean()))


def linear_trend_rate(values: pd.Series) -> float:
    """운영 순서 기준 단순 추세 기울기를 평균 대비 비율로 반환합니다."""
    vals = pd.to_numeric(values, errors="coerce").dropna().astype(float)
    if len(vals) < 3 or vals.mean() == 0:
        return 0.0
    x = pd.Series(range(len(vals)), dtype=float)
    x_mean, y_mean = x.mean(), vals.mean()
    denom = float(((x - x_mean) ** 2).sum())
    if denom == 0:
        return 0.0
    slope = float(((x - x_mean) * (vals.reset_index(drop=True) - y_mean)).sum() / denom)
    return slope / abs(y_mean)


def insight_confidence(sample_size: int) -> str:
    if sample_size >= 5:
        return "높음"
    if sample_size >= 3:
        return "보통"
    return "참고"


def make_product_history_table(row: pd.Series, history: pd.DataFrame, limit: int = 10) -> pd.DataFrame:
    """현재 건을 포함한 동일 상품 최근 발송 이력을 화면 표시용으로 생성합니다."""
    hist = product_history_including_current(row, history).copy()
    if hist.empty:
        return pd.DataFrame()
    hist = hist.sort_values("_date", ascending=False).head(limit).copy()
    hist["발송일"] = hist["_date"].dt.strftime("%Y-%m-%d")
    hist["타겟"] = hist.apply(target_label, axis=1)
    hist["프로모션"] = hist.apply(promotion_label, axis=1)
    if "발송일 최저가" in hist.columns:
        hist["발송일 최저가 여부"] = hist.apply(
            lambda r: (
                "O" if float(r.get("발송일 최저가", 0) or 0) > 0
                and float(r.get("멤버십혜택가", 0) or 0) <= float(r.get("발송일 최저가", 0) or 0)
                else ("X" if float(r.get("발송일 최저가", 0) or 0) > 0 else "-")
            ), axis=1,
        )
    else:
        hist["발송일 최저가 여부"] = "-"
    cols = ["발송일", "타겟", "주문금액"]
    if "멤버십혜택가" in hist.columns:
        cols.append("멤버십혜택가")
    cols += ["발송일 최저가 여부", "프로모션"]
    view = hist[[c for c in cols if c in hist.columns]].copy()
    for col in ["주문금액", "멤버십혜택가"]:
        if col in view.columns:
            view[col] = view[col].map(format_integer_price)
    return view.reset_index(drop=True)


def issue_storage_key(row: pd.Series) -> str:
    date_value = pd.to_datetime(row.get("_date"), errors="coerce")
    date_text = date_value.strftime("%Y-%m-%d") if pd.notna(date_value) else "no-date"
    product_key = clean_identifier_value(row.get("쇼라코드", "")) or clean_identifier_value(row.get("알파코드", "")) or str(row.get("상품명", "")).strip()
    target = target_label(row)
    return f"{date_text}|{product_key}|{target}"


def get_saved_issue(row: pd.Series) -> dict:
    issues = st.session_state.setdefault("daily_operation_issues", {})
    return issues.get(issue_storage_key(row), {})


def save_operation_issue(row: pd.Series, issue_type: str, memo: str) -> None:
    issues = st.session_state.setdefault("daily_operation_issues", {})
    key = issue_storage_key(row)
    normalized_type = "" if issue_type == "선택 안 함" else issue_type
    if normalized_type or memo.strip():
        issues[key] = {
            "유형": [normalized_type] if normalized_type else [],
            "메모": memo.strip(),
            "저장일시": datetime.now().strftime("%Y-%m-%d %H:%M"),
        }
    elif key in issues:
        del issues[key]


def delete_operation_issue(row: pd.Series) -> None:
    st.session_state.setdefault("daily_operation_issues", {}).pop(issue_storage_key(row), None)




def _daily_finalize_insight_lines(lines, row=None, history=None):
    """
    일일실적 최종 정리.
    - 신규/첫 운영이면 '신규 첫 운영'을 첫 줄에 배치
    - 동일 의미의 '금번 성과' 중복 제거
    - 과거 이력이 있는 상품은 고성과 회차와 금번 조건을 비교해
      가격/타겟/운영간격/프로모션 등 확인 가능한 요인만 '성과 하락 요인 추정'으로 제시
    - 원인 근거가 없으면 억지 추정하지 않음
    - 다음 운영 제안은 마지막 유지
    """
    if not lines:
        return _daily_finalize_insight_lines(lines, row=row, history=history)

    out = [str(x).strip() for x in lines if str(x).strip()]

    # 신규/첫 운영 문장 우선 배치
    first_ops = [x for x in out if ("신규 첫 운영" in x or "첫 운영" in x and ("신규" in x or "첫 TEST" in x))]
    if first_ops:
        first = first_ops[0]
        out = [x for x in out if x != first]
        out.insert(0, first)

    # 금번 성과 중복 제거: 첫 번째 핵심 성과만 유지
    seen_perf = False
    dedup = []
    for x in out:
        if "금번 성과" in x:
            if seen_perf:
                continue
            seen_perf = True
        dedup.append(x)
    out = dedup

    # 신규 첫 운영 문장이 이미 성과를 포함하면 별도 금번 성과 제거
    if out and ("신규 첫 운영" in out[0] or "첫 TEST" in out[0]):
        out = [out[0]] + [x for x in out[1:] if "금번 성과" not in x]

    # 과거 이력 기반 원인 추정
    cause_line = None
    try:
        if history is not None and len(history) >= 2 and row is not None:
            h = history.copy()

            # 컬럼 후보
            amount_col = "주문금액" if "주문금액" in h.columns else None
            price_col = "멤버십 혜택가" if "멤버십 혜택가" in h.columns else ("혜택가" if "혜택가" in h.columns else None)
            target_col = "타겟" if "타겟" in h.columns else None
            date_col = "_date" if "_date" in h.columns else ("발송일" if "발송일" in h.columns else None)
            promo_col = "프로모션" if "프로모션" in h.columns else None

            if amount_col:
                hh = h.copy()
                hh[amount_col] = pd.to_numeric(hh[amount_col], errors="coerce").fillna(0)
                current_amt = float(pd.to_numeric(pd.Series([row.get(amount_col, 0)]), errors="coerce").fillna(0).iloc[0])
                past = hh[hh[amount_col].notna()].copy()

                # 현재 행과 동일 레코드가 포함될 수 있으므로 날짜/금액 기준으로 가능한 범위에서 제외
                if date_col and date_col in past.columns and row.get(date_col) is not None:
                    try:
                        curd = pd.to_datetime(row.get(date_col), errors="coerce")
                        pdts = pd.to_datetime(past[date_col], errors="coerce")
                        past = past[pdts < curd]
                    except Exception:
                        pass

                if not past.empty:
                    best = past.sort_values(amount_col, ascending=False).iloc[0]
                    best_amt = float(best.get(amount_col, 0) or 0)

                    # 의미 있는 하락일 때만 원인 추정
                    if best_amt > 0 and current_amt < best_amt * 0.7:
                        causes = []

                        # 가격 변화
                        if price_col:
                            cur_price = pd.to_numeric(pd.Series([row.get(price_col)]), errors="coerce").iloc[0]
                            best_price = pd.to_numeric(pd.Series([best.get(price_col)]), errors="coerce").iloc[0]
                            if pd.notna(cur_price) and pd.notna(best_price) and cur_price > best_price:
                                diff = int(round(cur_price - best_price))
                                causes.append(f"고성과 당시 대비 혜택가 {diff:,}원 상승")

                        # 타겟 변화
                        if target_col:
                            cur_t = str(row.get(target_col, "") or "").strip()
                            best_t = str(best.get(target_col, "") or "").strip()
                            if cur_t and best_t and cur_t != best_t:
                                causes.append(f"고성과 타겟 {best_t} 대비 금번 {cur_t}로 변경")

                        # 프로모션 변화
                        if promo_col:
                            cur_p = str(row.get(promo_col, "") or "").strip()
                            best_p = str(best.get(promo_col, "") or "").strip()
                            if cur_p and best_p and cur_p != best_p:
                                causes.append("고성과 당시와 프로모션 조건 상이")

                        # 최근 동일 타겟에서도 부진했는지 확인 -> 타겟만의 문제로 단정 방지
                        same_target_weak = False
                        if target_col:
                            cur_t = str(row.get(target_col, "") or "").strip()
                            if cur_t:
                                same = past[past[target_col].astype(str) == cur_t].sort_values(date_col if date_col else amount_col)
                                if not same.empty:
                                    recent_same = float(same.iloc[-1][amount_col])
                                    if recent_same < 1_000_000:
                                        same_target_weak = True
                                        causes.append("최근 동일 타겟에서도 100만원 미만으로 상품 반응 약화")

                        if causes:
                            cause_text = "·".join(causes[:3])
                            cause_line = (
                                f"• 성과 하락 요인 추정 : 과거 최고 {best_amt/1_000_000:.1f}백만원 대비 금번 성과 하락 > "
                                f"{cause_text} 등 복합 영향 가능성 점검 필요"
                            )
                        else:
                            cause_line = (
                                f"• 성과 하락 요인 점검 : 과거 최고 {best_amt/1_000_000:.1f}백만원 대비 금번 성과가 낮으나 "
                                f"현재 데이터에서 가격·타겟 등 주요 조건의 뚜렷한 차이 확인 어려움 > "
                                f"상품 피로도·노출 조건·시즌 수요 등 추가 확인 필요"
                            )
    except Exception:
        cause_line = None

    if cause_line and not any(("성과 하락 요인" in x or "성과 하락 원인" in x) for x in out):
        # 다음 운영 제안 직전에 삽입
        pos = next((i for i, x in enumerate(out) if "다음 운영 제안" in x), len(out))
        out.insert(pos, cause_line)

    # 다음 운영 제안은 마지막으로
    actions = [x for x in out if "다음 운영 제안" in x]
    out = [x for x in out if "다음 운영 제안" not in x]
    if actions:
        out.append(actions[0])

    return _daily_select_insights(out, max_core=3)


def _daily_marketing_season_context(product_name: str, current_date) -> dict:
    """상품명 + 운영월 기준 일일실적용 시즌/마케팅 캘린더 맥락."""
    name = str(product_name or "").lower()
    dt = pd.to_datetime(current_date, errors="coerce")
    month = int(dt.month) if pd.notna(dt) else 0
    rules = [
        ({6,7,8}, ["냉감","쿨링","듀라론","인견","아이스","쿨매트","냉감패드","냉감이불","베개커버"],
         "폭염·열대야 수요가 집중되는 여름 냉감 침구 시즌", "냉감 소재·세탁 편의성·통기성·간편 교체 등 체감 기능이 명확한 상품"),
        ({5,6,7,8}, ["선풍기","써큘","서큘","에어컨","냉풍기"],
         "기온 상승과 폭염 대비 수요가 확대되는 여름 냉방가전 시즌", "저소음·리모컨·공기순환·공간효율 등 사용 편의성이 강화된 상품"),
        ({6,7,8}, ["우산","우양산","양산"], "장마와 폭염 수요가 동시에 발생하는 여름 시즌", "경량·암막·휴대성·자동 기능 등 계절 사용성이 명확한 상품"),
        ({6,7,8}, ["선크림","선케어","자외선","썬크림"], "자외선 노출과 야외활동이 증가하는 여름 선케어 시즌", "휴대성·사용감·자외선 차단 기능이 명확한 상품"),
        ({7,8,9}, ["갈치","생선","수산","전복"], "휴가철·보양식·가정식 수요가 함께 움직이는 여름 식품 시즌", "손질 편의성·중량·구성 혜택이 명확한 상품"),
        ({9,10,11}, ["홍삼","영양제","유산균"], "환절기 건강관리와 명절 선물 수요가 확대되는 가을 시즌", "복용 편의성·구성 혜택·선물 적합성이 명확한 상품"),
        ({11,12,1,2}, ["온열","전기요","히터","난방","기모"], "기온 하락으로 보온·난방 수요가 집중되는 겨울 시즌", "보온성·안전성·전력 효율·사용 편의성이 명확한 상품"),
    ]
    for months, keywords, context, attributes in rules:
        if month in months and any(k in name for k in keywords):
            return {"context": context, "attributes": attributes}
    return {}


def _daily_price_competitiveness(current_price: float, lowest: float) -> dict:
    """최저가 대비 차이를 차이율로 판정."""
    if current_price <= 0 or lowest <= 0:
        return {}
    advantage = (lowest-current_price)/lowest
    if advantage >= .05: return {"level":"strong","rate":advantage}
    if advantage >= .01: return {"level":"moderate","rate":advantage}
    if advantage >= -.01: return {"level":"same","rate":advantage}
    return {"level":"weak","rate":advantage}


def generate_insight_report(row: pd.Series, history: pd.DataFrame, issue: dict | None = None) -> dict:
    """상품별 핵심 인사이트를 생성합니다. 운영 이슈가 있으면 성과 판단보다 우선 반영합니다."""
    name = str(row.get("상품명", "")).strip()
    amount = float(row.get("주문금액", 0) or 0)
    current_target = target_label(row)
    current_date = pd.to_datetime(row.get("_date"), errors="coerce")
    current_price = float(row.get("멤버십혜택가", 0) or 0)
    grade = product_grade(amount)
    season_ctx = _daily_marketing_season_context(name, current_date)
    summary = product_history_summary(row, history)
    prior = summary["과거이력"].copy().sort_values("_date")
    same_target = summary["동일타겟이력"].copy()
    cumulative = product_history_including_current(row, history).copy().sort_values("_date")
    insights: list[tuple[int, str, str, str, str]] = []
    risks: list[str] = []
    issue = issue or {}
    issue_types = set(issue.get("유형", []))
    critical_issue = bool(issue_types.intersection({"판매중단", "가격오류"}))

    def add(priority: int, category: str, sentence: str, evidence: str = "", confidence: str = "보통"):
        if sentence and sentence not in [item[2] for item in insights]:
            insights.append((priority, category, sentence, evidence, confidence))

    if issue_types:
        issue_text = "·".join(sorted(issue_types))
        detail = issue.get("메모", "")
        sentence = f"금번 운영에서 {issue_text} 이슈가 등록되어 주문금액만으로 정상적인 상품 반응을 판단하기 어렵습니다."
        if detail:
            sentence += f" ({detail})"
        add(120, "운영 이슈", sentence, "운영 이슈 등록", "높음")

    # 현재 주문금액 등급에 따른 기본 평가
    # 같은 의미라도 상품·타겟 조건별로 표현을 달리해 문장 반복을 줄입니다.
    sentence_key = f"{name}|{current_target}|{grade}|{int(amount)}"
    if not critical_issue:
        if amount < 1_000_000:
            sentence = stable_variant(sentence_key, [
                f"금번 {compact_money(amount)}으로 100만원 미만의 부진 상품 수준을 기록해 기대 대비 매우 저조한 실적입니다.",
                f"금번 주문금액은 {compact_money(amount)}으로, MMS 메인 상품으로 활용하기에는 반응이 제한적이었습니다.",
                f"금번 {compact_money(amount)}에 그치며 가격·타겟 조건 대비 구매 반응이 충분히 확보되지 않았습니다.",
                f"금번 주문금액이 {compact_money(amount)}으로 100만원 미만에 머물러 상품 적합도 재검토가 필요합니다.",
            ])
            add(97, "금번 성과", sentence, "현재 주문금액 기준", "높음")
        elif amount < 2_000_000:
            sentence = stable_variant(sentence_key, [
                f"금번 {compact_money(amount)}으로 관찰 상품 수준을 기록해 기대 대비 다소 아쉬운 실적입니다.",
                f"금번 주문금액은 {compact_money(amount)}으로 목표 수준에는 다소 미치지 못했습니다.",
                f"금번 {compact_money(amount)}을 기록해 기본 수요는 확인했으나 메인 상품 성과로는 다소 제한적이었습니다.",
                f"금번 주문금액이 {compact_money(amount)}으로 200만원 미만에 머물러 추가 조건 검증이 필요합니다.",
            ])
            add(91, "금번 성과", sentence, "현재 주문금액 기준", "높음")
        elif amount < 3_000_000:
            sentence = stable_variant(sentence_key, [
                f"금번 {compact_money(amount)}으로 안정 상품 수준을 기록해 상품당 목표 250만원에 근접한 성과입니다.",
                f"금번 주문금액은 {compact_money(amount)}으로 목표 수준 전후의 안정적인 성과를 확보했습니다.",
                f"금번 {compact_money(amount)}을 기록해 추가 운영 여부를 판단할 수 있는 기본 성과를 확보했습니다.",
            ])
            add(75, "금번 성과", sentence, "현재 주문금액 기준", "높음")
        elif amount < 5_000_000:
            sentence = stable_variant(sentence_key, [
                f"금번 {compact_money(amount)}으로 우수 상품 수준의 성과를 확보했습니다.",
                f"금번 주문금액은 {compact_money(amount)}으로 목표를 상회하는 양호한 판매 성과를 기록했습니다.",
                f"금번 {compact_money(amount)}을 기록해 재편성 검토가 가능한 우수 성과를 확인했습니다.",
            ])
            add(80, "금번 성과", sentence, "현재 주문금액 기준", "높음")
        else:
            sentence = stable_variant(sentence_key, [
                f"금번 {compact_money(amount)}으로 핵심 상품 수준의 매우 우수한 성과를 기록했습니다.",
                f"금번 주문금액은 {compact_money(amount)}으로 핵심 매출 견인 상품 수준의 성과를 확보했습니다.",
                f"금번 {compact_money(amount)}을 기록하며 차주 핵심 편성 후보로 검토할 수 있는 성과를 확인했습니다.",
            ])
            add(90, "금번 성과", sentence, "현재 주문금액 기준", "높음")

    # 성과 및 추세
    if not prior.empty:
        past_max, past_avg = float(prior["주문금액"].max()), float(prior["주문금액"].mean())
        if amount > past_max and amount >= 3_000_000:
            add(100, "성과", f"금번 {current_target or '운영 타겟'}에서 {compact_money(amount)}을 기록하며 역대 최고 실적을 경신했습니다.", f"과거 최고 {compact_money(past_max)}", insight_confidence(len(prior)))
        elif past_avg > 0 and amount >= past_avg * 1.5 and amount >= 3_000_000:
            add(90, "성과", f"과거 평균 대비 주문금액이 {((amount/past_avg)-1)*100:.0f}% 증가해 거래액이 크게 성장했습니다.", f"과거 평균 {compact_money(past_avg)}", insight_confidence(len(prior)))
        elif past_avg > 0 and amount <= past_avg * 0.7:
            decline_pct = (1 - amount / past_avg) * 100
            decline_sentence = stable_variant(sentence_key + "|decline", [
                f"금번 주문금액이 과거 평균 대비 {decline_pct:.0f}% 낮아 성과 둔화가 확인됩니다.",
                f"과거 평균 대비 주문금액이 {decline_pct:.0f}% 감소해 최근 판매 흐름이 약화되었습니다.",
                f"금번 실적은 과거 평균의 {amount/past_avg*100:.0f}% 수준으로, 이전 운영 대비 반응이 제한적이었습니다.",
                f"과거 평균 대비 {decline_pct:.0f}% 낮은 성과를 기록해 운영 조건 재점검이 필요합니다.",
            ])
            add(88, "성과", decline_sentence, f"과거 평균 {compact_money(past_avg)}", insight_confidence(len(prior)))
            risks.append("최근 성과 둔화")

    recent3 = cumulative.tail(3)
    if len(recent3) == 3:
        vals = recent3["주문금액"].astype(float).tolist()
        growth = vals[2] / vals[0] - 1 if vals[0] > 0 else 0
        if vals[0] < vals[1] < vals[2] and growth >= 0.15:
            add(95, "성장 추세", f"최근 3회 주문금액이 {compact_money(vals[0])} → {compact_money(vals[1])} → {compact_money(vals[2])}으로 연속 성장했습니다.", f"첫 회 대비 {growth*100:.0f}% 증가", "보통")
        elif vals[0] > vals[1] > vals[2] and vals[0] > 0 and vals[2] <= vals[0] * 0.8:
            add(89, "성장 추세", f"최근 3회 주문금액이 연속 감소해 운영 조건 재점검이 필요합니다.", f"첫 회 대비 {(1-vals[2]/vals[0])*100:.0f}% 감소", "보통")
            risks.append("최근 3회 연속 하락")
        elif min(vals) >= 3_000_000 and coefficient_of_variation(recent3["주문금액"]) <= 0.25:
            add(82, "운영 안정성", "최근 3회 모두 300만원 이상을 기록하고 실적 편차가 제한적이어서 안정적인 판매 흐름이 확인됩니다.", f"변동계수 {coefficient_of_variation(recent3['주문금액']):.2f}", "보통")

    if len(cumulative) >= 5:
        recent5 = cumulative.tail(5)
        cv = coefficient_of_variation(recent5["주문금액"])
        if float(recent5["주문금액"].mean()) >= 3_000_000 and cv <= 0.30:
            add(87, "운영 안정성", f"최근 5회 평균 {compact_money(recent5['주문금액'].mean())}을 기록하고 변동성이 낮아 반복 검증된 안정형 상품입니다.", f"5회 변동계수 {cv:.2f}", "높음")
        elif cv >= 0.75:
            add(70, "운영 위험", "회차별 주문금액 편차가 커 편성 조건에 따른 성과 변동성이 높은 상품입니다.", f"5회 변동계수 {cv:.2f}", "높음")
            risks.append("성과 변동성 높음")

    # 동일 주차 중복 제거 후 편성 횟수
    if "주차" in cumulative.columns and "주차" in row.index:
        week_value = str(row.get("주차", "")).strip()
        week_rows = cumulative[cumulative["주차"].astype(str).eq(week_value)].copy()
        if not week_rows.empty:
            unique_keys = [c for c in ["_date", "시간대", "캠페인명", "소재", "성별", "연령", "SEG"] if c in week_rows.columns]
            week_unique = week_rows.drop_duplicates(unique_keys) if unique_keys else week_rows
            if len(week_unique) >= 2 and (week_unique["주문금액"] >= 3_000_000).all():
                add(96, "성과", f"금주 총 {len(week_unique)}회 편성되며 모든 운영에서 300만원 이상의 주문금액을 기록했습니다.", "주차 내 고유 발송 기준", "높음")

    # 타겟 적합도 및 확장성
    if not prior.empty:
        target_df = prior.assign(
            _target=prior.apply(target_label, axis=1),
            _base_target=prior.apply(base_target_label, axis=1),
        )
        target_stats = target_df.groupby("_target")["주문금액"].agg(["mean", "sum", "count"]).sort_values("mean", ascending=False)
        overall_avg = float(prior["주문금액"].mean())
        current_base_target = base_target_label(row)
        if len(same_target) >= 2 and float(same_target["주문금액"].mean()) >= overall_avg * 1.15:
            add(93, "타겟 적합도", f"{current_target} 과거 평균은 {compact_money(same_target['주문금액'].mean())}으로 전체 평균보다 높아 핵심 타겟 적합도가 확인됩니다.", f"동일 타겟 {len(same_target)}회", insight_confidence(len(same_target)))
        elif same_target.empty and amount >= 3_000_000:
            same_base_prior = target_df[target_df["_base_target"].eq(current_base_target)] if current_base_target else target_df.iloc[0:0]
            if not same_base_prior.empty:
                seg_value = clean_identifier_value(row.get("SEG", ""))
                seg_text = f"SEG{seg_value}" if seg_value else "신규 SEG"
                add(86, "타겟 확장성", f"{current_base_target} 내 {seg_text} 첫 TEST에서 {compact_money(amount)}을 기록해 동일 타겟군 내 미발송 SEG 확장 가능성이 확인됩니다.", f"동일 성별·연령 과거 운영 {len(same_base_prior)}회 / 해당 SEG 첫 운영", "참고")
            else:
                add(86, "타겟 확장성", f"{current_target} 첫 TEST에서 {compact_money(amount)}을 기록해 신규 타겟 확장 가능성이 확인됩니다.", "신규 타겟 1회", "참고")
        if len(target_stats) >= 2 and target_stats["sum"].sum() > 0:
            top = target_stats.iloc[0]
            second = target_stats.iloc[1]
            if top["count"] >= 2 and second["count"] >= 2 and top["mean"] >= second["mean"] * 1.5:
                add(92, "타겟 적합도", f"{target_stats.index[0]} 회당 평균이 {compact_money(top['mean'])}으로 차순위 타겟보다 압도적으로 높아 타겟 확장보다 핵심 타겟 집중 운영이 효율적입니다.", f"차순위 평균 {compact_money(second['mean'])}", "높음")
            top_share = float(target_stats.iloc[0]["sum"] / target_stats["sum"].sum())
            if len(prior) >= 3 and int(top["count"]) >= 2 and top_share >= 0.7:
                add(67, "타겟 위험", f"누적 주문금액의 {top_share*100:.0f}%가 {target_stats.index[0]}에 집중되어 타겟 편중 여부를 함께 관리해야 합니다.", f"동일 상품 과거 {len(prior)}회 / 해당 타겟 {int(top['count'])}회", insight_confidence(int(top["count"])))
                risks.append("특정 타겟 편중")

    # 프로모션 의존도: 발송일 기준 실제 프로모션명과 일반기간을 비교
    if "프로모션명" in cumulative.columns and len(cumulative) >= 4:
        promo_mask = cumulative["프로모션명"].fillna("-").astype(str).ne("-")
        promo_rows, normal_rows = cumulative[promo_mask], cumulative[~promo_mask]
        if len(promo_rows) >= 2 and len(normal_rows) >= 2:
            promo_avg = float(promo_rows["주문금액"].mean())
            normal_avg = float(normal_rows["주문금액"].mean())
            current_promo = promotion_label(row)
            promo_name = current_promo if current_promo != "-" else str(promo_rows["프로모션명"].mode().iloc[0])
            if normal_avg > 0 and promo_avg >= normal_avg * 1.3:
                add(84, "프로모션", f"{promo_name} 평균 {compact_money(promo_avg)} 대비 일반 운영 기간 평균은 {compact_money(normal_avg)}으로 프로모션 의존도가 높은 상품입니다.", f"프로모션 {len(promo_rows)}회 / 일반 {len(normal_rows)}회", "보통")
                risks.append("프로모션 의존")
            elif promo_avg > 0 and normal_avg >= promo_avg * 0.85:
                add(76, "프로모션", f"일반 운영 기간에도 평균 {compact_money(normal_avg)}을 기록해 프로모션 의존도가 낮은 상품입니다.", f"프로모션 {len(promo_rows)}회 / 일반 {len(normal_rows)}회", "높음" if len(normal_rows) >= 3 else "보통")

    # 가격 경쟁력 및 탄력성: 비율보다 고객 체감 차액을 우선 표시하고 성과와 교차 해석
    lowest = float(row.get("발송일 최저가", 0) or 0)
    price_eval = _daily_price_competitiveness(current_price, lowest)
    if price_eval:
        price_diff = lowest - current_price
        if price_eval["level"] == "strong":
            if amount < 2_000_000:
                add(90, "가격·성과", f"멤버십 혜택가 {fmt_num(current_price)}원으로 발송일 비교 최저가보다 {fmt_num(price_diff)}원 저렴한 가격 경쟁력을 확보했음에도 금번 {compact_money(amount)}에 그쳐, 가격 외 상품·타겟 적합도 점검이 필요합니다.", "발송일 최저가 및 현재 주문금액 기준", "높음")
            elif not prior.empty and float(prior["주문금액"].mean()) > 0 and amount <= float(prior["주문금액"].mean()) * 0.7:
                past_avg_price_cross = float(prior["주문금액"].mean())
                add(90, "가격·성과", f"발송일 비교 최저가보다 {fmt_num(price_diff)}원 저렴한 가격 경쟁력을 확보했음에도 금번 {compact_money(amount)}으로 과거 평균 {compact_money(past_avg_price_cross)} 대비 성과가 낮아, 추가 할인보다 운영 간격·타겟 적합도 점검이 우선입니다.", "가격 경쟁력 + 과거 평균 성과 교차 기준", "높음")
            else:
                add(88, "가격", f"멤버십 혜택가 {fmt_num(current_price)}원으로 발송일 비교 최저가보다 {fmt_num(price_diff)}원 저렴해 높은 가격 경쟁력을 확보했습니다.", "발송일 최저가 기준", "높음")
        elif price_eval["level"] == "moderate":
            add(82, "가격", f"멤버십 혜택가 {fmt_num(current_price)}원은 발송일 비교 최저가보다 {fmt_num(price_diff)}원 저렴해 가격 우위는 있으나 차별화 폭은 제한적입니다.", "발송일 최저가 기준", "높음")
        elif price_eval["level"] == "same":
            diff_abs = abs(price_diff)
            add(84 if amount < 2_000_000 else 72, "가격", f"멤버십 혜택가 {fmt_num(current_price)}원은 발송일 비교 최저가와 {fmt_num(diff_abs)}원 차이로 사실상 동일한 수준이어서 가격 차별화는 제한적입니다.", "발송일 최저가 ±1% 이내", "높음")
            if amount < 2_000_000:
                risks.append("가격 차별화 제한")
        elif amount < 2_000_000:
            add(80, "가격 위험", f"멤버십 혜택가 {fmt_num(current_price)}원은 발송일 비교 최저가보다 {fmt_num(abs(price_diff))}원 높고 성과도 제한적이어서 가격 조건 재점검이 필요합니다.", "발송일 최저가 대비 가격 열위", "높음")
            risks.append("가격 경쟁력 미확보")
    if not prior.empty and current_price > 0:
        last = prior.iloc[-1]
        last_price, last_amount = float(last.get("멤버십혜택가", 0) or 0), float(last.get("주문금액", 0) or 0)
        if last_price > 0 and last_amount > 0 and current_price > last_price and amount >= last_amount * 0.9:
            add(85, "가격 탄력성", f"직전 대비 혜택가가 {fmt_num(current_price-last_price)}원 상승했으나 주문금액은 직전의 {amount/last_amount*100:.0f}% 수준을 유지해 가격 민감도가 낮은 흐름입니다.", "직전 운영 비교", "보통")
        price_rows = cumulative[(pd.to_numeric(cumulative.get("멤버십혜택가", 0), errors="coerce") > 0) & (cumulative["주문금액"] > 0)] if "멤버십혜택가" in cumulative.columns else pd.DataFrame()
        if len(price_rows) >= 4 and price_rows["멤버십혜택가"].nunique() >= 2:
            corr = price_rows[["멤버십혜택가", "주문금액"]].corr().iloc[0,1]
            if pd.notna(corr) and corr <= -0.6:
                add(74, "가격 탄력성", "가격 상승 구간에서 주문금액이 함께 하락하는 경향이 뚜렷해 가격 민감형 상품으로 판단됩니다.", f"가격-매출 상관계수 {corr:.2f}", "보통")
                risks.append("가격 민감형")

    # 피로도·희소성
    if not prior.empty and pd.notna(current_date):
        last = prior.iloc[-1]
        gap, last_amount = int((current_date-last["_date"]).days), float(last.get("주문금액", 0) or 0)
        if gap <= 21 and last_amount > 0:
            ratio = amount / last_amount
            past_avg_for_fatigue = float(prior["주문금액"].mean()) if not prior.empty else 0
            if ratio >= 0.8 and amount >= 2_000_000 and (past_avg_for_fatigue <= 0 or amount >= past_avg_for_fatigue * 0.7):
                add(80, "운영 피로도", f"직전 운영 후 {gap}일 만에 재편성했음에도 주문금액이 직전의 {ratio*100:.0f}% 수준을 유지해 단기 반복에 따른 추가 하락은 제한적입니다.", "직전 운영 비교", "보통")
            elif ratio < 0.75:
                add(83, "운영 위험", f"직전 운영 후 {gap}일 만의 재편성에서 주문금액이 직전 대비 {(1-ratio)*100:.0f}% 감소해 미편성 기간 부여가 필요합니다.", "직전 운영 비교", "보통")
                risks.append("단기 반복 피로도")
        elif gap >= 45 and amount >= 3_000_000:
            add(81, "운영 희소성", f"직전 운영 후 {gap}일 만의 재편성에서 {compact_money(amount)}을 기록해 장기간 미운영 후에도 우수한 반응이 확인됩니다.", "재편성 간격 기준", "보통")
    recent90 = cumulative[cumulative["_date"] >= current_date-pd.Timedelta(days=90)] if pd.notna(current_date) else cumulative
    if summary["운영횟수"] >= 1 and 2 <= len(recent90) <= 3 and float(recent90["주문금액"].mean()) >= 3_000_000:
        add(78, "운영 희소성", f"최근 3개월간 {len(recent90)}회 제한적으로 운영했음에도 평균 {compact_money(recent90['주문금액'].mean())}을 기록해 추가 운영 여력이 있습니다.", "최근 90일 기준", "보통")

    # 시즌·생애주기·포지션
    season_words = ["선풍기", "에어컨", "서큘레이터", "우양산", "래쉬가드", "삼계탕", "장어", "아이스크림", "제습기", "냉감"]
    if any(word in name for word in season_words) and grade in ["핵심 상품", "우수 상품"]:
        add(71, "시즌", "시즌 수요가 반영된 우수 성과로 수요가 유지되는 기간 내 추가 운영을 검토할 수 있습니다.", "상품명 시즌 키워드 기준", "참고")
    avg_all = float(cumulative["주문금액"].mean()) if not cumulative.empty else amount
    if len(cumulative) >= 5:
        trend = linear_trend_rate(cumulative.tail(8)["주문금액"])
        if trend >= 0.08:
            add(79, "생애주기", "중기 추세가 상승하는 성장기 상품으로 운영 비중 확대 검토가 가능합니다.", f"최근 최대 8회 추세율 {trend:.2f}", "보통")
        elif trend <= -0.08:
            add(77, "생애주기", "중기 추세가 하락하는 성숙·하락 전환 구간으로 운영 조건 재설계가 필요합니다.", f"최근 최대 8회 추세율 {trend:.2f}", "보통")
            risks.append("생애주기 하락 전환")
    if len(cumulative) >= 5 and avg_all >= 5_000_000:
        add(91, "상품 포지션", f"누적 {len(cumulative)}회 평균 {compact_money(avg_all)}을 기록한 반복 검증형 대표 매출 견인 상품입니다.", f"누적 {len(cumulative)}회", "높음")
    elif summary["운영횟수"] == 0 and amount >= 3_000_000:
        add(83, "상품 포지션", "첫 운영에서 우수한 주문금액을 기록한 신규 성장 상품으로 추가 TEST가 필요합니다.", "신규 1회", "참고")

    # MMS 메인 상품 적합도는 반복 부진 근거가 충분할 때만 강하게 판단
    if not critical_issue and len(cumulative) >= 3:
        recent_normal = cumulative[~cumulative.apply(is_promotional, axis=1)].tail(3)
        if len(recent_normal) >= 3 and float(recent_normal["주문금액"].mean()) < 1_000_000:
            add(108, "상품 적합도", f"최근 일반기간 3회 평균이 {compact_money(recent_normal['주문금액'].mean())}으로 반복 운영에서도 성과 개선이 제한적이어서 MMS 메인 상품으로는 적합도가 낮은 것으로 판단됩니다.", "일반기간 최근 3회", "높음")
            risks.append("MMS 메인 적합도 낮음")

    # 시즌성·마케팅 캘린더
    if season_ctx:
        add(86 if summary["운영횟수"] == 0 else 78, "시즌성",
            f"{season_ctx['context']}에 해당하는 상품입니다. {season_ctx['attributes']} 중심으로 상품 매력을 강화해 운영하는 것이 적절합니다.",
            f"{pd.to_datetime(current_date).month if pd.notna(current_date) else '-'}월 마케팅 캘린더 및 상품 속성 기준", "보통")

    # 다음 운영 제안: 분석 결과를 실제 편성 액션으로 연결합니다.
    current_promo_name = promotion_label(row)
    recent_gap = None
    if not prior.empty and pd.notna(current_date):
        recent_gap = int((current_date - prior.iloc[-1]["_date"]).days)

    if critical_issue:
        action_sentence = "운영 이슈 해소 후 동일 상품·동일 타겟으로 재TEST하여 정상 성과를 다시 확인하는 것이 필요합니다."
        action_evidence = "운영 이슈 영향으로 성과 판단 보류"
    elif "MMS 메인 적합도 낮음" in risks:
        action_sentence = "반복 부진이 확인된 만큼 MMS 메인 편성은 축소하고, 가격·구성 또는 타겟 조건이 개선된 경우에만 선택적으로 재TEST하는 것이 필요합니다."
        action_evidence = "일반기간 최근 3회 반복 부진"
    elif "가격 경쟁력 미확보" in risks:
        action_sentence = "현재 가격 조건에서는 재편성 우선순위를 낮추고, 발송일 비교 최저가 대비 유의미한 가격 경쟁력 또는 구성 혜택을 확보한 뒤 재TEST하는 것이 필요합니다."
        action_evidence = "발송일 최저가 대비 1% 초과 가격 열위"
    elif summary["운영횟수"] == 0 and amount < 1_000_000:
        if price_eval and price_eval.get("level") in {"same","moderate"}:
            action_sentence = "첫 운영 부진을 상품 적합도만의 문제로 단정하기 어렵습니다. 실질적인 가격 차별화 또는 구성 혜택을 보강하고 타겟을 변경해 1회 재TEST한 뒤 가격·타겟 요인을 구분해 판단하는 것이 필요합니다."
            action_evidence = "신규 첫 운영 + 100만원 미만 + 가격 차별화 제한"
        else:
            action_sentence = "첫 운영만으로 상품 적합도를 단정하기 어려우므로 가격·구성·타겟 중 최소 한 가지 조건을 보완해 1회 재TEST 후 판단하는 것이 필요합니다."
            action_evidence = "신규 첫 운영 100만원 미만"
    elif summary["운영횟수"] == 0 and amount >= 2_000_000 and season_ctx:
        action_sentence = f"신규 첫 운영에서 {compact_money(amount)}을 기록해 기본 판매 가능성을 확인했습니다. {season_ctx['context']}인 만큼 동일 타겟 또는 미발송 SEG로 1회 추가 TEST해 시즌 내 성과 확장 여부를 확인하는 것이 좋습니다."
        action_evidence = "신규 첫 운영 + 안정 이상 성과 + 시즌 적합성"
    elif "프로모션 의존" in risks:
        action_sentence = f"일반기간 확대보다 {current_promo_name if current_promo_name != '-' else '프로모션'} 기간 우선 편성하고, 일반기간 운영은 가격·구성 보강 시 선택적으로 검토하는 것이 필요합니다."
        action_evidence = "프로모션·일반기간 평균 비교"
    elif "최근 성과 둔화" in risks and amount < 2_000_000:
        action_sentence = "직전 회차 대비 유지 여부와 별개로 절대 매출과 과거 평균 대비 성과가 모두 낮은 만큼 단기 반복 편성은 제한하고, 과거 고성과 타겟·가격·시즌 조건을 재확인한 후 재운영하는 것이 적절합니다."
        action_evidence = "절대 성과 + 과거 평균 동시 부진"
    elif "단기 반복 피로도" in risks or (recent_gap is not None and recent_gap <= 21 and amount < 3_000_000):
        action_sentence = "단기 반복 편성은 줄이고 최소 2~3주 미편성 기간 후 기존 우수 타겟 중심으로 재편성하는 것이 필요합니다."
        action_evidence = "최근 운영 간격 및 성과 하락 기준"
    elif amount >= 5_000_000:
        if same_target.empty:
            current_base_target = base_target_label(row)
            same_base_prior = prior[prior.apply(base_target_label, axis=1).eq(current_base_target)] if (not prior.empty and current_base_target) else prior.iloc[0:0]
            if not same_base_prior.empty:
                seg_value = clean_identifier_value(row.get("SEG", ""))
                seg_text = f"SEG{seg_value}" if seg_value else "신규 SEG"
                action_sentence = f"금번 {current_base_target} 내 {seg_text}의 우수 반응을 바탕으로 동일 SEG 1회 추가 검증 후 다른 미발송 SEG 확대를 검토하는 것이 좋습니다."
                action_evidence = "동일 타겟군 내 신규 SEG 핵심 상품 성과"
            else:
                action_sentence = f"금번 {current_target or '신규 타겟'}의 우수 반응을 바탕으로 동일 타겟 1회 추가 검증 후 미발송 SEG 확대를 검토하는 것이 좋습니다."
                action_evidence = "신규 타겟 핵심 상품 성과"
        else:
            action_sentence = f"{current_target or '현재 우수 타겟'} 중심의 핵심 재편성 상품으로 유지하고, 운영 간격을 관리하며 미발송 SEG 확대를 검토하는 것이 좋습니다."
            action_evidence = "핵심 상품 및 타겟 실적 기준"
    elif amount >= 3_000_000:
        action_sentence = f"{current_target or '현재 타겟'}에 우선 재편성해 성과를 한 차례 더 확인하고, 유사 실적이 유지되면 운영 비중 확대를 검토하는 것이 좋습니다."
        action_evidence = "우수 상품 기준"
    elif amount >= 2_000_000:
        if price_eval and price_eval.get("level") == "strong":
            action_sentence = "가격 경쟁력을 확보한 상태에서 안정 수준의 성과를 기록한 만큼 동일 조건 또는 미발송 SEG로 1회 추가 TEST해 250만원 이상 성과 재현 여부를 확인하는 것이 좋습니다."
            action_evidence = "안정 상품 + 유의미한 가격 경쟁력"
        else:
            action_sentence = "동일 조건으로 1회 추가 TEST하되 가격·타겟 조건을 함께 점검하고 250만원 이상 성과가 재현되는지 확인한 뒤 확대 여부를 판단하는 것이 좋습니다."
            action_evidence = "안정 상품 및 목표 250만원 기준"
    elif amount >= 1_000_000:
        action_sentence = "가격·타겟·전시순서 중 한 가지 조건을 조정해 선택적으로 재TEST하고, 200만원 이상 회복 여부를 확인하는 것이 필요합니다."
        action_evidence = "관찰 상품 기준"
    else:
        action_sentence = "현재 조건의 반복 편성은 지양하고, 가격·구성·타겟 중 개선 가능한 조건을 먼저 확보한 뒤 재TEST 여부를 판단하는 것이 필요합니다."
        action_evidence = "부진 상품 기준"

    # 중요도·중복 제어: 분석 5개 + 다음 운영 제안 1개로 최대 6개를 유지합니다.
    # 역대 최고/연속 성장처럼 더 강한 성과 해석이 있으면 단순 "금번 성과" 문장은 중복 제거
    has_strong_performance_story = any(
        category in {"성과", "성장 추세"} and ("역대 최고" in sentence or "연속 성장" in sentence)
        for _, category, sentence, _, _ in insights
    )
    if has_strong_performance_story:
        insights = [item for item in insights if item[1] != "금번 성과"]
    insights = sorted(insights, key=lambda x: (-x[0], x[1]))
    selected, category_count = [], {}
    for _, category, sentence, evidence, confidence in insights:
        if any(item["sentence"] == sentence for item in selected) or category_count.get(category, 0) >= 1:
            continue
        item_type = "risk" if category in {"운영 위험", "가격 위험", "타겟 위험", "상품 적합도", "운영 이슈"} or "낮" in sentence or "저조" in sentence or "아쉬" in sentence else "fact"
        selected.append({"category": category, "sentence": sentence, "evidence": evidence, "confidence": confidence, "type": item_type})
        category_count[category] = 1
        if len(selected) >= 5:
            break
    if not selected:
        selected.append({"category": "운영", "sentence": f"금번 {current_target or '운영 타겟'}에서 {compact_money(amount)}을 기록했으며 추가 이력 축적 후 판단이 필요합니다.", "evidence": "현재 1회", "confidence": "참고", "type": "fact"})

    selected.append({
        "category": "다음 운영 제안",
        "sentence": f"다음 운영 제안: {action_sentence}",
        "evidence": action_evidence,
        "confidence": "높음" if critical_issue or summary["운영횟수"] >= 3 else "참고",
        "type": "action",
    })

    return {
        "상품명": name,
        "상품등급": grade,
        "인사이트": selected,
        "위험요인": sorted(set(risks)),
        "발송이력": make_product_history_table(row, history, limit=5),
        "운영이슈": issue,
    }



def format_daily_insight_item(item: dict) -> str:
    """일일실적 상품 인사이트를 간결한 실무 보고체로 표시합니다."""
    category = str(item.get("category", "") or "").strip()
    sentence = str(item.get("sentence", "") or "").strip()
    evidence = str(item.get("evidence", "") or "").strip()

    # 짧은 항목명
    if category == "다음 운영 제안" or sentence.startswith("다음 운영 제안:"):
        title = "다음 운영 제안"
        sentence = sentence.replace("다음 운영 제안:", "", 1).strip()
    elif "역대 최고" in sentence:
        title = "역대 최고 경신"
    elif category == "성장 추세":
        title = "성과 성장"
    elif category == "가격·성과":
        title = "가격 대비 성과 점검"
    elif "가격" in category:
        title = "가격 경쟁력" if "위험" not in category else "가격 조건 점검"
    elif category == "타겟 확장성":
        title = "타겟·SEG 확장"
    elif category == "타겟 적합도":
        title = "핵심 타겟 적합도"
    elif category == "타겟 위험":
        title = "타겟 편중 점검"
    elif category in {"운영 위험", "운영 이슈"}:
        title = "운영 조건 점검"
    elif category in {"금번 성과", "성과"}:
        title = "금번 성과"
    else:
        title = category or "운영 인사이트"

    # 서술형 종결 → 간결한 실무 보고체
    replacements = [
        ("기록하며 역대 최고 실적을 경신했습니다.", "기록, 역대 최고 실적 경신"),
        ("역대 최고 실적을 경신했습니다.", "역대 최고 실적 경신"),
        ("기록했습니다.", "기록"),
        ("확인되었습니다.", "확인"),
        ("확인됩니다.", "확인"),
        ("확보했습니다.", "확보"),
        ("유지했습니다.", "유지"),
        ("개선되었습니다.", "개선"),
        ("감소했습니다.", "감소"),
        ("증가했습니다.", "증가"),
        ("하락했습니다.", "하락"),
        ("상승했습니다.", "상승"),
        ("가능성이 확인됩니다.", "가능성 확인"),
        ("흐름이 확인됩니다.", "흐름 확인"),
        ("재편성이 가능합니다.", "재편성 가능"),
        ("검토하는 것이 좋습니다.", "검토 필요"),
        ("검토하는 것이 필요합니다.", "검토 필요"),
        ("검토할 필요가 있습니다.", "검토 필요"),
        ("확인하는 것이 필요합니다.", "확인 필요"),
        ("확인하는 것이 좋습니다.", "확인 필요"),
        ("판단하는 것이 필요합니다.", "판단 필요"),
        ("것이 적절합니다.", "적절"),
        ("것이 필요합니다.", "필요"),
        ("것이 좋습니다.", "검토 필요"),
        ("필요합니다.", "필요"),
        ("수준입니다.", "수준"),
        ("성과입니다.", "성과"),
    ]
    for old, new in replacements:
        sentence = sentence.replace(old, new)

    sentence = re.sub(r"(했습니다|됩니다|있습니다)\.$", "", sentence).rstrip(".").strip()

    # 괄호형 근거는 중복/기준성 문구를 줄이고 필요한 근거만 ' > '로 연결
    noisy_evidence = {
        "현재 주문금액 기준", "발송일 최저가 기준", "직전 운영 비교",
        "핵심 상품 기준", "우수 상품 기준", "관찰 상품 기준", "부진 상품 기준",
    }
    if evidence and evidence not in noisy_evidence and evidence not in sentence and len(evidence) <= 48:
        sentence = f"{sentence} > {evidence}"

    return f"• {title} : {sentence}"


def make_insight(row: pd.Series, history: pd.DataFrame) -> str:
    """상품구분·상품분석·PPT에서 사용할 한 줄형 호환 함수입니다."""
    report = generate_insight_report(row, history, get_saved_issue(row))
    sentences = [item["sentence"] for item in report["인사이트"]]
    return f"[{report['상품명']}] " + " > ".join(sentences)

def build_ppt(title: str, lines: list[str], table_df: pd.DataFrame | None = None) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(.6), Inches(.35), Inches(12), Inches(.65))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True

    body = slide.shapes.add_textbox(Inches(.7), Inches(1.2), Inches(12), Inches(5.7))
    tf = body.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(14)
        para.space_after = Pt(7)

    if table_df is not None and not table_df.empty:
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        head = slide2.shapes.add_textbox(Inches(.6), Inches(.3), Inches(12), Inches(.6))
        hp = head.text_frame.paragraphs[0]
        hp.text = "상세 실적"
        hp.font.size = Pt(24)
        hp.font.bold = True

        view = table_df.head(14).copy()
        rows, cols = len(view) + 1, len(view.columns)
        table = slide2.shapes.add_table(
            rows, cols, Inches(.35), Inches(1), Inches(12.6), Inches(5.9)
        ).table
        for j, c in enumerate(view.columns):
            table.cell(0, j).text = str(c)
        for i, (_, r) in enumerate(view.iterrows(), start=1):
            for j, c in enumerate(view.columns):
                table.cell(i, j).text = str(r[c])

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def append_total_and_change_rows(raw: pd.DataFrame, mode: str) -> pd.DataFrame:
    """상세 기간 행 아래에 총합계와 최신 기간 증감을 추가합니다."""
    if raw.empty:
        return raw.copy()

    d = raw.copy()
    total = {
        "_label": "총합계",
        "연도": "",
        "월": "",
        "발송횟수": d["발송횟수"].sum(),
        "상품수": d["상품수"].sum(),
        "URL": d["URL"].sum(),
        "발송건수": d["발송건수"].sum(),
        "클릭수": d["클릭수"].sum(),
        "주문건수": d["주문건수"].sum(),
        "주문수량": d["주문수량"].sum(),
        "주문금액": d["주문금액"].sum(),
    }
    total["반응율(Uniq CTR)"] = total["클릭수"] / total["발송건수"] if total["발송건수"] else 0
    total["객단가"] = total["주문금액"] / total["주문건수"] if total["주문건수"] else 0
    total["클릭 CVR"] = total["주문건수"] / total["클릭수"] if total["클릭수"] else 0
    total["발송 CVR"] = total["주문건수"] / total["발송건수"] if total["발송건수"] else 0
    total["클릭당매출(RPC)"] = total["주문금액"] / total["클릭수"] if total["클릭수"] else 0
    total["발송대비매출(SPM)"] = total["주문금액"] / total["발송건수"] if total["발송건수"] else 0
    total["발송당매출(발송횟수)"] = total["주문금액"] / total["발송횟수"] if total["발송횟수"] else 0

    change = {"_label": "증감", "연도": "", "월": ""}
    numeric_cols = [
        "발송횟수", "상품수", "URL", "발송건수", "클릭수",
        "주문건수", "주문수량", "주문금액", "객단가",
        "클릭당매출(RPC)", "발송대비매출(SPM)", "발송당매출(발송횟수)"
    ]
    rate_cols = ["반응율(Uniq CTR)", "클릭 CVR", "발송 CVR"]

    if len(d) >= 2:
        cur, prev = d.iloc[-1], d.iloc[-2]
        for col in numeric_cols:
            if col in d.columns:
                change[col] = (cur[col] - prev[col]) / abs(prev[col]) if prev[col] else pd.NA
        for col in rate_cols:
            if col in d.columns:
                change[col] = cur[col] - prev[col]
    else:
        for col in numeric_cols + rate_cols:
            change[col] = pd.NA

    return pd.concat([d, pd.DataFrame([total, change])], ignore_index=True)


def format_home_table_with_summary(df: pd.DataFrame, mode: str) -> pd.DataFrame:
    """홈 표 전용: 일반행 + 총합계 + 증감 행."""
    raw = append_total_and_change_rows(df, mode)
    if raw.empty:
        return raw

    label_name = {"Monthly": "구분", "Weekly": "구분", "Daily": "구분"}[mode]
    raw = raw.rename(columns={"_label": label_name})

    if mode == "Monthly":
        order = [
            label_name, "연도", "월", "발송횟수", "상품수", "URL", "발송건수", "클릭수",
            "반응율(Uniq CTR)", "주문건수", "주문수량", "주문금액", "객단가",
            "클릭 CVR", "발송 CVR", "클릭당매출(RPC)",
            "발송대비매출(SPM)", "발송당매출(발송횟수)"
        ]
    elif mode == "Weekly":
        order = [
            label_name, "연도", "발송횟수", "상품수", "URL", "발송건수", "클릭수",
            "반응율(Uniq CTR)", "주문건수", "주문수량", "주문금액", "객단가",
            "클릭 CVR", "발송 CVR", "클릭당매출(RPC)",
            "발송대비매출(SPM)", "발송당매출(발송횟수)"
        ]
    else:
        order = [
            label_name, "연도", "월", "발송횟수", "상품수", "URL", "발송건수", "클릭수",
            "반응율(Uniq CTR)", "주문건수", "주문수량", "주문금액", "객단가",
            "클릭 CVR", "발송 CVR", "클릭당매출(RPC)", "발송대비매출(SPM)"
        ]

    view = raw[[c for c in order if c in raw.columns]].copy()

    # pandas 최신 버전에서는 숫자형 열에 "5.7%" 같은 문자열을 다시 대입할 수 없으므로
    # 화면 표시용 데이터프레임 전체를 object 형식으로 변환합니다.
    view = view.astype("object")

    change_mask = view[label_name].astype(str).eq("증감")

    # 일반행/총합계 포맷
    for col in ["반응율(Uniq CTR)", "클릭 CVR", "발송 CVR"]:
        if col in view.columns:
            view.loc[~change_mask, col] = view.loc[~change_mask, col].map(fmt_pct)
            view.loc[change_mask, col] = view.loc[change_mask, col].map(
                lambda x: change_label(x, pp=True)
            )

    for col in [
        "발송횟수", "상품수", "URL", "발송건수", "클릭수",
        "주문건수", "주문수량", "주문금액", "객단가",
        "클릭당매출(RPC)", "발송당매출(발송횟수)"
    ]:
        if col in view.columns:
            view.loc[~change_mask, col] = view.loc[~change_mask, col].map(fmt_num)
            view.loc[change_mask, col] = view.loc[change_mask, col].map(change_label)

    if "발송대비매출(SPM)" in view.columns:
        view.loc[~change_mask, "발송대비매출(SPM)"] = (
            view.loc[~change_mask, "발송대비매출(SPM)"].map(lambda x: f"{float(x):.1f}")
        )
        view.loc[change_mask, "발송대비매출(SPM)"] = (
            view.loc[change_mask, "발송대비매출(SPM)"].map(change_label)
        )

    return view


def merge_lowest_price(product_df: pd.DataFrame, lowest_df: pd.DataFrame | None = None) -> pd.DataFrame:
    """상품 RAW의 '발송일 최저가' 컬럼을 직접 사용합니다."""
    d = product_df.copy()

    if "발송일 최저가" in d.columns:
        d["최저가"] = num(d["발송일 최저가"])
        d.loc[d["최저가"] <= 0, "최저가"] = pd.NA
        d["가격차이"] = d["최저가"] - d["멤버십혜택가"]
        d["최저가 확보"] = d.apply(
            lambda r: (
                "확보" if pd.notna(r.get("최저가")) and r["멤버십혜택가"] < r["최저가"]
                else (
                    "동일가" if pd.notna(r.get("최저가")) and r["멤버십혜택가"] == r["최저가"]
                    else ("미확보" if pd.notna(r.get("최저가")) else "")
                )
            ),
            axis=1,
        )
        return d

    # 과거 파일 호환: 발송일 최저가 컬럼이 없으면 공란 처리
    d["최저가"] = pd.NA
    d["가격차이"] = pd.NA
    d["최저가 확보"] = ""
    return d


def weekly_product_chart(sw: pd.DataFrame) -> go.Figure:
    f = sw.sort_values("_date").copy()
    labels = [
        f"{r.get('요일','')}<br>{r.get('시간대','')}<br>{r.get('소재','')}"
        for _, r in f.iterrows()
    ]
    fig = make_subplots(specs=[[{"secondary_y": True}]])
    fig.add_trace(
        go.Bar(
            x=labels, y=f["주문금액"], name="주문금액",
            marker_color="#70ad47",
            text=[fmt_num(v) for v in f["주문금액"]],
            textposition="inside",
        ),
        secondary_y=False,
    )
    fig.add_trace(
        go.Scatter(
            x=labels, y=f["주문수량"], name="주문수량",
            mode="lines+markers+text",
            line=dict(color="#f4b000", width=3),
            text=[fmt_num(v) for v in f["주문수량"]],
            textposition="top center",
        ),
        secondary_y=True,
    )
    fig.add_trace(
        go.Scatter(
            x=labels, y=f["상품수"], name="상품수",
            mode="lines+markers+text",
            line=dict(color="#ed7d31", width=2),
            text=[fmt_num(v) for v in f["상품수"]],
            textposition="bottom center",
        ),
        secondary_y=True,
    )
    fig.update_yaxes(tickformat=",", gridcolor="#ddd", secondary_y=False)
    fig.update_yaxes(tickformat=",", showgrid=False, secondary_y=True)
    fig.update_layout(
        title=dict(text="MMS 상품 실적", x=0.5, xanchor="center", font=dict(size=23)),
        height=560,
        margin=dict(l=60, r=70, t=70, b=150),
        plot_bgcolor="#ffffff",
        barmode="overlay",
        xaxis=dict(automargin=True, tickfont=dict(size=10)),
        legend=dict(orientation="h", y=-.24),
    )
    return fig


def weekly_send_chart(sw: pd.DataFrame) -> go.Figure:
    f = sw.sort_values("_date").copy()
    labels = [
        f"{r.get('요일','')}<br>{r.get('시간대','')}<br>{r.get('소재','')}"
        for _, r in f.iterrows()
    ]
    click_all = first_col(f, ["클릭 수", "클릭 수(uniq)"])
    click_uniq = first_col(f, ["클릭 수(uniq)", "클릭 수"])
    ctr_all = first_col(f, ["반응율", "반응율(uniq)"])
    ctr_uniq = first_col(f, ["반응율(uniq)", "반응율"])

    fig = make_subplots(specs=[[{"secondary_y": True}]])
    fig.add_trace(
        go.Bar(
            x=labels, y=f[click_all], name="클릭 수",
            marker_color="#5b9bd5",
            text=[fmt_num(v) for v in f[click_all]],
            textposition="inside",
        ),
        secondary_y=False,
    )
    fig.add_trace(
        go.Bar(
            x=labels, y=f[click_uniq], name="클릭 수(uniq)",
            marker_color="#a5a5a5",
            text=[fmt_num(v) for v in f[click_uniq]],
            textposition="inside",
        ),
        secondary_y=False,
    )
    fig.add_trace(
        go.Scatter(
            x=labels, y=f[ctr_all] * 100, name="반응율",
            mode="lines+markers+text",
            line=dict(color="#ed7d31", width=3),
            text=[f"{v*100:.0f}%" for v in f[ctr_all]],
            textposition="top center",
        ),
        secondary_y=True,
    )
    fig.add_trace(
        go.Scatter(
            x=labels, y=f[ctr_uniq] * 100, name="반응율(uniq)",
            mode="lines+markers+text",
            line=dict(color="#f4b000", width=3),
            text=[f"{v*100:.0f}%" for v in f[ctr_uniq]],
            textposition="bottom center",
        ),
        secondary_y=True,
    )
    fig.update_yaxes(tickformat=",", gridcolor="#ddd", secondary_y=False)
    fig.update_yaxes(ticksuffix="%", showgrid=False, secondary_y=True)
    fig.update_layout(
        title=dict(text="MMS 발송 통계", x=0.5, xanchor="center", font=dict(size=23)),
        height=560,
        margin=dict(l=60, r=70, t=70, b=150),
        plot_bgcolor="#ffffff",
        barmode="group",
        xaxis=dict(automargin=True, tickfont=dict(size=10)),
        legend=dict(orientation="h", y=-.24),
    )
    return fig


def grouped_send_table(sw: pd.DataFrame, keys: list[str]) -> pd.DataFrame:
    send_col = first_col(sw, ["발송 성공 건수", "총 발송 건수"])
    click_col = first_col(sw, ["클릭 수(uniq)", "클릭 수"])
    g = sw.groupby(keys, as_index=False).agg(
        발송횟수=("소재", "size"),
        발송건수=(send_col, "sum"),
        클릭수=(click_col, "sum"),
        주문건수=("주문건수", "sum"),
        주문금액=("주문금액", "sum"),
    )
    g["CTR(uniq)"] = safe_div(g["클릭수"], g["발송건수"])
    g["CVR(클릭>구매)"] = safe_div(g["주문건수"], g["클릭수"])
    g["객단가"] = safe_div(g["주문금액"], g["주문건수"])
    g["SPM"] = safe_div(g["주문금액"], g["발송건수"])
    g["발송당매출(발송횟수)"] = safe_div(g["주문금액"], g["발송횟수"])
    return g.fillna(0)


def add_total_row(df: pd.DataFrame, label_col: str, label="총합계") -> pd.DataFrame:
    if df.empty:
        return df
    total = {}
    for c in df.columns:
        if c == label_col:
            total[c] = label
        elif pd.api.types.is_numeric_dtype(df[c]):
            total[c] = df[c].sum()
        else:
            total[c] = ""
    return pd.concat([df, pd.DataFrame([total])], ignore_index=True)


def weekly_display_format(df: pd.DataFrame) -> pd.DataFrame:
    """주간 표 표시 형식을 안전하게 통일합니다."""
    out = df.copy()

    def format_percent_value(x):
        if pd.isna(x) or str(x).strip() in ["", "nan", "None"]:
            return ""
        if isinstance(x, str):
            value = x.strip()
            if value.endswith("%"):
                return value
            value = value.replace(",", "")
        else:
            value = x
        try:
            numeric = float(value)
            return fmt_pct(numeric)
        except (TypeError, ValueError):
            return str(x)

    def format_number_value(x):
        if pd.isna(x) or str(x).strip() in ["", "nan", "None"]:
            return ""
        if isinstance(x, str):
            value = x.strip()
            if value.endswith("%"):
                return value
            value = value.replace(",", "")
        else:
            value = x
        try:
            return fmt_num(float(value))
        except (TypeError, ValueError):
            return str(x)

    def format_spm_value(x):
        if pd.isna(x) or str(x).strip() in ["", "nan", "None"]:
            return ""
        if isinstance(x, str):
            value = x.strip().replace(",", "")
        else:
            value = x
        try:
            return f"{float(value):.1f}"
        except (TypeError, ValueError):
            return str(x)

    percent_cols = [
        "편성비중", "주문비중", "CTR(uniq)", "CVR(클릭>구매)",
        "CTR", "CVR"
    ]
    for c in percent_cols:
        if c in out.columns:
            out[c] = out[c].map(format_percent_value)

    if "SPM" in out.columns:
        out["SPM"] = out["SPM"].map(format_spm_value)

    number_cols = [
        "발송횟수", "상품수", "URL", "발송성공건수", "발송건수",
        "클릭수(uniq)", "클릭수", "객단가", "주문건수", "주문수량",
        "주문금액", "발송당매출(발송횟수)", "멤버십혜택가",
        "최저가", "가격차이", "전시순서"
    ]
    for c in number_cols:
        if c in out.columns:
            out[c] = out[c].map(format_number_value)

    return out


def category_summary_table(
    pw: pd.DataFrame,
    category_col: str,
    week: str,
    year: int,
) -> pd.DataFrame:
    """편성비중·주문비중·주문금액과 총합계를 생성합니다."""
    cat = (
        pw.groupby(category_col, dropna=False, as_index=False)
        .agg(
            편성수=("상품명", "size"),
            주문금액=("주문금액", "sum"),
        )
    )
    cat[category_col] = cat[category_col].fillna("미분류").astype(str)
    total_count = cat["편성수"].sum()
    total_amount = cat["주문금액"].sum()

    cat["편성비중"] = cat["편성수"] / total_count if total_count else 0
    cat["주문비중"] = cat["주문금액"] / total_amount if total_amount else 0

    # 주문비중 큰 순서
    cat = cat.sort_values(
        ["주문비중", "주문금액", category_col],
        ascending=[False, False, True],
    ).reset_index(drop=True)

    view = cat[[category_col, "편성비중", "주문비중", "주문금액"]].copy()
    view = view.rename(columns={category_col: "행 레이블"})
    view.insert(0, "주차", week)
    view.insert(0, "연도", year)

    total_row = pd.DataFrame([{
        "연도": "",
        "주차": "",
        "행 레이블": "총합계",
        "편성비중": 1.0 if total_count else 0,
        "주문비중": 1.0 if total_amount else 0,
        "주문금액": total_amount,
    }])
    return pd.concat([view, total_row], ignore_index=True)


def category_pie_chart(
    table: pd.DataFrame,
    title: str,
) -> go.Figure:
    """주문비중 큰 순서부터 시계방향으로 표시합니다."""
    data = table[table["행 레이블"] != "총합계"].copy()
    data = data.sort_values(
        ["주문비중", "주문금액"],
        ascending=[False, False],
    ).reset_index(drop=True)

    # 음수 주문금액은 파이에서 표현할 수 없어 0으로 처리하되 표에는 원값 유지
    values = pd.to_numeric(data["주문금액"], errors="coerce").fillna(0).clip(lower=0)

    fig = go.Figure(
        go.Pie(
            labels=data["행 레이블"],
            values=values,
            sort=False,
            direction="clockwise",
            rotation=0,
            textinfo="label+percent",
            textposition="inside",
            hovertemplate="%{label}<br>주문금액 %{value:,.0f}원<br>비중 %{percent}<extra></extra>",
        )
    )
    fig.update_layout(
        title=dict(text=title, x=0.5, xanchor="center"),
        height=560,
        margin=dict(l=40, r=40, t=70, b=40),
        uniformtext_minsize=8,
        uniformtext_mode="show",
    )
    return fig


def clean_identifier_value(x):
    if pd.isna(x) or str(x).strip() in ["", "nan", "None"]:
        return ""
    value = str(x).strip().replace(",", "")
    if value.endswith(".0"):
        value = value[:-2]
    return value


def clean_identifier_columns(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    for c in ["연도", "월", "일", "일자", "연령", "SEG", "전시순서", "쇼라코드", "알파코드"]:
        if c in out.columns:
            out[c] = out[c].map(clean_identifier_value)
    return out


def style_weekly_product_rows(formatted_df: pd.DataFrame, raw_amounts: list):
    """주간 상품실적: 총합계는 배경색을 건드리지 않고 Bold만 적용."""
    styles = pd.DataFrame("", index=formatted_df.index, columns=formatted_df.columns)

    for idx, amount in enumerate(raw_amounts):
        if idx >= len(formatted_df):
            break

        row_values = [_clean_text_value(v) for v in formatted_df.iloc[idx].tolist()]
        if any(v == "총합계" for v in row_values):
            # 배경색 지정 금지: Streamlit 기본 흰 배경 유지
            styles.iloc[idx, :] = "font-weight: 800 !important;"
            continue

        try:
            value = float(amount)
        except (TypeError, ValueError):
            continue

        if value >= 3_000_000:
            styles.iloc[idx, :] = "background-color: #fff2cc;"
        elif value < 1_000_000:
            styles.iloc[idx, :] = "background-color: #e7e6e6;"

    return styles


def weekly_delta(cur: float, prev: float, pp: bool = False) -> str:
    if pd.isna(cur) or pd.isna(prev):
        return "-"
    if pp:
        return change_label(cur - prev, pp=True)
    if prev == 0:
        return "-"
    return change_label((cur - prev) / abs(prev))


def build_weekly_detail_analysis(
    week: str,
    year: int,
    pw: pd.DataFrame,
    sw: pd.DataFrame,
    products_all: pd.DataFrame,
    sends_all: pd.DataFrame,
) -> str:
    """선택 주차의 실제 수치만 사용해 주간 분석 문구를 생성합니다."""
    send_col = first_col(sw, ["발송 성공 건수", "총 발송 건수"])
    click_col = first_col(sw, ["클릭 수(uniq)", "클릭 수"])

    send_count = float(sw[send_col].sum())
    click_count = float(sw[click_col].sum())
    order_count = float(sw["주문건수"].sum())
    qty = float(sw["주문수량"].sum())
    amount = float(sw["주문금액"].sum())
    ctr = click_count / send_count if send_count else 0
    cvr = order_count / click_count if click_count else 0
    aov = amount / order_count if order_count else 0
    spm = amount / send_count if send_count else 0

    all_weeks = (
        sends_all[sends_all["_year"] == year]
        .groupby("주차")["_date"].min()
        .sort_values()
    )
    week_names = [str(x) for x in all_weeks.index]
    prev_sw = pd.DataFrame()
    if week in week_names and week_names.index(week) > 0:
        prev_week = week_names[week_names.index(week) - 1]
        prev_sw = sends_all[
            (sends_all["_year"] == year)
            & (sends_all["주차"].astype(str) == prev_week)
        ]

    if not prev_sw.empty:
        psend = float(prev_sw[send_col].sum())
        pclick = float(prev_sw[click_col].sum())
        porders = float(prev_sw["주문건수"].sum())
        pqty = float(prev_sw["주문수량"].sum())
        pamount = float(prev_sw["주문금액"].sum())
        pctr = pclick / psend if psend else 0
        pcvr = porders / pclick if pclick else 0
        paov = pamount / porders if porders else 0
        pspm = pamount / psend if psend else 0
        prev_compare = (
            f"전주 대비 발송횟수 {weekly_delta(len(sw), len(prev_sw))} / "
            f"상품수 {len(pw):,}건 {weekly_delta(len(pw), float(prev_sw['상품수'].sum()))} / "
            f"발송건수 {weekly_delta(send_count, psend)} 운영\n"
            f"주문건수 {int(order_count):,}건 {weekly_delta(order_count, porders)} / "
            f"주문수량 {int(qty):,}건 {weekly_delta(qty, pqty)} / "
            f"주문금액 {compact_money(amount)} {weekly_delta(amount, pamount)} 기록\n"
            f"CTR {ctr*100:.1f}% {weekly_delta(ctr, pctr, pp=True)} / "
            f"CVR {cvr*100:.1f}% {weekly_delta(cvr, pcvr, pp=True)} / "
            f"객단가 {int(aov):,}원 {weekly_delta(aov, paov)} / "
            f"SPM {spm:.1f} {weekly_delta(spm, pspm)} 기록"
        )
    else:
        prev_compare = (
            f"발송횟수 {len(sw)}회 / 상품수 {len(pw)}건 / 발송건수 {int(send_count):,}건 운영\n"
            f"주문건수 {int(order_count):,}건 / 주문수량 {int(qty):,}건 / "
            f"주문금액 {compact_money(amount)} 기록\n"
            f"CTR {ctr*100:.1f}% / CVR {cvr*100:.1f}% / "
            f"객단가 {int(aov):,}원 / SPM {spm:.1f} 기록"
        )

    product_rank = (
        pw.groupby("상품명", as_index=False)
        .agg(주문금액=("주문금액", "sum"))
        .sort_values("주문금액", ascending=False)
    )
    top = product_rank.head(1)
    over5 = product_rank[product_rank["주문금액"] >= 5_000_000]["상품명"].tolist()
    over3 = product_rank[
        (product_rank["주문금액"] >= 3_000_000)
        & (product_rank["주문금액"] < 5_000_000)
    ]["상품명"].tolist()
    under1 = product_rank[product_rank["주문금액"] < 1_000_000]["상품명"].tolist()

    top_line = ""
    if not top.empty:
        top_line = (
            f"[{top.iloc[0]['상품명']}] {compact_money(top.iloc[0]['주문금액'])}으로 "
            "금주 최고 매출 기록"
        )

    send_stats = sw.copy()
    send_stats["_SPM"] = safe_div(send_stats["주문금액"], send_stats[send_col])
    send_stats["_CTR"] = safe_div(send_stats[click_col], send_stats[send_col])
    max_spm = send_stats.loc[send_stats["_SPM"].idxmax()]
    max_ctr = send_stats.loc[send_stats["_CTR"].idxmax()]
    min_spm = send_stats.loc[send_stats["_SPM"].idxmin()]

    big_cat = pw.groupby("대카", as_index=False)["주문금액"].sum().sort_values("주문금액", ascending=False)
    mid_cat = pw.groupby("중카", as_index=False)["주문금액"].sum().sort_values("주문금액", ascending=False)
    big_total = big_cat["주문금액"].sum()
    big_lines = ", ".join(
        f"{r['대카']} {r['주문금액']/big_total*100:.1f}%"
        for _, r in big_cat.head(3).iterrows()
    )
    mid_lines = " > ".join(mid_cat.head(5)["중카"].astype(str).tolist())

    seg = grouped_send_table(sw, ["성별", "연령"])
    best_seg_spm = seg.loc[seg["SPM"].idxmax()]
    best_seg_amt = seg.loc[seg["주문금액"].idxmax()]

    weekday = grouped_send_table(sw, ["요일"])
    time_df = grouped_send_table(sw, ["시간대"])
    best_day = weekday.loc[weekday["SPM"].idxmax()]
    best_time = time_df.loc[time_df["SPM"].idxmax()]

    product_insights = "\n".join(
        make_insight(r, products_all) for _, r in pw.sort_values("주문금액", ascending=False).head(6).iterrows()
    )

    price_df = merge_lowest_price(pw)
    unavailable = price_df[price_df["최저가 확보"] == "미확보"]["상품명"].dropna().astype(str).unique().tolist()
    price_line = (
        " " + " / ".join(unavailable)
        if unavailable else "발송일 최저가가 입력된 상품 중 미확보 사례 없음"
    )

    lines = [
        "■ 주간 요약",
        prev_compare,
        "",
        "■ MMS 상품 실적",
        top_line,
        ("[5백만원 이상] " + " / ".join(over5)) if over5 else "5백만원 이상 상품 없음",
        ("[3백만원 이상] " + " / ".join(over3)) if over3 else "3~5백만원 상품 없음",
        ("[1백만원 미만] " + " / ".join(under1)) if under1 else "1백만원 미만 상품 없음",
        "",
        "■ MMS 발송 통계",
        f"{max_spm.get('요일','')}({max_spm.get('시간대','')}) {max_spm.get('소재','')} "
        f"SPM {max_spm['_SPM']:.1f} 금주 최고 효율 기록",
        f"{max_ctr.get('요일','')}({max_ctr.get('시간대','')}) {max_ctr.get('소재','')} "
        f"CTR {max_ctr['_CTR']*100:.1f}% 금주 최고 반응 기록",
        f"{min_spm.get('요일','')}({min_spm.get('시간대','')}) {min_spm.get('소재','')} "
        f"SPM {min_spm['_SPM']:.1f} 금주 최저 효율 기록",
        "",
        "■ 카테고리 분석",
        f"대카테고리 기준 {big_lines} 매출 비중 차지",
        f"중카테고리 기준 {mid_lines} 중심 매출 구성",
        "",
        "■ SEG 분석",
        f"{best_seg_spm['성별']}{best_seg_spm['연령']} SPM {best_seg_spm['SPM']:.1f}로 최고 효율 기록",
        f"{best_seg_amt['성별']}{best_seg_amt['연령']} 주문금액 {compact_money(best_seg_amt['주문금액'])}으로 최고 매출 기여",
        "",
        "■ 요일·시간대 분석",
        f"{best_day['요일']}요일 SPM {best_day['SPM']:.1f}로 최고 효율 기록",
        f"{best_time['시간대']} 시간대 SPM {best_time['SPM']:.1f}로 시간대 최고 효율 기록",
        "",
        "■ 상품 인사이트",
        product_insights,
        "",
        "■ 최저가 미확보 상품",
        price_line,
    ]
    return "\n".join(str(x) for x in lines if x is not None)


def _weekly_plain_delta(cur: float, prev: float, pp: bool = False) -> str:
    if pd.isna(cur) or pd.isna(prev):
        return "-"
    if pp:
        diff = (cur - prev) * 100
        if abs(diff) < 0.05:
            diff = 0.0
        return f"{diff:+.1f}%p" if diff != 0 else "0.0%p"
    if prev == 0:
        return "-"
    diff = (cur - prev) / abs(prev) * 100
    if abs(diff) < 0.05:
        diff = 0.0
    return f"{diff:+.1f}%" if diff != 0 else "0.0%"


def _season_recommendations(month: int, operated_text: str):
    season_map = {
        1: ("한파·설 명절", ["온열가전", "겨울 침구", "보온의류", "명절 식품·선물세트"]),
        2: ("환절기·신학기", ["건강식품", "공기청정·청소가전", "신학기 생활용품", "간편식"]),
        3: ("봄·신학기", ["봄 의류", "청소·생활가전", "건강식품", "나들이용품"]),
        4: ("봄 나들이·가정의달 준비", ["선케어", "나들이 식품", "건강식품", "선물형 소형가전"]),
        5: ("가정의달·초여름", ["건강식품", "뷰티 선물", "선풍기·서큘레이터", "여름 의류"]),
        6: ("장마·초여름", ["제습·건조가전", "선풍기·서큘레이터", "여름 침구·쿨링용품", "간편식"]),
        7: ("폭염·휴가 시즌", ["캐리어", "여행용 소형가전", "선케어", "기능성·냉감 의류", "보양식·간편식"]),
        8: ("폭염·휴가·추석 준비", ["냉방가전", "여행용품", "선케어", "보양식", "명절 선물 사전수요 상품"]),
        9: ("추석·환절기", ["명절 식품·선물세트", "건강식품", "환절기 의류", "생활가전"]),
        10: ("가을·동절기 준비", ["온열가전", "가을 의류", "건강식품", "침구"]),
        11: ("초겨울·연말 준비", ["온열가전", "겨울 침구", "보온의류", "연말 선물형 상품"]),
        12: ("한파·연말", ["온열가전", "겨울 의류", "홈파티 식품", "선물형 가전·뷰티"]),
    }
    theme, groups = season_map.get(month, ("시즌 수요", ["시즌 상품"]))
    fresh = [g for g in groups if g.split("·")[0] not in operated_text][:3] or groups[:3]
    return theme, "·".join(fresh)



def _short_weekly_product_name(name: str) -> str:
    """주간 화면용 상품명 축약. 브랜드 + 핵심 상품 + 핵심 구성은 남기고 원본 데이터는 변경하지 않음."""
    original = str(name or "").strip()
    s = original

    # 앞쪽 운영 태그/무료배송 태그 제거
    s = re.sub(r"^\s*(?:\[(?:M|무료배송|쇼라 단독 특가|멤버스특가)[^\]]*\]\s*)+", "", s, flags=re.I)
    s = re.sub(r"^\s*\(M\)\s*", "", s, flags=re.I)

    # 브랜드 대괄호는 텍스트로 살림
    s = re.sub(r"\[([^\]]+)\]", r"\1", s)

    # 불필요한 운영 메모만 제거
    s = re.sub(r"\s*\((?:재고부족|편성\s*\d+회|추가\s*멤포\s*상품)[^)]*\)", "", s, flags=re.I)

    # 샘플/쇼핑백/증정 부가구성 제거
    s = re.sub(r"\s*\((?:샘플[^)]*|쇼핑백[^)]*|증정[^)]*)\)", "", s, flags=re.I)

    # 슬래시 옵션은 무조건 자르지 않음. 제품 유형이 드러나도록 대표어 정리.
    if "붕어싸만코" in s and "아이스크림" in s:
        qty = re.search(r"(?:총\s*)?(\d+)\s*개", s)
        s = "빙그레 붕어싸만코 아이스크림" + (f" {qty.group(1)}개" if qty else "")
    elif "이지프로" in s and "면도기" in s:
        s = re.sub(r"S\d+/\d+", "", s)
        s = re.sub(r"\s+", " ", s).strip()
    elif "디올" in s and "립" in s:
        gram = re.search(r"(\d+(?:\.\d+)?)\s*g", s, re.I)
        s = "디올 어딕트 립 글로우 립밤" + (f" {gram.group(1)}g" if gram else "")

    # NEW/특별혜택가 등 수식어 정리
    s = re.sub(r"\bNEW\b|\(NEW\)|특별혜택가|신형", "", s, flags=re.I)
    s = re.sub(r"\s+", " ", s).strip(" -/")

    # 너무 길 때는 핵심 수량/유형을 훼손하지 않도록 마지막 괄호 부가설명부터 제거
    if len(s) > 48:
        s2 = re.sub(r"\s*\([^)]{8,}\)\s*$", "", s).strip()
        if len(s2) >= 12:
            s = s2
    if len(s) > 58:
        # 최후 수단: 단어 단위 절단. ellipsis는 사용하되 제품 유형이 이미 앞에 남는 경우만.
        cut = s[:58]
        if " " in cut:
            cut = cut.rsplit(" ", 1)[0]
        s = cut.rstrip() + "…"

    return s or original




def _extract_unit_count_from_name(name: str):
    """상품명에서 총 수량/매수 추출. 2+1, 3+3, 24롤×2팩, 본품+리필 등 복합 구성을 우선 해석."""
    s = str(name or "").replace("×", "x").replace("X", "x")

    # 1) 3+3, 2+1 등 합산형
    m = re.search(r"(\d+)\s*\+\s*(\d+)", s)
    if m:
        a, b = int(m.group(1)), int(m.group(2))
        total = a + b
        if 1 <= total <= 500:
            return total

    # 2) 24롤 x 2팩 / 12개 x 2박스
    m = re.search(r"(\d+)\s*(롤|개|매|봉|캔|팩)\s*x\s*(\d+)\s*(팩|박스|세트)?", s, re.I)
    if m:
        total = int(m.group(1)) * int(m.group(3))
        if 1 <= total <= 1000:
            return total

    # 3) 본품 + 리필 N개
    base = 0
    if re.search(r"본품", s):
        base = 1
    refill = re.search(r"리필[^0-9]{0,10}(\d+)\s*개", s)
    if refill:
        total = base + int(refill.group(1))
        if 1 <= total <= 500:
            return total

    # 4) 총 N개/매/봉 등 명시
    m = re.search(r"(?:총\s*)?(\d+)\s*(개|매|봉|캔|팩|롤|병|포|입)", s)
    if m:
        v = int(m.group(1))
        if 1 <= v <= 1000:
            return v

    # 5) 일반 패턴 중 최대값
    patterns = [
        r"(\d+)\s*매",
        r"(\d+)\s*개",
        r"(\d+)\s*봉",
        r"(\d+)\s*캔",
        r"(\d+)\s*팩",
        r"(\d+)\s*롤",
        r"(\d+)\s*병",
        r"(\d+)\s*포",
        r"(\d+)\s*입",
    ]
    vals = []
    for p in patterns:
        for mm in re.finditer(p, s):
            try:
                v = int(mm.group(1))
                if 1 <= v <= 1000:
                    vals.append(v)
            except Exception:
                pass
    return max(vals) if vals else None



def _unit_price_phrase(name: str, sale_price: float):
    cnt = _extract_unit_count_from_name(name)
    if not cnt or not sale_price or pd.isna(sale_price):
        return None
    unit = sale_price / cnt
    if "매" in str(name):
        label = "장당"
    else:
        label = "개당"
    return f"{label} {unit:,.0f}원"


def _weekly_product_history_stats(product_name: str, all_products: pd.DataFrame, week_end):
    h = all_products[all_products["상품명"].astype(str) == str(product_name)].copy()
    if h.empty:
        return None
    h["_date2"] = pd.to_datetime(h["_date"], errors="coerce")
    h = h[h["_date2"].notna() & (h["_date2"] <= week_end)]
    if h.empty:
        return None
    h = h.sort_values("_date2")
    amt = pd.to_numeric(h["주문금액"], errors="coerce").fillna(0)
    sale_col = first_col(h, ["멤버십 혜택가", "행사가", "판매가", "혜택가"])
    sale_vals = pd.to_numeric(h[sale_col], errors="coerce") if sale_col else pd.Series(dtype=float)
    return {
        "count": len(h),
        "avg": float(amt.mean()) if len(amt) else 0,
        "max": float(amt.max()) if len(amt) else 0,
        "ge3m": int((amt >= 3_000_000).sum()),
        "ge5m": int((amt >= 5_000_000).sum()),
        "last_date": h["_date2"].max(),
        "recent3": amt.tail(3).tolist(),
        "sale_last": float(sale_vals.iloc[-1]) if len(sale_vals) and pd.notna(sale_vals.iloc[-1]) else None,
        "sale_min": float(sale_vals.min()) if len(sale_vals) and sale_vals.notna().any() else None,
        "sale_max": float(sale_vals.max()) if len(sale_vals) and sale_vals.notna().any() else None,
    }


def _dominant_target_products(pw: pd.DataFrame, gender: str, age: str, topn: int = 3):
    sub = pw[
        (pw["성별"].astype(str) == str(gender))
        & (pw["연령"].astype(str).map(clean_identifier_value) == clean_identifier_value(age))
    ].copy()
    if sub.empty:
        return []
    g = sub.groupby("상품명", as_index=False)["주문금액"].sum().sort_values("주문금액", ascending=False)
    return [_short_weekly_product_name(x) for x in g.head(topn)["상품명"].astype(str)]


def _weekly_food_underperformers(pw: pd.DataFrame):
    sub = pw[pw["대카"].astype(str).str.contains("식품|건강", na=False)].copy()
    if sub.empty:
        return []
    g = sub.groupby("상품명", as_index=False)["주문금액"].sum().sort_values("주문금액")
    return [_short_weekly_product_name(x) for x in g[g["주문금액"] < 1_000_000].head(4)["상품명"].astype(str)]



def _recent_4week_time_pattern(current_week: str, year: int, sends_all: pd.DataFrame):
    """최근 4주 요일/시간대 SPM 반복성. 3주 이상 동일 우위일 때만 인사이트 후보 반환."""
    send_col = first_col(sends_all, ["발송 성공 건수", "총 발송 건수"])
    if not send_col or sends_all.empty:
        return None

    wk_dates = (
        sends_all[sends_all["_year"] == year]
        .groupby("주차")["_date"].min()
        .sort_values()
    )
    weeks = [str(x) for x in wk_dates.index]
    if current_week not in weeks:
        return None
    idx = weeks.index(current_week)
    selected = weeks[max(0, idx-3):idx+1]
    if len(selected) < 3:
        return None

    time_winners = []
    day_winners = []
    for w in selected:
        sub = sends_all[(sends_all["_year"] == year) & (sends_all["주차"].astype(str) == w)].copy()
        if sub.empty:
            continue
        tg = grouped_send_table(sub, ["시간대"])
        dg = grouped_send_table(sub, ["요일"])
        if not tg.empty:
            time_winners.append(str(tg.loc[tg["SPM"].idxmax(), "시간대"]))
        if not dg.empty:
            day_winners.append(str(dg.loc[dg["SPM"].idxmax(), "요일"]))

    def winner_summary(vals):
        if not vals:
            return None
        vc = pd.Series(vals).value_counts()
        winner = str(vc.index[0])
        cnt = int(vc.iloc[0])
        return winner, cnt, len(vals)

    return {
        "time": winner_summary(time_winners),
        "day": winner_summary(day_winners),
        "weeks": selected,
    }


def _detect_current_product_status(product_name: str, all_products: pd.DataFrame):
    """최신 행의 상품명/상태성 컬럼에서 판매 가능성 확인. 불명확하면 None."""
    h = all_products[all_products["상품명"].astype(str) == str(product_name)].copy()
    if h.empty:
        return None
    h["_date2"] = pd.to_datetime(h["_date"], errors="coerce")
    h = h.sort_values("_date2")
    row = h.iloc[-1]

    status_cols = [c for c in ["상태","판매상태","상품상태","전시상태","재고상태"] if c in h.columns]
    raw = " ".join(str(row.get(c, "")) for c in status_cols)
    raw += " " + str(row.get("상품명",""))

    negative = ["판매중지","판매종료","품절","재고부족","전시종료","중단"]
    positive = ["판매중","판매가능","정상판매","재고있음"]

    if any(k in raw for k in negative):
        return False
    if any(k in raw for k in positive):
        return True
    return None


def _promotion_performance_stats(product_name: str, all_products: pd.DataFrame):
    """프로모션 컬럼이 있으면 일반/프로모션 평균 분리."""
    if "프로모션" not in all_products.columns:
        return None
    h = all_products[all_products["상품명"].astype(str) == str(product_name)].copy()
    if h.empty:
        return None
    h["주문금액"] = pd.to_numeric(h["주문금액"], errors="coerce").fillna(0)
    promo_mask = h["프로모션"].fillna("").astype(str).str.strip().replace("-", "").ne("")
    promo = h[promo_mask]
    normal = h[~promo_mask]
    if promo.empty or normal.empty:
        return None
    return {
        "promo_avg": float(promo["주문금액"].mean()),
        "normal_avg": float(normal["주문금액"].mean()),
        "promo_n": len(promo),
        "normal_n": len(normal),
    }


def _latest_and_high_perf_price(product_name: str, all_products: pd.DataFrame):
    """현재/최근 가격과 과거 고성과 운영 가격 비교."""
    h = all_products[all_products["상품명"].astype(str) == str(product_name)].copy()
    if h.empty:
        return None
    price_col = first_col(h, ["멤버십 혜택가","행사가","판매가","혜택가"])
    if not price_col:
        return None
    h["_date2"] = pd.to_datetime(h["_date"], errors="coerce")
    h["_price"] = pd.to_numeric(h[price_col], errors="coerce")
    h["_amt"] = pd.to_numeric(h["주문금액"], errors="coerce").fillna(0)
    h = h[h["_price"].notna()].sort_values("_date2")
    if h.empty:
        return None
    latest = h.iloc[-1]
    high = h[h["_amt"] >= 5_000_000]
    hp = float(high["_price"].mean()) if not high.empty else None
    return {
        "latest_price": float(latest["_price"]),
        "high_perf_avg_price": hp,
        "latest_date": latest["_date2"],
    }



def _product_target_strength_analysis(product_name: str, all_products: pd.DataFrame, week_end):
    """
    상품별 성별→연령→SEG 강세와 반복 운영 패턴을 실제 이력으로 판정.
    프로모션은 이 판정에 사용하지 않음.
    근거: 운영횟수, 주문금액, 평균매출, 300/500만원 이상 횟수, SPM(가능한 경우).
    """
    h = all_products[all_products["상품명"].astype(str) == str(product_name)].copy()
    if h.empty:
        return None

    h["_date2"] = pd.to_datetime(h["_date"], errors="coerce")
    h = h[h["_date2"].notna() & (h["_date2"] <= week_end)].sort_values("_date2")
    if len(h) < 2:
        return {"type": "insufficient", "count": len(h)}

    h["_amt"] = pd.to_numeric(h["주문금액"], errors="coerce").fillna(0)
    send_col = first_col(h, ["발송 성공 건수", "총 발송 건수", "발송건수"])
    if send_col:
        h["_send"] = pd.to_numeric(h[send_col], errors="coerce").fillna(0)
        h["_spm_calc"] = h["_amt"] / h["_send"].replace(0, pd.NA)
    elif "SPM" in h.columns:
        h["_spm_calc"] = pd.to_numeric(h["SPM"], errors="coerce")
    else:
        h["_spm_calc"] = pd.NA

    def stats(group_cols):
        g = h.groupby(group_cols, dropna=False).agg(
            운영횟수=("_amt","size"),
            평균매출=("_amt","mean"),
            최고매출=("_amt","max"),
            삼백이상=("_amt", lambda x: int((x >= 3_000_000).sum())),
            오백이상=("_amt", lambda x: int((x >= 5_000_000).sum())),
            평균SPM=("_spm_calc","mean"),
        ).reset_index()
        g["고성과율"] = g["오백이상"] / g["운영횟수"]
        return g

    gender = stats(["성별"]) if "성별" in h.columns else pd.DataFrame()
    age = stats(["성별","연령"]) if {"성별","연령"}.issubset(h.columns) else pd.DataFrame()
    seg = stats(["성별","연령","SEG"]) if {"성별","연령","SEG"}.issubset(h.columns) else pd.DataFrame()

    # 동일 타겟 반복 피로도: 최근 동일 성별/연령/SEG 3회 이상 연속 감소 + 최초 대비 30% 이상 감소
    fatigue = None
    if {"성별","연령","SEG"}.issubset(h.columns):
        for keys, sub in h.groupby(["성별","연령","SEG"], dropna=False):
            sub = sub.sort_values("_date2").tail(4)
            vals = sub["_amt"].tolist()
            if len(vals) >= 3:
                recent = vals[-3:]
                decreasing = all(recent[i] < recent[i-1] for i in range(1, len(recent)))
                drop = (recent[0] - recent[-1]) / recent[0] if recent[0] > 0 else 0
                if decreasing and drop >= 0.30:
                    fatigue = {
                        "target": keys, "vals": recent, "drop": drop,
                        "count": len(sub)
                    }
                    break

    # 안정 반복: 동일 성별/연령에서 최근 3회 이상 모두 300만원 이상, SEG 2개 이상이면 강한 근거
    stable = None
    if {"성별","연령"}.issubset(h.columns):
        for keys, sub in h.groupby(["성별","연령"], dropna=False):
            sub = sub.sort_values("_date2").tail(4)
            vals = sub["_amt"].tolist()
            seg_n = sub["SEG"].nunique(dropna=True) if "SEG" in sub.columns else 0
            if len(vals) >= 3 and all(v >= 3_000_000 for v in vals):
                stable = {"target": keys, "vals": vals, "seg_n": int(seg_n), "count": len(sub)}
                break

    # 성별 강세: 양쪽 2회 이상 + 평균매출 1.5배 이상 + SPM도 열위가 아니어야 강세 판정
    gender_strength = None
    if len(gender) >= 2:
        eligible = gender[gender["운영횟수"] >= 2].sort_values("평균매출", ascending=False)
        if len(eligible) >= 2:
            a, b = eligible.iloc[0], eligible.iloc[1]
            ratio = a["평균매출"] / b["평균매출"] if b["평균매출"] > 0 else float("inf")
            spm_ok = True
            if pd.notna(a["평균SPM"]) and pd.notna(b["평균SPM"]) and b["평균SPM"] > 0:
                spm_ok = a["평균SPM"] >= b["평균SPM"]
            if ratio >= 1.5 and spm_ok:
                gender_strength = (a, b, ratio)

    # 연령 강세: 같은 성별 안에서 3040/5060 등 2회 이상씩 비교
    age_strength = None
    if not age.empty:
        for gender_name, gg in age.groupby("성별"):
            eligible = gg[gg["운영횟수"] >= 2].sort_values("평균매출", ascending=False)
            if len(eligible) >= 2:
                a, b = eligible.iloc[0], eligible.iloc[1]
                ratio = a["평균매출"] / b["평균매출"] if b["평균매출"] > 0 else float("inf")
                spm_ok = True
                if pd.notna(a["평균SPM"]) and pd.notna(b["평균SPM"]) and b["평균SPM"] > 0:
                    spm_ok = a["평균SPM"] >= b["평균SPM"]
                if ratio >= 1.5 and spm_ok:
                    age_strength = (a, b, ratio)
                    break

    return {
        "type": "ok",
        "fatigue": fatigue,
        "stable": stable,
        "gender_strength": gender_strength,
        "age_strength": age_strength,
        "gender": gender,
        "age": age,
        "seg": seg,
    }


def _target_strength_sentence(product_name: str, analysis):
    if not analysis or analysis.get("type") != "ok":
        return None
    short = _short_weekly_product_name(product_name)

    # 1순위: 동일 타겟 반복 하락
    f = analysis.get("fatigue")
    if f:
        g, a, seg = f["target"]
        vals = " → ".join(compact_money(v) for v in f["vals"])
        return (
            f"• {_with_topic(short)} 동일 타겟 {g}{clean_identifier_value(a)} {seg}에서 최근 3회 주문금액이 "
            f"{vals}으로 연속 하락해 최초 대비 {f['drop']*100:.0f}% 감소했습니다. "
            f"타겟 변경에 따른 차이가 아닌 동일 조건 반복 운영에서 성과 둔화가 확인된 만큼, "
            f"즉시 동일 SEG 재편성보다 최근 미발송 SEG로 전환 TEST하거나 일정 기간 미편성 후 재운영하는 것이 적절합니다."
        )

    # 2순위: 특정 성별/연령 강세
    ag = analysis.get("age_strength")
    if ag:
        a, b, ratio = ag
        spm_txt = ""
        if pd.notna(a["평균SPM"]) and pd.notna(b["평균SPM"]):
            spm_txt = f", 평균 SPM도 {a['평균SPM']:.1f} vs {b['평균SPM']:.1f}"
        return (
            f"• {_with_topic(short)} {a['성별']}{clean_identifier_value(a['연령'])}에서 {int(a['운영횟수'])}회 평균 "
            f"{compact_money(a['평균매출'])}, 500만원 이상 {int(a['오백이상'])}회를 기록한 반면 "
            f"{b['성별']}{clean_identifier_value(b['연령'])}은 {int(b['운영횟수'])}회 평균 {compact_money(b['평균매출'])}"
            f"{spm_txt}으로 차이가 확인됐습니다. 평균매출이 {ratio:.1f}배 높은 {a['성별']}{clean_identifier_value(a['연령'])}을 "
            f"우선 편성하되, 해당 연령대 내 고성과 SEG와 미발송 SEG를 순차 TEST하는 것이 적절합니다."
        )

    gs = analysis.get("gender_strength")
    if gs:
        a, b, ratio = gs
        spm_txt = ""
        if pd.notna(a["평균SPM"]) and pd.notna(b["평균SPM"]):
            spm_txt = f", 평균 SPM {a['평균SPM']:.1f} vs {b['평균SPM']:.1f}"
        return (
            f"• {_with_topic(short)} {a['성별']} 타겟 {int(a['운영횟수'])}회 평균 {compact_money(a['평균매출'])}, "
            f"500만원 이상 {int(a['오백이상'])}회로 {b['성별']} 평균 {compact_money(b['평균매출'])} 대비 {ratio:.1f}배 높았고"
            f"{spm_txt}로 효율도 함께 확인됐습니다. 단일 회차가 아닌 반복 이력에서 성별 강세가 확인된 만큼 "
            f"{a['성별']} 중심으로 편성하되 연령·SEG별 성과를 기준으로 세부 타겟을 좁히는 것이 적절합니다."
        )

    # 3순위: 안정적 반복 운영
    s = analysis.get("stable")
    if s:
        g, a = s["target"]
        vals = " → ".join(compact_money(v) for v in s["vals"])
        seg_text = f", {s['seg_n']}개 SEG에서 운영" if s["seg_n"] >= 2 else ""
        return (
            f"• {_with_topic(short)} {g}{clean_identifier_value(a)}에서 최근 {len(s['vals'])}회 주문금액이 {vals}으로 "
            f"모두 300만원 이상을 유지했고{seg_text}했습니다. 동일 성별·연령에서 반복 성과가 유지돼 "
            f"해당 타겟 적합도가 확인된 상품으로, 고성과 SEG 순환과 미발송 SEG 확대 TEST를 병행할 수 있습니다."
        )
    return None



def _next_week_action_candidates(pw, products_all, week_end):
    """근거가 충분한 차주 운영 제안 후보를 점수화. 타겟/가격/성과/미편성/카테고리 근거 사용."""
    actions=[]; current=set(pw["상품명"].dropna().astype(str))
    wk=pw.groupby("상품명",as_index=False).agg(주문금액=("주문금액","sum"),금주운영횟수=("상품명","size")).sort_values("주문금액",ascending=False)

    # 즉시 재편성 + 타겟 강세
    for _,r in wk[wk["주문금액"]>=5_000_000].head(5).iterrows():
        pname=str(r["상품명"]); hist=_weekly_product_history_stats(pname,products_all,week_end)
        if not hist: continue
        ta=_product_target_strength_analysis(pname,products_all,week_end)
        s=f"{_with_topic(_short_weekly_product_name(pname))} 금주 {compact_money(r['주문금액'])}, 누적 {hist['count']}회 운영 중 500만원 이상 {hist['ge5m']}회를 기록해 차주 재편성 우선 후보입니다."
        if ta and ta.get("age_strength"):
            a,b,ratio=ta["age_strength"]
            s+=f" {a['성별']}{clean_identifier_value(a['연령'])} 평균매출이 비교 연령대 대비 {ratio:.1f}배 높아 해당 타겟 내 고성과 SEG와 미발송 SEG를 순차 편성하는 것이 적절합니다."
        elif ta and ta.get("gender_strength"):
            a,b,ratio=ta["gender_strength"]
            s+=f" {a['성별']} 평균매출이 {b['성별']} 대비 {ratio:.1f}배 높아 {a['성별']} 중심 편성이 적절합니다."
        else:
            s+=" 최근 고성과 타겟을 우선 유지하되 동일 SEG 과다 반복은 피하는 것이 적절합니다."
        actions.append((100+float(r["주문금액"])/1e6,"즉시 재편성",s))

    histdf=products_all.copy()
    histdf["_date2"]=pd.to_datetime(histdf["_date"],errors="coerce")

    # 가격 회복 시 재운영
    for pname,h in histdf.groupby("상품명"):
        pname=str(pname)
        if pname in current: continue
        h=h[h["_date2"].notna()&(h["_date2"]<=week_end)]
        amt=pd.to_numeric(h["주문금액"],errors="coerce").fillna(0)
        if len(h)<2 or (amt>=5_000_000).sum()<1: continue
        pi=_latest_and_high_perf_price(pname,products_all)
        if not pi or not pi["high_perf_avg_price"]: continue
        diff=(pi["latest_price"]-pi["high_perf_avg_price"])/pi["high_perf_avg_price"]*100
        if diff>=10:
            actions.append((70,"가격 조건",f"{_with_topic(_short_weekly_product_name(pname))} 과거 500만원 이상 고성과 이력이 있으나 최신 혜택가 {pi['latest_price']:,.0f}원으로 고성과 당시 평균 {pi['high_perf_avg_price']:,.0f}원 대비 {diff:.1f}% 높습니다. 현재 조건에서는 우선순위를 낮추고 과거 고성과 가격대에 근접할 경우 재운영을 검토하는 것이 적절합니다."))

    # 최근 미편성 고성과
    past=histdf[~histdf["상품명"].astype(str).isin(current)].copy()
    if not past.empty:
        dorm=past.groupby("상품명",as_index=False).agg(
            운영횟수=("상품명","size"),평균매출=("주문금액","mean"),최고매출=("주문금액","max"),
            최근발송일=("_date2","max"),고성과횟수=("주문금액",lambda x:int((pd.to_numeric(x,errors="coerce").fillna(0)>=5_000_000).sum())))
        dorm["미편성일수"]=(week_end.normalize()-dorm["최근발송일"].dt.normalize()).dt.days
        dorm=dorm[(dorm["운영횟수"]>=2)&(dorm["평균매출"]>=3_000_000)&(dorm["고성과횟수"]>=1)&dorm["미편성일수"].between(14,120)]
        for _,r in dorm.sort_values(["고성과횟수","평균매출"],ascending=False).head(5).iterrows():
            pname=str(r["상품명"])
            if _detect_current_product_status(pname,products_all) is False: continue
            pi=_latest_and_high_perf_price(pname,products_all)
            if pi and pi["high_perf_avg_price"]:
                diff=(pi["latest_price"]-pi["high_perf_avg_price"])/pi["high_perf_avg_price"]*100
                if diff>10: continue
            actions.append((80+r["고성과횟수"]*2,"최근 미편성",f"{_with_topic(_short_weekly_product_name(pname))} 과거 {int(r['운영횟수'])}회 중 {int(r['고성과횟수'])}회 500만원 이상, 평균 {compact_money(r['평균매출'])}을 기록했고 최근 {int(r['미편성일수'])}일간 미편성 상태입니다. 현재 판매 가능 여부와 최신 가격 조건을 확인한 뒤, 과거 고성과 당시와 유사한 조건이 유지되면 차주 재편성 후보로 검토하는 것이 적절합니다."))

    # 반복 부진 상품이지만 카테고리는 강한 경우 → 상품 교체
    total=float(pw["주문금액"].sum())
    cat=pw.groupby("대카",as_index=False)["주문금액"].sum()
    shares={str(r["대카"]):float(r["주문금액"])/total*100 for _,r in cat.iterrows()} if total else {}
    for pname,h in products_all.groupby("상품명"):
        pname=str(pname)
        if pname not in current: continue
        hh=h[pd.to_datetime(h["_date"],errors="coerce")<=week_end]
        vals=pd.to_numeric(hh["주문금액"],errors="coerce").fillna(0)
        if len(vals)>=2 and int((vals>=3_000_000).sum())==0:
            wkrow=pw[pw["상품명"].astype(str)==pname]
            if wkrow.empty: continue
            catname=str(wkrow["대카"].iloc[0]); share=shares.get(catname,0)
            if share>=20:
                actions.append((65,"상품 교체",f"{_with_topic(_short_weekly_product_name(pname))} 과거 {len(vals)}회 운영 중 300만원 이상 달성 이력이 없지만 {catname} 카테고리는 금주 전체 주문금액의 {share:.1f}%를 차지했습니다. 카테고리 비중을 줄이기보다 동일 카테고리 내 과거 300만원 이상 성과가 반복된 검증 상품으로 교체하는 것이 적절합니다."))

    seen=set(); out=[]
    for score,kind,s in sorted(actions,key=lambda x:x[0],reverse=True):
        if s in seen: continue
        seen.add(s); out.append((score,kind,s))
    return out


def _season_gap_action(pw, products_all, week_end):
    """최근 4주 시즌 상품군별 실제 편성 횟수 차이가 클 때만 공백 제안."""
    start=week_end-pd.Timedelta(days=27)
    d=pd.to_datetime(products_all["_date"],errors="coerce")
    recent=products_all[(d>=start)&(d<=week_end)].copy()
    if int(week_end.month) not in [6,7,8]: return None
    groups={"냉방가전":r"선풍기|써큘|서큘|에어컨|냉풍","여행용품":r"캐리어|여행|파우치|보스턴백","선케어":r"선크림|선스틱|선쿠션|자외선","냉감·기능성 의류":r"냉감|쿨링|드라이셀|기능성|언더셔츠"}
    names=recent["상품명"].fillna("").astype(str)
    counts={k:int(names.str.contains(p,case=False,regex=True).sum()) for k,p in groups.items()}
    mx=max(counts,key=counts.get); mn=min(counts,key=counts.get)
    if counts[mx]>=max(3,counts[mn]*2+1):
        detail=" / ".join(f"{k} {v}회" for k,v in counts.items())
        return f"최근 4주 시즌 상품 편성은 {detail}로 구성됐습니다. {mx} 편성이 상대적으로 많은 반면 {mn} 편성이 적어, 동일 상품군 반복 확대보다 {mn} 신규·유사신규 상품 TEST를 우선 검토해 시즌 매출원을 분산할 필요가 있습니다."
    return None



def _seasonal_last_year_evidence(products_all: pd.DataFrame, week_end):
    """
    시즌성 근거 탐색 우선순위:
    1) 전년 동일시점 ±35일
    2) 전년 동일 시즌 월(기준월 ±1개월)
    3) 보유 데이터 전체 과거 동일 시즌 월
    상품명 + 대카 + 중카를 함께 탐지하며 실제 300만원 이상 이력이 있는 상품만 반환.
    """
    if pd.isna(week_end):
        return []

    df = products_all.copy()
    df["_date2"] = pd.to_datetime(df["_date"], errors="coerce")
    df = df[df["_date2"].notna() & (df["_date2"] < week_end)].copy()
    if df.empty:
        return []

    price_col = first_col(df, ["멤버십 혜택가", "멤버십혜택가", "혜택가", "최종혜택가", "행사가", "판매가", "MMS혜택가", "MMS 혜택가", "실판매가"])
    df["_amt"] = pd.to_numeric(df["주문금액"], errors="coerce").fillna(0)
    df["_price"] = pd.to_numeric(df[price_col], errors="coerce") if price_col else pd.NA

    prior_year = int(week_end.year) - 1
    exact_start = (week_end - pd.DateOffset(years=1) - pd.Timedelta(days=35)).normalize()
    exact_end = (week_end - pd.DateOffset(years=1) + pd.Timedelta(days=35)).normalize()

    season_months = {
        1:[12,1,2], 2:[1,2,3], 3:[2,3,4], 4:[3,4,5],
        5:[4,5,6], 6:[5,6,7], 7:[6,7,8], 8:[7,8,9],
        9:[8,9,10], 10:[9,10,11], 11:[10,11,12], 12:[11,12,1]
    }[int(week_end.month)]

    scopes = [
        ("전년 동시점", df[df["_date2"].between(exact_start, exact_end)].copy()),
        ("전년 동시즌", df[(df["_date2"].dt.year == prior_year) & (df["_date2"].dt.month.isin(season_months))].copy()),
        ("과거 동시즌", df[df["_date2"].dt.month.isin(season_months)].copy()),
    ]

    labels = pd.Series([""] * len(df), index=df.index)
    if "대카" in df.columns:
        labels = labels + " " + df["대카"].fillna("").astype(str)
    if "중카" in df.columns:
        labels = labels + " " + df["중카"].fillna("").astype(str)

    patterns = _season_keyword_match_mask(df["상품명"], labels)

    for scope_name, scope_df in scopes:
        if scope_df.empty:
            continue
        scope_labels = pd.Series([""] * len(scope_df), index=scope_df.index)
        if "대카" in scope_df.columns:
            scope_labels = scope_labels + " " + scope_df["대카"].fillna("").astype(str)
        if "중카" in scope_df.columns:
            scope_labels = scope_labels + " " + scope_df["중카"].fillna("").astype(str)

        combined = scope_df["상품명"].fillna("").astype(str) + " " + scope_labels
        results = []

        for group_name, pat_kw in patterns.items():
            sub = scope_df[combined.str.contains(pat_kw, case=False, regex=True, na=False)].copy()
            if sub.empty:
                continue

            for pname, h in sub.groupby("상품명"):
                h = h.sort_values("_date2")
                n = len(h)
                avg_amt = float(h["_amt"].mean())
                max_amt = float(h["_amt"].max())
                ge3 = int((h["_amt"] >= 3_000_000).sum())
                ge5 = int((h["_amt"] >= 5_000_000).sum())
                if ge3 < 1:
                    continue

                avg_price = float(h["_price"].mean()) if pd.to_numeric(h["_price"], errors="coerce").notna().any() else None
                hp3 = h[h["_amt"] >= 3_000_000]
                hp5 = h[h["_amt"] >= 5_000_000]
                hp3_price = float(pd.to_numeric(hp3["_price"], errors="coerce").mean()) if not hp3.empty and pd.to_numeric(hp3["_price"], errors="coerce").notna().any() else None
                hp5_price = float(pd.to_numeric(hp5["_price"], errors="coerce").mean()) if not hp5.empty and pd.to_numeric(hp5["_price"], errors="coerce").notna().any() else None

                target_cols = [c for c in ["성별","연령","SEG"] if c in h.columns]
                target_text = ""
                target_amt = None
                if target_cols:
                    tg = (
                        h.groupby(target_cols, dropna=False)["_amt"]
                        .mean()
                        .reset_index()
                        .sort_values("_amt", ascending=False)
                    )
                    if not tg.empty:
                        r0 = tg.iloc[0]
                        bits = []
                        if "성별" in target_cols:
                            v = _clean_text_value(r0["성별"])
                            if v: bits.append(v)
                        if "연령" in target_cols:
                            v = _clean_text_value(clean_identifier_value(r0["연령"]))
                            if v: bits.append(v)
                        if "SEG" in target_cols:
                            v = _clean_text_value(r0["SEG"])
                            if v: bits.append(v)
                        target_text = " ".join(bits)
                        target_amt = float(r0["_amt"])

                results.append({
                    "scope": scope_name,
                    "group": group_name,
                    "product": str(pname),
                    "count": n,
                    "avg_amt": avg_amt,
                    "max_amt": max_amt,
                    "ge3": ge3,
                    "ge5": ge5,
                    "avg_price": avg_price,
                    "hp3_price": hp3_price,
                    "hp5_price": hp5_price,
                    "target": target_text,
                    "target_avg": target_amt,
                })

        if results:
            return sorted(results, key=lambda x: (x["ge5"], x["ge3"], x["avg_amt"]), reverse=True)

    return []




def _normalize_season_group(product_name: str, current_group: str) -> str:
    """상품명 기반 시즌 상품군 보정. 명확한 상품 키워드를 일반 키워드보다 우선."""
    name = _clean_text_value(product_name).lower()
    if any(k in name for k in ["우양산", "양산", "우산"]):
        return "우양산"
    if any(k in name for k in ["캐리어", "여행가방", "트렁크"]):
        return "캐리어·여행용품"
    if any(k in name for k in ["선스틱", "선크림", "선쿠션", "자외선"]):
        return "선케어"
    if any(k in name for k in ["냉감", "쿨링", "드라이셀", "흡습속건"]):
        return "냉감·기능성 의류"
    if any(k in name for k in ["선풍기", "서큘", "써큘", "에어써큘", "냉풍기"]):
        return "냉방가전"
    return current_group

def _season_specific_action(group_name: str) -> str:
    actions = {
        "우양산": "1만원 내외 가격대의 경량·휴대용·암막 기능을 갖춘 우양산",
        "캐리어·여행용품": "기내용·경량·확장형 등 휴가 수요가 명확한 캐리어·여행용품",
        "여행용 소형가전": "휴대성·소형·멀티전압 등 여행 편의성이 명확한 소형가전",
        "선케어": "휴대성·간편 도포·높은 자외선 차단 지수를 갖춘 선스틱·선쿠션·선크림",
        "냉감·기능성 의류": "냉감·흡습속건·통기성 기능이 명확한 여름 기능성 의류",
        "냉방가전": "리모컨·저소음·공기순환·공간효율 등 사용 편의성이 강화된 선풍기·서큘레이터",
        "보양식·간편식": "간편 조리·대중성·가격 경쟁력이 검증된 여름 보양식·간편식",
    }
    return actions.get(group_name, "동일 시즌 수요와 기능성이 명확한 상품")



def _marketing_calendar_reason(group_name: str, ref_date=None) -> str:
    """월별/계절 마케팅 캘린더 기반 '왜 지금인가' 근거."""
    try:
        month = int(pd.Timestamp(ref_date).month) if ref_date is not None else int(pd.Timestamp.today().month)
    except Exception:
        month = int(pd.Timestamp.today().month)

    reasons = {
        1: {
            "default": "신년·설 준비 수요가 확대되는 1월",
            "건강": "신년 건강관리와 설 선물 수요가 함께 확대되는 1월",
        },
        2: {"default": "설 이후 신학기·봄 준비 수요가 시작되는 2월"},
        3: {"default": "신학기·입학·이사와 봄맞이 수요가 집중되는 3월"},
        4: {"default": "환절기·봄나들이·야외활동 수요가 확대되는 4월"},
        5: {"default": "가정의달 선물·나들이 수요가 집중되는 5월"},
        6: {
            "default": "초여름·장마 준비와 휴가 사전 수요가 시작되는 6월",
            "우양산": "장마 시작과 자외선 차단 수요가 함께 커지는 6월",
            "냉방가전": "기온 상승으로 냉방가전 수요가 본격화되는 6월",
        },
        7: {
            "default": "장마·폭염·여름휴가 수요가 동시에 집중되는 7월",
            "우양산": "장마·폭염이 겹쳐 우천 대응과 자외선 차단 수요가 동시에 커지는 7월",
            "냉방가전": "폭염·열대야로 냉방가전 수요가 집중되는 7월",
            "캐리어·여행용품": "본격적인 여름휴가·여행 수요가 확대되는 7월",
            "선케어": "휴가·야외활동과 강한 자외선으로 선케어 수요가 집중되는 7월",
            "냉감·기능성 의류": "폭염과 야외활동 증가로 냉감·흡습속건 의류 수요가 커지는 7월",
            "보양식·간편식": "초복·중복 등 보양식 수요가 집중되는 7월",
        },
        8: {"default": "폭염·휴가 후반과 개학 준비 수요가 이어지는 8월"},
        9: {"default": "추석 선물·귀성 및 환절기 수요가 확대되는 9월"},
        10: {"default": "가을 나들이·환절기·겨울 준비 수요가 시작되는 10월"},
        11: {"default": "겨울 준비·김장·연말 쇼핑 수요가 확대되는 11월"},
        12: {"default": "한파·크리스마스·연말 선물 수요가 집중되는 12월"},
    }
    month_map = reasons.get(month, {})
    return month_map.get(group_name) or month_map.get("default") or f"{month}월 시즌 수요가 형성되는 시기"


def _season_single_or_repeat_sentence(x: dict) -> str:
    name = _safe_product_label(x["product"])
    subject = _with_topic(name)
    scope = x["scope"]
    price_bits = []
    if int(x.get("count", 0)) == 1:
        one_price = x.get("avg_price")
        if one_price is None:
            one_price = x.get("hp5_price") if x.get("hp5_price") is not None else x.get("hp3_price")
        if one_price is not None:
            price_bits.append(f"당시 혜택가 {one_price:,.0f}원")
    else:
        if x.get("avg_price") is not None:
            label = "전년 평균 혜택가" if scope.startswith("전년") else "과거 시즌 평균 혜택가"
            price_bits.append(f"{label} {x['avg_price']:,.0f}원")
        if x.get("hp5_price") is not None:
            price_bits.append(f"500만원 이상 고성과 회차 평균 혜택가 {x['hp5_price']:,.0f}원")
        elif x.get("hp3_price") is not None:
            price_bits.append(f"300만원 이상 고성과 회차 평균 혜택가 {x['hp3_price']:,.0f}원")

    target = _clean_text_value(x.get("target"))
    target_part = ""
    if target:
        target_part = f", 주요 고성과 타겟 {target}"
        if x.get("target_avg") is not None:
            target_part += f" 평균 {compact_money(x['target_avg'])}"

    price_part = f", {', '.join(price_bits)}" if price_bits else ""
    season_group = _normalize_season_group(x["product"], x["group"])
    action_product = _season_specific_action(season_group)
    calendar_reason = _marketing_calendar_reason(season_group, x.get("ref_date"))

    if int(x["count"]) == 1:
        # 1회 성과는 '검증'이 아니라 '고성과 사례'로만 표현
        return (
            f"{subject} {scope} 1회 운영에서 {compact_money(x['max_amt'])}을 기록한 고성과 사례이며{price_part}"
            f"{target_part}. 단일 운영 사례인 만큼 반복 성과가 검증된 상품으로 단정할 수는 없으나, "
            f"{calendar_reason} 시즌 수요와 과거 고성과가 함께 확인된 만큼 당시와 유사한 가격 조건을 확보한 "
            f"{action_product}의 신규·유사신규 TEST를 검토할 필요가 있습니다."
        )

    # 2회 이상: 반복 성과 수준을 수치로 구분
    if x["ge5"] >= 2 or (x["ge3"] >= 2 and x["ge3"] / max(x["count"], 1) >= 0.5):
        return (
            f"{subject} {scope} {x['count']}회 운영 중 300만원 이상 {x['ge3']}회"
            + (f", 500만원 이상 {x['ge5']}회" if x["ge5"] else "")
            + f", 평균 {compact_money(x['avg_amt'])}, 최고 {compact_money(x['max_amt'])}{target_part}{price_part}의 성과를 기록해 "
              f"동시즌 반복 성과가 확인됐습니다. 당시와 유사한 가격 조건 확보 시 동일 상품 재운영을 우선 검토하고, "
              f"{action_product}으로 신규·유사신규 TEST를 확장하는 것이 적절합니다."
        )

    return (
        f"{subject} {scope} {x['count']}회 운영 중 300만원 이상 {x['ge3']}회, 평균 {compact_money(x['avg_amt'])}, "
        f"최고 {compact_money(x['max_amt'])}{target_part}{price_part}의 성과가 확인됐습니다. 반복 고성과로 단정하기에는 표본이 제한적이므로 "
        f"동일 상품 또는 {action_product}을 추가 TEST해 재현 여부를 확인하는 것이 적절합니다."
    )

def _seasonal_action_sentence(products_all: pd.DataFrame, week_end):
    """실제 시즌 고성과 상품을 1회 성공 사례와 반복 검증 사례로 구분해 구체적 상품 조건까지 제안."""
    items = _seasonal_last_year_evidence(products_all, week_end)
    if not items:
        return None

    selected = []
    used_groups = set()
    for x in items:
        if x["group"] in used_groups:
            continue
        selected.append(x)
        used_groups.add(x["group"])
        if len(selected) >= 2:
            break

    if not selected:
        return None

    _season_lines = []
    _season_seen = set()
    for x in selected:
        _line = "• " + _clean_seg_display_text(_season_single_or_repeat_sentence(x))
        _key = re.sub(r"\s+", " ", _line).strip()
        if _key not in _season_seen:
            _season_seen.add(_key)
            _season_lines.append(_line)
    return "\n".join(_season_lines)



def _md_recommendation_tables(products_all: pd.DataFrame, week_df: pd.DataFrame, week_end):
    """MD 의사결정용: 재편성 추천 / 신규·유사신규 소싱 제안 데이터."""
    rec_rows = []
    sourcing_rows = []

    if products_all is None or products_all.empty:
        return pd.DataFrame(), pd.DataFrame()

    pcol = first_col(products_all, ["상품명", "MMS 상품명", "상품"])
    acol = first_col(products_all, ["주문금액", "거래액", "매출"])
    dcol = first_col(products_all, ["발송일", "발송일자", "일자", "날짜"])
    price_col = first_col(products_all, ["멤버십 혜택가", "멤버십혜택가", "혜택가", "최종혜택가", "행사가", "판매가", "MMS혜택가", "MMS 혜택가", "실판매가"])
    if not pcol or not acol:
        return pd.DataFrame(), pd.DataFrame()

    tmp = products_all.copy()
    tmp[acol] = pd.to_numeric(tmp[acol], errors="coerce").fillna(0)
    if dcol:
        tmp[dcol] = pd.to_datetime(tmp[dcol], errors="coerce")

    _week_product_col = first_col(week_df, ["상품명", "MMS 상품명", "상품"]) if week_df is not None else None
    _week_products = set(week_df[_week_product_col].dropna().astype(str)) if _week_product_col and not week_df.empty else set()

    tmp = _attach_product_master_keys(tmp)
    for _master_key, g in tmp.groupby("_product_master_key", dropna=True):
        pname = str(g[pcol].dropna().iloc[-1]) if pcol in g.columns and g[pcol].notna().any() else str(_master_key)
        if not _clean_text_value(pname):
            continue
        if _week_products and str(pname) not in _week_products:
            continue
        cnt = len(g)
        ge3 = int((g[acol] >= 3_000_000).sum())
        ge5 = int((g[acol] >= 5_000_000).sum())
        avg_amt = float(g[acol].mean()) if cnt else 0
        max_amt = float(g[acol].max()) if cnt else 0
        last_date = g[dcol].max() if dcol else pd.NaT
        days = None
        if dcol and pd.notna(last_date):
            try:
                days = max(0, (pd.Timestamp(week_end).normalize() - pd.Timestamp(last_date).normalize()).days)
            except Exception:
                days = None
        if ge5 >= 1 or ge3 >= 2:
            price_txt = ""
            if price_col:
                vals = pd.to_numeric(g[price_col], errors="coerce").dropna()
                if not vals.empty:
                    price_txt = f"{vals.mean():,.0f}원"
            rec_rows.append({
                "상품": _safe_product_label(pname),
                "운영횟수": cnt,
                "300만원↑": ge3,
                "500만원↑": ge5,
                "평균매출": compact_money(avg_amt),
                "최고매출": compact_money(max_amt),
                "최근 미편성": f"{days}일" if days is not None and days > 0 else "-",
                "과거 평균 혜택가": price_txt or "-"
            })

    season_items = _seasonal_last_year_evidence(products_all, week_end)
    for x in season_items[:10]:
        sg = _normalize_season_group(x["product"], x["group"])
        reason = _marketing_calendar_reason(sg, week_end)
        price = x.get("avg_price") or x.get("hp5_price") or x.get("hp3_price")
        sourcing_rows.append({
            "시즌/상품군": sg,
            "과거 고성과 사례": _safe_product_label(x["product"]),
            "성과": f"{x['scope']} {x['count']}회 / 평균 {compact_money(x['avg_amt'])} / 최고 {compact_money(x['max_amt'])}",
            "당시 가격": f"{price:,.0f}원" if price is not None else "-",
            "왜 지금": reason,
            "소싱 방향": _season_specific_action(sg)
        })

    rec_df = pd.DataFrame(rec_rows)
    if not rec_df.empty:
        rec_df = rec_df.sort_values(["500만원↑", "300만원↑", "운영횟수"], ascending=False).head(15)
    src_df = pd.DataFrame(sourcing_rows).drop_duplicates(subset=["시즌/상품군", "과거 고성과 사례"]).head(10)
    return rec_df, src_df



def _extract_product_from_insight_line(line: str) -> str:
    m = re.match(r"\s*\[([^\]]+)\]", str(line))
    return _clean_text_value(m.group(1)) if m else ""

def _consolidate_product_detail_insights(detail_text: str, week_df: pd.DataFrame, all_products: pd.DataFrame) -> str:
    """
    동일 상품의 발송 건별 상세 인사이트를 상품별 1개로 통합.
    기존 문장에 의존하지 않고 실제 금주/누적 데이터로 핵심 판정을 재구성한다.
    """
    if not detail_text or week_df is None or getattr(week_df, "empty", True):
        return detail_text

    pcol = first_col(week_df, ["상품명", "MMS 상품명", "상품"])
    acol = first_col(week_df, ["주문금액", "거래액", "매출"])
    dcol = first_col(week_df, ["발송일", "발송일자", "일자", "날짜"])
    target_col = first_col(week_df, ["타겟", "성별연령", "성별/연령", "성연령"])
    seg_col = first_col(week_df, ["SEG", "세그", "세그먼트"])
    price_col = first_col(week_df, ["멤버십 혜택가", "멤버십혜택가", "혜택가", "최종혜택가", "행사가", "판매가"])
    lowest_col = first_col(week_df, ["발송일 비교 최저가", "비교최저가", "최저가", "네이버최저가"])

    if not pcol or not acol:
        return detail_text

    w = _attach_product_master_keys(week_df.copy())
    w[acol] = pd.to_numeric(w[acol], errors="coerce").fillna(0)
    if dcol:
        w[dcol] = pd.to_datetime(w[dcol], errors="coerce")

    # Preserve product order from the original detail text where possible.
    original_order = []
    for ln in str(detail_text).splitlines():
        pn = _extract_product_from_insight_line(ln)
        if pn and pn not in original_order:
            original_order.append(pn)
    for pn in w[pcol].dropna().astype(str).unique():
        if pn not in original_order:
            original_order.append(pn)

    out = []
    # 상품명 순서 대신 master key 기준으로 통합
    _master_order = []
    if "_product_master_key" in w.columns:
        for _mk in w["_product_master_key"].dropna().astype(str):
            if _mk not in _master_order:
                _master_order.append(_mk)

    for _master_key in _master_order:
        g = w[w["_product_master_key"].astype(str) == str(_master_key)].copy()
        if g.empty:
            continue
        pname = str(g[pcol].dropna().iloc[-1]) if pcol in g.columns and g[pcol].notna().any() else str(_master_key)

        g = g.sort_values(dcol) if dcol else g
        week_total = float(g[acol].sum())
        week_max = float(g[acol].max())
        week_count = len(g)
        ge3 = int((g[acol] >= 3_000_000).sum())
        ge5 = int((g[acol] >= 5_000_000).sum())

        # Cumulative history
        hist = pd.DataFrame()
        if all_products is not None and not getattr(all_products, "empty", True):
            apcol = first_col(all_products, ["상품명", "MMS 상품명", "상품"])
            aacol = first_col(all_products, ["주문금액", "거래액", "매출"])
            if apcol and aacol:
                _all_master = _attach_product_master_keys(all_products.copy())
                hist = _all_master[_all_master["_product_master_key"].astype(str) == str(_master_key)].copy()
                hist[aacol] = pd.to_numeric(hist[aacol], errors="coerce").fillna(0)
        hist_count = len(hist)
        hist_avg = float(hist[aacol].mean()) if not hist.empty else None
        hist_max = float(hist[aacol].max()) if not hist.empty else None

        parts = []
        subject = f"[{_safe_product_label(pname)}]"

        # 1. Current week result
        if week_count > 1:
            parts.append(
                f"금주 {week_count}회 편성에서 누적 {compact_money(week_total)}, "
                f"회차 최고 {compact_money(week_max)}을 기록했으며 300만원 이상 {ge3}회"
                + (f", 500만원 이상 {ge5}회" if ge5 else "") + "의 성과가 확인됐습니다."
            )
        else:
            parts.append(f"금번 {compact_money(week_total)}을 기록했습니다.")

        # 2. Historical meaning
        if hist_count >= 2 and hist_avg is not None:
            if hist_max is not None and week_max >= hist_max - 1:
                parts.append(f"누적 {hist_count}회 평균 {compact_money(hist_avg)} 대비 금번 회차 최고 {compact_money(week_max)}으로 역대 최고 수준의 성과를 기록했습니다.")
            elif week_total >= 5_000_000:
                parts.append(f"누적 {hist_count}회 평균 {compact_money(hist_avg)}을 기록한 반복 운영 상품으로, 금주에도 핵심 매출 기여 수준의 성과를 확보했습니다.")

        # 3. Target/SEG comparison in current week
        group_cols = [c for c in [target_col, seg_col] if c]
        target_summary = None
        if group_cols and len(g) >= 2:
            tg = g.groupby(group_cols, dropna=False)[acol].agg(["count","mean","sum"]).reset_index().sort_values("mean", ascending=False)
            if len(tg) >= 2:
                top, second = tg.iloc[0], tg.iloc[1]
                def _label(row):
                    vals = []
                    for c in group_cols:
                        v = _clean_text_value(row[c])
                        if v and v.lower() != "nan":
                            vals.append(v)
                    return " ".join(vals)
                l1, l2 = _label(top), _label(second)
                if l1 and l2 and float(second["mean"]) > 0:
                    ratio = float(top["mean"]) / float(second["mean"])
                    if ratio >= 1.5:
                        target_summary = (
                            f"{l1} 평균 {compact_money(float(top['mean']))}으로 "
                            f"{l2} 평균 {compact_money(float(second['mean']))} 대비 {ratio:.1f}배 높아 "
                            f"{l1}의 상대적 강세가 확인됩니다."
                        )
        if target_summary:
            parts.append(target_summary)

        # 4. Actual repeat trend, no generic "check later" sentence.
        if week_count >= 3:
            vals = g[acol].tolist()
            seq = " → ".join(compact_money(float(v)) for v in vals[-3:])
            if vals[-3] > vals[-2] > vals[-1]:
                parts.append(f"최근 3회 주문금액은 {seq}으로 2회 연속 하락해 동일 조건의 즉시 반복 편성보다 운영 간격 또는 타겟 전환이 필요합니다.")
            elif vals[-3] < vals[-2] < vals[-1]:
                parts.append(f"최근 3회 주문금액은 {seq}으로 연속 상승해 현재까지 반복 운영에 따른 성과 둔화는 확인되지 않습니다.")
            else:
                parts.append(f"최근 3회 주문금액은 {seq}으로 연속 하락 흐름은 확인되지 않아 단기 재편성 여력이 있습니다.")

        # 5. Price evidence only when actual columns exist.
        if price_col:
            prices = pd.to_numeric(g[price_col], errors="coerce").dropna()
            lows = pd.to_numeric(g[lowest_col], errors="coerce").dropna() if lowest_col else pd.Series(dtype=float)
            if not prices.empty and not lows.empty:
                curp, lowp = float(prices.iloc[-1]), float(lows.iloc[-1])
                if curp <= lowp:
                    parts.append(f"멤버십 혜택가 {curp:,.0f}원으로 발송일 비교 최저가 {lowp:,.0f}원 이하의 가격 경쟁력을 확보했습니다.")
                else:
                    parts.append(f"멤버십 혜택가 {curp:,.0f}원으로 발송일 비교 최저가 {lowp:,.0f}원 대비 가격 우위가 없어 가격 조건 재점검이 필요합니다.")

        # 6. First test terminology: product / target / SEG are explicitly separated.
        if hist_count == 1:
            parts.append("MMS 첫 운영 상품으로 추가 운영을 통해 재현성 검증이 필요합니다.")
        elif target_col:
            # Compare current target against prior history if possible.
            apcol = first_col(all_products, ["상품명", "MMS 상품명", "상품"]) if all_products is not None else None
            atcol = first_col(all_products, ["타겟", "성별연령", "성별/연령", "성연령"]) if all_products is not None else None
            aseg = first_col(all_products, ["SEG", "세그", "세그먼트"]) if all_products is not None else None
            if apcol and atcol and not hist.empty:
                current_targets = g[target_col].dropna().astype(str).unique().tolist()
                for ct in current_targets:
                    hct = hist[hist[atcol].astype(str) == str(ct)] if atcol in hist.columns else pd.DataFrame()
                    if len(hct) == 1:
                        parts.append(f"{ct} 타겟 첫 TEST 성과로, 동일 타겟 1회 추가 검증 후 확대 여부를 판단하는 것이 적절합니다.")
                        break

        # 7. One final action, aligned with evidence.
        action = ""
        if week_count >= 3:
            vals = g[acol].tolist()
            declining = vals[-3] > vals[-2] > vals[-1]
        else:
            declining = False

        if target_summary:
            # pull top label from sentence
            top_label = target_summary.split(" 평균 ")[0]
        else:
            top_label = ""

        if declining:
            action = "다음 운영 제안: 최근 연속 하락이 확인된 동일 조건의 즉시 반복 편성은 보수적으로 운영하고, 최근 미발송 타겟·SEG 전환 또는 일정 기간 미편성 후 재검토하는 것이 적절합니다."
        elif week_total >= 5_000_000 or ge5 >= 1:
            if top_label:
                action = f"다음 운영 제안: {top_label}을 우선 재편성 대상으로 검토하고, 동일 조건 1회 추가 검증 후 미발송 SEG 확대 여부를 판단하는 것이 적절합니다."
            else:
                action = "다음 운영 제안: 금주 고성과 조건을 우선 유지해 단기 재편성을 검토하되, 동일 SEG 과다 반복은 피하고 미발송 SEG 확대 여부를 함께 검토하는 것이 적절합니다."
        else:
            action = "다음 운영 제안: 금주 성과와 과거 평균을 함께 비교해 재편성 우선순위를 결정하고, 근거가 충분하지 않은 타겟 확대는 추가 TEST 후 판단하는 것이 적절합니다."

        parts.append(action)
        out.append(subject + " " + " > ".join(parts))

    return "\n".join(out)


def _normalize_reco_product_key(name: str) -> str:
    """DEPRECATED: 추천 중복 판정은 상품 마스터키 기준으로 전환. 이름 유사도 병합 금지."""
    return _clean_text_value(name)

def _extract_bracket_or_subject_product(sentence: str) -> str:
    s = _clean_text_value(sentence)
    m = re.search(r"\[([^\]]+)\]", s)
    if m:
        return _clean_text_value(m.group(1))
    # bullet 뒤 첫 서술부를 상품명 후보로 사용
    s = re.sub(r"^[•\-\s]+", "", s)
    for token in [" 상품은 ", "은 ", "는 ", "이 ", "가 "]:
        if token in s:
            return _clean_text_value(s.split(token, 1)[0])
    return ""

def _recommendation_priority(sentence: str) -> int:
    """같은 상품이 여러 추천 규칙에 걸릴 때 가장 정보량 높은 문장을 우선."""
    s = str(sentence)
    score = 0
    if "동시즌" in s: score += 8
    if "최근" in s and "미편성" in s: score += 5
    if "누적" in s: score += 4
    if "500만원 이상" in s: score += 4
    if "평균" in s and "최고" in s: score += 3
    if "혜택가" in s or "가격" in s: score += 3
    if "타겟" in s or "SEG" in s: score += 2
    if "신규·유사신규" in s: score += 1
    return score

def _clean_seg_display_text(s: str) -> str:
    """SEG 숫자가 1.0/2.0/3.0으로 노출되는 문제 및 조사 오류 보정."""
    s = str(s)
    s = re.sub(r"\b(남성|여성)\s*(3040|5060)\s+([123])\.0\b", r"\1 \2 SEG\3", s)
    s = re.sub(r"\bSEG\s*([123])\.0\b", r"SEG\1", s, flags=re.I)
    s = re.sub(r"선풍기·서큘레이터으로", "선풍기·서큘레이터로", s)
    s = re.sub(r"선풍기·서큘레이터을", "선풍기·서큘레이터를", s)
    return s

def _dedupe_next_week_recommendations(points):
    """
    모든 주차 공통 적용.
    - 완전 동일 문장 제거
    - 동일 상품이 시즌/미편성/재편성 규칙에 중복 포착되면 1개만 유지
    - 시즌+미편성은 정보량 높은 시즌 문장을 우선
    """
    if not points:
        return points

    # 문자열/리스트 모두 지원
    was_string = isinstance(points, str)
    rows = [x.strip() for x in str(points).splitlines() if x.strip()] if was_string else [str(x).strip() for x in points if str(x).strip()]

    # exact dedupe
    exact = []
    seen = set()
    for r in rows:
        c = _clean_seg_display_text(r)
        key = re.sub(r"\s+", " ", c).strip()
        if key not in seen:
            seen.add(key)
            exact.append(c)

    # product-level dedupe
    product_best = {}
    non_product = []
    order = []
    for r in exact:
        pname = _extract_bracket_or_subject_product(r)
        if not pname:
            non_product.append(r)
            continue
        key = _normalize_reco_product_key(pname)
        if not key:
            non_product.append(r)
            continue
        if key not in product_best:
            product_best[key] = r
            order.append(key)
        else:
            old = product_best[key]
            if _recommendation_priority(r) > _recommendation_priority(old):
                product_best[key] = r

    merged = [product_best[k] for k in order] + non_product

    # second exact pass after cleanup
    final, seen2 = [], set()
    for r in merged:
        k = re.sub(r"\s+", " ", r).strip()
        if k not in seen2:
            seen2.add(k)
            final.append(r)

    return "\n".join(final) if was_string else final


def _style_total_row(df: pd.DataFrame):
    """총합계 별도 배경색 없이 기본 표 스타일 유지."""
    return df


def _normalize_core_product_name(name: str) -> str:
    """
    코드 변경 전후 동일 상품 판정을 위한 핵심 상품명 정규화.
    단순 접두어/프로모션 문구만 제거하며 모델/용량/수량/본품 구성은 유지한다.
    """
    s = _clean_text_value(name).lower()
    s = re.sub(r"\[(?:m|단독|무료배송|쇼라[^\]]*|특별혜택가)\]", " ", s)
    s = re.sub(r"★\s*단독", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def _extract_model_tokens(name: str):
    """모델명 후보 추출. 영문+숫자 조합을 우선 식별."""
    s = _clean_text_value(name).upper()
    toks = re.findall(r"\b[A-Z]{1,6}[-_/]?[A-Z0-9]{2,}\b", s)
    # 지나치게 일반적인 토큰 제거
    bad = {"MMS", "NEW", "SET", "EA", "ML", "KG"}
    return tuple(sorted(set(t for t in toks if t not in bad)))


def _extract_quantity_signature(name: str):
    """
    용량/수량/구성 시그니처.
    120포, 3통, 500ml, 20캔, 36개, 1+1, 3+3 등 핵심 구성 차이를 유지한다.
    """
    s = _clean_text_value(name).lower()
    pats = [
        r"\d+(?:\.\d+)?\s*(?:ml|l|g|kg|포|통|캔|개|매|롤|팩|봉|병|입|개월|인치)",
        r"\d+\s*[x×]\s*\d+",
        r"\d+\s*\+\s*\d+",
    ]
    vals = []
    for p in pats:
        vals.extend(re.findall(p, s))
    return tuple(sorted(set(re.sub(r"\s+", "", v) for v in vals)))


def _extract_gift_signature(name: str):
    """
    증정/추가 구성 식별.
    증정품이 달라도 본품은 동일할 수 있으므로 상품 마스터는 연결 가능하되,
    운영 구성은 별도 variant로 기록한다.
    """
    s = _clean_text_value(name).lower()
    gift_markers = ["증정", "추가증정", "+", "보조배터리", "리필", "사은품"]
    hits = [m for m in gift_markers if m in s]
    return tuple(sorted(set(hits)))


def _get_product_code_columns(df: pd.DataFrame):
    shora = first_col(df, ["쇼라코드", "쇼핑라운지코드", "상품코드", "샵바이코드"])
    alpha = first_col(df, ["알파코드", "알파상품코드"])
    return shora, alpha


def _product_master_key_from_row(row, shora_col=None, alpha_col=None, name_col=None):
    """
    상품 마스터키 판정:
    1) 동일 코드면 동일상품
    2) 코드가 달라도 모델명 + 핵심 상품명 + 용량/수량 시그니처가 동일하면 과거 이력 연결
    3) 구성/모델/용량이 다르면 별도 상품
    """
    name = _clean_text_value(row.get(name_col, "")) if name_col else ""
    core = _normalize_core_product_name(name)
    models = _extract_model_tokens(name)
    qty = _extract_quantity_signature(name)

    shora = _clean_text_value(row.get(shora_col, "")) if shora_col else ""
    alpha = _clean_text_value(row.get(alpha_col, "")) if alpha_col else ""

    # 코드가 있으면 우선 코드 기반 키
    if shora:
        code_key = f"SHORA:{shora}"
    elif alpha:
        code_key = f"ALPHA:{alpha}"
    else:
        code_key = ""

    # 실질 동일성 비교용 fingerprint
    # 모델명이 있으면 모델+수량을 강하게 사용
    if models:
        fp = f"MODEL:{'|'.join(models)}::QTY:{'|'.join(qty)}"
    else:
        # 모델명이 없으면 핵심명+수량
        fp = f"NAME:{core}::QTY:{'|'.join(qty)}"

    return {
        "code_key": code_key,
        "fingerprint": fp,
        "core_name": core,
        "models": models,
        "qty": qty,
        "gift_sig": _extract_gift_signature(name),
    }


def _attach_product_master_keys(df: pd.DataFrame):
    """
    각 행에 _product_master_key / _product_variant_key 추가.
    코드가 바뀌었더라도 fingerprint가 동일하면 동일 master로 묶는다.
    증정/추가구성은 variant에서 분리한다.
    """
    if df is None or df.empty:
        return df

    out = df.copy()
    name_col = first_col(out, ["상품명", "MMS 상품명", "상품"])
    if not name_col:
        out["_product_master_key"] = ""
        out["_product_variant_key"] = ""
        return out

    shora_col, alpha_col = _get_product_code_columns(out)

    meta = []
    for _, row in out.iterrows():
        meta.append(_product_master_key_from_row(row, shora_col, alpha_col, name_col))

    tmp = pd.DataFrame(meta, index=out.index)

    # 동일 fingerprint가 여러 코드에 걸쳐 있으면 같은 master로 연결
    fp_to_master = {}
    for idx, r in tmp.iterrows():
        fp = r["fingerprint"]
        code = r["code_key"]
        if fp not in fp_to_master:
            fp_to_master[fp] = code or fp

    out["_product_master_key"] = [fp_to_master.get(r["fingerprint"], r["code_key"] or r["fingerprint"]) for _, r in tmp.iterrows()]
    out["_product_variant_key"] = [
        f"{out.loc[idx, '_product_master_key']}::GIFT:{'|'.join(r['gift_sig'])}"
        for idx, r in tmp.iterrows()
    ]
    return out


def _same_product_master(df: pd.DataFrame, row_a_idx, row_b_idx) -> bool:
    d = _attach_product_master_keys(df)
    return _clean_text_value(d.loc[row_a_idx, "_product_master_key"]) == _clean_text_value(d.loc[row_b_idx, "_product_master_key"])


def _sentence_product_master_key(sentence: str, products_all: pd.DataFrame):
    """
    추천 문장에 포함된 실제 상품명을 찾아 상품 마스터키 반환.
    상품명 유사도 병합은 사용하지 않고, 실제 데이터에 존재하는 상품명 일치 → master key 매핑만 사용.
    """
    if products_all is None or products_all.empty:
        return "", ""

    d = _attach_product_master_keys(products_all.copy())
    pcol = first_col(d, ["상품명", "MMS 상품명", "상품"])
    if not pcol or "_product_master_key" not in d.columns:
        return "", ""

    s = str(sentence)
    candidates = []
    for _, r in d[[pcol, "_product_master_key"]].dropna(subset=[pcol]).drop_duplicates().iterrows():
        pname = _clean_text_value(r[pcol])
        if pname and pname in s:
            candidates.append((len(pname), pname, _clean_text_value(r["_product_master_key"])))

    if not candidates:
        return "", ""
    candidates.sort(reverse=True)
    _, pname, key = candidates[0]
    return key, pname


def _extract_recent_unassigned_days(sentence: str):
    m = re.search(r"최근\s*(\d+)일간\s*미편성", str(sentence))
    return int(m.group(1)) if m else None


def _merge_same_product_recommendations(sentences, products_all: pd.DataFrame):
    """
    모든 주차 공통:
    동일 상품 마스터가 '즉시 재편성 / 최근 미편성 / 동시즌 / 시즌성'에 중복 등장하면 1개로 통합.
    코드가 다르더라도 실질 동일 product master면 통합하고, 다른 master는 절대 합치지 않음.
    """
    if not sentences:
        return sentences

    groups = {}
    no_product = []
    order = []

    for raw in sentences:
        s = _clean_seg_display_text(str(raw).strip())
        if not s:
            continue
        key, pname = _sentence_product_master_key(s, products_all)
        if not key:
            no_product.append(s)
            continue
        if key not in groups:
            groups[key] = []
            order.append(key)
        groups[key].append((s, pname))

    merged = []
    for key in order:
        items = groups[key]
        if len(items) == 1:
            merged.append(items[0][0])
            continue

        # 가장 근거가 풍부한 문장을 기본으로 선택
        def score(x):
            s = x[0]
            val = 0
            if "동시즌" in s: val += 10
            if "300만원 이상" in s: val += 4
            if "500만원 이상" in s: val += 5
            if "평균" in s and "최고" in s: val += 4
            if "혜택가" in s or "가격" in s: val += 4
            if "타겟" in s or "SEG" in s: val += 3
            if "미편성" in s: val += 2
            return val

        items = sorted(items, key=score, reverse=True)
        base_sentence = items[0][0]

        # 다른 중복 문장에서 최근 미편성 기간만 유용한 근거로 흡수
        days = None
        for s, _ in items:
            d = _extract_recent_unassigned_days(s)
            if d is not None:
                days = max(days or 0, d)

        if days is not None and "미편성" not in base_sentence:
            # 마지막 액션 문장 전에 자연스럽게 삽입
            marker_candidates = [
                "현재 판매 가능 여부",
                "당시와 유사한 가격 조건",
                "동일 상품 재운영",
                "신규·유사신규 TEST",
            ]
            inserted = False
            for marker in marker_candidates:
                pos = base_sentence.find(marker)
                if pos > 0:
                    prefix = base_sentence[:pos].rstrip()
                    suffix = base_sentence[pos:]
                    if prefix.endswith("."):
                        base_sentence = prefix + f" 최근 {days}일간 미편성된 상태이며, " + suffix
                    else:
                        base_sentence = prefix + f", 최근 {days}일간 미편성된 상태이며, " + suffix
                    inserted = True
                    break
            if not inserted:
                base_sentence += f" 최근 {days}일간 미편성된 상태입니다."

        merged.append(base_sentence)

    # 상품을 특정하지 않는 카테고리/신규소싱 제안은 유지하되 완전 동일 문장 제거
    merged.extend(no_product)
    final, seen = [], set()
    for s in merged:
        k = re.sub(r"\s+", " ", s).strip()
        if k not in seen:
            seen.add(k)
            final.append(s)
    return final

def _weekly_table_title(title: str):
    """주간실적의 표 제목을 카드 폭 기준 가운데 정렬합니다."""
    st.markdown(
        (
            "<div style='width:100%; text-align:center; "
            "font-weight:700; font-size:1.05rem; "
            "margin:0.45rem 0 0.35rem 0;'>"
            f"{title}</div>"
        ),
        unsafe_allow_html=True,
    )



def _style_weekly_category_total(df: pd.DataFrame):
    """대/중카테고리 총합계: 배경색을 건드리지 않고 Bold만 적용."""
    def _row_style(row):
        values = [_clean_text_value(v) for v in row.tolist()]
        if any(v == "총합계" for v in values):
            # 배경색 지정 금지: Streamlit 기본 흰 배경 유지
            return ["font-weight: 800 !important;" for _ in row]
        return ["" for _ in row]

    try:
        return df.style.apply(_row_style, axis=1)
    except Exception:
        return df



def _get_query_param(name: str, default: str = "") -> str:
    """Streamlit query param을 문자열로 안전하게 읽습니다."""
    try:
        value = st.query_params.get(name, default)
        if isinstance(value, list):
            return str(value[-1]) if value else default
        return str(value) if value is not None else default
    except Exception:
        return default


def _set_weekly_deeplink(year: int, week: str) -> None:
    """현재 주간실적 선택값을 URL에 반영합니다."""
    try:
        st.query_params["menu"] = "weekly"
        st.query_params["year"] = str(int(year))
        st.query_params["week"] = str(week)
    except Exception:
        pass


def _get_secret_value(*names):
    """Streamlit Secrets → 환경변수 순으로 안전하게 인증값 조회."""
    for name in names:
        try:
            value = st.secrets.get(name)
            if value:
                return str(value).strip()
        except Exception:
            pass
        value = os.getenv(name)
        if value:
            return str(value).strip()
    return None


def _naver_trend_credentials():
    """
    NAVER API HUB 인증정보.
    권장 Secrets:
      NAVER_API_HUB_CLIENT_ID
      NAVER_API_HUB_CLIENT_SECRET
    """
    cid = _get_secret_value(
        "NAVER_API_HUB_CLIENT_ID",
        "NAVER_API_CLIENT_ID",
    )
    secret = _get_secret_value(
        "NAVER_API_HUB_CLIENT_SECRET",
        "NAVER_API_CLIENT_SECRET",
    )
    return cid, secret


def _naver_trend_seed_catalog(month: int):
    """
    Shopping Insight는 '임의의 실시간 인기 키워드 목록'을 반환하는 API가 아니라
    지정한 키워드의 쇼핑 클릭 추이를 조회하는 API이므로,
    현재 월/시즌에 맞는 후보군을 코드 내부에서 자동 선정해 상승폭을 비교한다.
    """
    common = [
        ("생필품", "50000008", ["휴지", "세제", "샴푸", "치약", "생수"]),
        ("건강식품", "50000006", ["유산균", "비타민", "오메가3", "단백질", "건강기능식품"]),
        ("디지털가전", "50000003", ["무선이어폰", "면도기", "소형가전", "마사지기", "로봇청소기"]),
    ]
    seasonal = {
        1: [("겨울", "50000003", ["온풍기", "전기요", "가습기", "히터", "온열매트"])],
        2: [("신학기", "50000003", ["태블릿", "노트북", "이어폰", "백팩", "책상"])],
        3: [("봄", "50000000", ["봄자켓", "트렌치코트", "운동화", "청소기", "공기청정기"])],
        4: [("나들이", "50000000", ["선크림", "운동화", "바람막이", "캠핑용품", "도시락"])],
        5: [("가정의달", "50000006", ["건강식품", "안마기", "화장품세트", "소형가전", "선물세트"])],
        6: [("초여름", "50000003", ["선풍기", "서큘레이터", "제습기", "선크림", "냉감티셔츠"])],
        7: [
            ("휴가", "50000005", ["캐리어", "여행가방", "여행용파우치", "기내용캐리어", "보스턴백"]),
            ("여름뷰티", "50000002", ["선스틱", "선크림", "선쿠션", "쿨링화장품", "데오드란트"]),
            ("여름가전", "50000003", ["휴대용선풍기", "서큘레이터", "제습기", "냉풍기", "미니선풍기"]),
            ("여름패션", "50000000", ["냉감티셔츠", "기능성티셔츠", "래쉬가드", "샌들", "아쿠아슈즈"]),
        ],
        8: [
            ("휴가", "50000005", ["캐리어", "여행가방", "여행용파우치", "기내용캐리어", "보스턴백"]),
            ("여름뷰티", "50000002", ["선스틱", "선크림", "선쿠션", "데오드란트", "쿨링화장품"]),
            ("여름가전", "50000003", ["휴대용선풍기", "서큘레이터", "제습기", "냉풍기", "미니선풍기"]),
        ],
        9: [("추석", "50000006", ["선물세트", "한우", "건강식품", "과일세트", "홍삼"])],
        10: [("가을", "50000000", ["가디건", "경량패딩", "등산복", "온열매트", "가습기"])],
        11: [("겨울준비", "50000003", ["온풍기", "전기요", "가습기", "패딩", "온열매트"])],
        12: [("연말", "50000006", ["선물세트", "홈파티", "와인잔", "향수", "소형가전"])],
    }
    return seasonal.get(month, []) + common


def _naver_shopping_keyword_scores(as_of_date, timeout=8):
    """
    NAVER API HUB Shopping Insight의 키워드별 트렌드 조회.
    최근 7일 평균 관심도 vs 직전 7일 평균 관심도를 비교해 상승률 계산.
    인증정보/통신/API 응답 문제 시 [] 반환 → 기존 대시보드에는 영향 없음.
    """
    cid, secret = _naver_trend_credentials()
    if not cid or not secret:
        return []

    end = pd.Timestamp(as_of_date).normalize()
    start = end - pd.Timedelta(days=13)

    url = "https://naverapihub.apigw.ntruss.com/shopping/v1/category/keywords"
    headers = {
        "X-NCP-APIGW-API-KEY-ID": cid,
        "X-NCP-APIGW-API-KEY": secret,
        "Content-Type": "application/json",
    }

    output = []
    for theme, category_id, keywords in _naver_trend_seed_catalog(int(end.month)):
        # API는 한 요청에 복수 keyword group을 받을 수 있으므로 최대 5개씩 묶음.
        keyword_groups = [
            {"name": kw, "param": [kw]}
            for kw in keywords[:5]
        ]
        payload = {
            "startDate": start.strftime("%Y-%m-%d"),
            "endDate": end.strftime("%Y-%m-%d"),
            "timeUnit": "date",
            "category": str(category_id),
            "keyword": keyword_groups,
        }
        try:
            resp = requests.post(url, headers=headers, json=payload, timeout=timeout)
            if resp.status_code != 200:
                continue
            data = resp.json()
        except Exception:
            continue

        for result in data.get("results", []):
            title = str(result.get("title", "")).strip()
            points = result.get("data", []) or []
            if not title or len(points) < 8:
                continue
            vals = []
            for p in points:
                try:
                    vals.append((pd.Timestamp(p.get("period")), float(p.get("ratio", 0))))
                except Exception:
                    pass
            vals = sorted(vals, key=lambda x: x[0])
            if len(vals) < 8:
                continue
            recent = [v for _, v in vals[-7:]]
            previous = [v for _, v in vals[-14:-7]] if len(vals) >= 14 else [v for _, v in vals[:-7]]
            if not previous:
                continue
            recent_avg = sum(recent) / len(recent)
            prev_avg = sum(previous) / len(previous)
            growth = ((recent_avg - prev_avg) / prev_avg * 100) if prev_avg > 0 else None
            output.append({
                "theme": theme,
                "keyword": title,
                "recent_avg": recent_avg,
                "prev_avg": prev_avg,
                "growth": growth,
            })

    # 최근 관심도와 상승률을 함께 반영하되, 전주 평균 0은 강한 트렌드로 단정하지 않음.
    output = [x for x in output if x["growth"] is not None]
    return sorted(
        output,
        key=lambda x: (x["growth"], x["recent_avg"]),
        reverse=True,
    )


def _match_trend_to_mms_history(keyword: str, products_all: pd.DataFrame, week_end):
    """
    외부 트렌드 키워드를 내부 MMS 과거 이력과 교차검증.
    키워드 직접 포함 + 일부 대표 동의어 매핑만 사용하며, 근거가 없으면 None.
    """
    aliases = {
        "기내용캐리어": ["캐리어"],
        "여행가방": ["캐리어", "여행가방"],
        "휴대용선풍기": ["선풍기", "써큘레이터", "서큘레이터"],
        "미니선풍기": ["선풍기"],
        "냉감티셔츠": ["냉감", "드라이셀", "기능성", "언더셔츠"],
        "기능성티셔츠": ["기능성", "드라이셀", "언더셔츠"],
        "선스틱": ["선스틱"],
        "선크림": ["선크림"],
        "선쿠션": ["선쿠션"],
        "무선이어폰": ["이어폰", "블루투스"],
        "면도기": ["면도기"],
        "유산균": ["유산균", "락토핏", "BNR17"],
    }
    terms = aliases.get(keyword, [keyword])
    pattern = "|".join(re.escape(x) for x in terms if x)
    if not pattern:
        return None

    df = products_all.copy()
    df["_date2"] = pd.to_datetime(df["_date"], errors="coerce")
    hist = df[
        df["상품명"].fillna("").astype(str).str.contains(pattern, case=False, regex=True)
        & df["_date2"].notna()
        & (df["_date2"] <= week_end)
    ].copy()
    if hist.empty:
        return None

    hist["_amt"] = pd.to_numeric(hist["주문금액"], errors="coerce").fillna(0)
    # 최소 한 번 300만원 이상이어야 'MMS 검증 근거'로 인정.
    if int((hist["_amt"] >= 3_000_000).sum()) < 1:
        return None

    price_col = first_col(hist, ["멤버십 혜택가", "행사가", "판매가", "혜택가"])
    if price_col:
        hist["_price"] = pd.to_numeric(hist[price_col], errors="coerce")
    else:
        hist["_price"] = pd.NA

    # 가장 성과가 좋은 실제 상품
    prod = (
        hist.groupby("상품명", as_index=False)
        .agg(
            운영횟수=("상품명", "size"),
            평균매출=("_amt", "mean"),
            최고매출=("_amt", "max"),
            고성과횟수=("_amt", lambda x: int((x >= 5_000_000).sum())),
            평균혜택가=("_price", "mean"),
        )
        .sort_values(["고성과횟수", "평균매출", "최고매출"], ascending=False)
    )
    if prod.empty:
        return None
    best = prod.iloc[0]
    pname = str(best["상품명"])
    ph = hist[hist["상품명"].astype(str) == pname].copy()

    # 해당 상품의 고성과 타겟
    target_cols = [c for c in ["성별", "연령", "SEG"] if c in ph.columns]
    target = ""
    target_amt = None
    if target_cols:
        tg = (
            ph.groupby(target_cols, dropna=False)["_amt"]
            .mean()
            .reset_index()
            .sort_values("_amt", ascending=False)
        )
        if not tg.empty:
            r = tg.iloc[0]
            bits = []
            if "성별" in target_cols:
                bits.append(str(r["성별"]))
            if "연령" in target_cols:
                bits.append(clean_identifier_value(r["연령"]))
            if "SEG" in target_cols:
                bits.append(str(r["SEG"]))
            target = " ".join(bits)
            target_amt = float(r["_amt"])

    hp = ph[ph["_amt"] >= 5_000_000]
    hp_price = (
        float(pd.to_numeric(hp["_price"], errors="coerce").mean())
        if not hp.empty and pd.to_numeric(hp["_price"], errors="coerce").notna().any()
        else None
    )

    return {
        "product": pname,
        "count": int(best["운영횟수"]),
        "avg_amt": float(best["평균매출"]),
        "max_amt": float(best["최고매출"]),
        "ge5": int(best["고성과횟수"]),
        "avg_price": float(best["평균혜택가"]) if pd.notna(best["평균혜택가"]) else None,
        "hp_price": hp_price,
        "target": target,
        "target_avg": target_amt,
    }


def _latest_trend_action_sentence(products_all: pd.DataFrame, week_end):
    """
    외부 최신 트렌드 + 내부 MMS 검증이 모두 있을 때만 차주 운영 제안 1건 생성.
    외부 API 실패/인증키 없음/내부 근거 없음 → None.
    """
    scores = _naver_shopping_keyword_scores(week_end)
    if not scores:
        return None

    # 상승률 20% 이상을 우선. 미달이면 '트렌드'로 강하게 표현하지 않음.
    for tr in scores[:12]:
        if tr["growth"] < 20:
            continue
        evidence = _match_trend_to_mms_history(tr["keyword"], products_all, week_end)
        if not evidence:
            continue

        price_bits = []
        if evidence["avg_price"] is not None:
            price_bits.append(f"과거 평균 혜택가 {evidence['avg_price']:,.0f}원")
        if evidence["hp_price"] is not None:
            price_bits.append(f"500만원 이상 고성과 회차 평균 {evidence['hp_price']:,.0f}원")
        price_text = ", ".join(price_bits)

        target_text = ""
        if evidence["target"]:
            target_text = f", 고성과 타겟 {evidence['target']}"
            if evidence["target_avg"] is not None:
                target_text += f" 평균 {compact_money(evidence['target_avg'])}"

        return (
            f"• NAVER 쇼핑 클릭 트렌드에서 최근 7일 '{tr['keyword']}' 관심도가 직전 7일 대비 "
            f"{tr['growth']:+.1f}% 상승했습니다. 내부 MMS 이력상 {_short_weekly_product_name(evidence['product'])}은 "
            f"{evidence['count']}회 운영, 평균 {compact_money(evidence['avg_amt'])}, 최고 {compact_money(evidence['max_amt'])}"
            + (f", 500만원 이상 {evidence['ge5']}회" if evidence["ge5"] else "")
            + target_text
            + (f", {price_text}" if price_text else "")
            + "의 성과가 확인됐습니다. 최신 관심 상승과 과거 MMS 성과가 동시에 확인된 만큼 "
              "유사 가격대·구성의 동일/유사 상품을 고성과 타겟 중심으로 신규·유사신규 TEST하는 것이 적절합니다."
        )
    return None



def _has_final_consonant(text_value: str) -> bool:
    s = str(text_value or "").strip()
    if not s:
        return False
    ch = s[-1]
    code = ord(ch)
    if 0xAC00 <= code <= 0xD7A3:
        return (code - 0xAC00) % 28 != 0
    return False


def _topic_particle(text_value: str) -> str:
    return "은" if _has_final_consonant(text_value) else "는"


def _subject_particle(text_value: str) -> str:
    return "이" if _has_final_consonant(text_value) else "가"


def _with_topic(text_value: str) -> str:
    s = str(text_value or "").strip()
    if not s:
        return s
    # 영문/숫자/모델명으로 끝나는 상품명은 조사를 직접 붙이지 않고 '상품은'으로 안전하게 처리
    last = s[-1]
    if not ("가" <= last <= "힣"):
        return f"{s} 상품은"
    return f"{s}{_topic_particle(s)}"


def _safe_product_label(name: str) -> str:
    """주간 화면용 축약명 + 조사 붙이기용 정리."""
    return _short_weekly_product_name(name).strip()


def _clean_text_value(value) -> str:
    """nan/None/NaT/빈 문자열을 화면에 노출하지 않음."""
    if value is None:
        return ""
    try:
        if pd.isna(value):
            return ""
    except Exception:
        pass
    s = str(value).strip()
    if s.lower() in {"nan", "none", "nat", "<na>"}:
        return ""
    return s


def _season_keyword_match_mask(series: pd.Series, labels: pd.Series | None = None):
    """
    상품명뿐 아니라 대카/중카까지 함께 보며 시즌 상품 탐지.
    상품명 키워드 누락을 보완하기 위한 공통 매칭용.
    """
    s = series.fillna("").astype(str)
    if labels is not None:
        s = s + " " + labels.fillna("").astype(str)
    patterns = {
        "캐리어·여행용품": r"캐리어|여행|파우치|보스턴백|트래블|기내용",
        "여행용 소형가전": r"여행용|휴대용|미니|면도기|드라이기|선풍기",
        "선케어": r"선크림|선스틱|선쿠션|자외선|선케어",
        "냉감·기능성 의류": r"냉감|쿨링|드라이셀|기능성|언더셔츠|에어리즘",
        "냉방가전": r"선풍기|써큘|서큘|냉풍|에어컨",
        "보양식·간편식": r"삼계탕|장어|갈비탕|국밥|간편식|즉석|냉동|보양식",
    }
    return patterns


def _repeat_operation_sentence(product_name: str, pw: pd.DataFrame):
    """금주 회차별 실적을 실제로 판정해 상품 운영 시사점 문장 생성."""
    sub = pw[pw["상품명"].astype(str) == str(product_name)].copy()
    if sub.empty or len(sub) < 2:
        return None
    sub["_date2"] = pd.to_datetime(sub["_date"], errors="coerce")
    if "시간대" in sub.columns:
        sub["_time_sort"] = sub["시간대"].astype(str)
    else:
        sub["_time_sort"] = ""
    sub = sub.sort_values(["_date2", "_time_sort"])
    vals = pd.to_numeric(sub["주문금액"], errors="coerce").fillna(0).tolist()
    if len(vals) < 2:
        return None

    seq_txt = " → ".join(compact_money(v) for v in vals)
    short = _safe_product_label(product_name)

    # 최근 3회 이상 연속 하락
    if len(vals) >= 3:
        recent3 = vals[-3:]
        decreasing = all(recent3[i] < recent3[i-1] for i in range(1, len(recent3)))
        drop = (recent3[0] - recent3[-1]) / recent3[0] if recent3[0] > 0 else 0
        if decreasing and drop >= 0.30:
            return (
                f"• {_with_topic(short)} 금주 {len(vals)}회 편성에서 회차별 주문금액이 {seq_txt}으로, "
                f"최근 3회 연속 하락하며 최초 대비 {drop*100:.0f}% 감소했습니다. "
                f"동일 상품의 반복 운영 성과 둔화가 확인돼 즉시 동일 조건 재편성보다 "
                f"최근 미발송 타겟·SEG 전환 TEST 또는 일정 기간 미편성 후 재운영하는 것이 적절합니다."
            )

    # 모두 300만원 이상이면 안정 유지
    if all(v >= 3_000_000 for v in vals):
        return (
            f"• {_with_topic(short)} 금주 {len(vals)}회 편성 모두 300만원 이상을 기록했고 회차별 주문금액은 "
            f"{seq_txt}으로 연속 하락은 확인되지 않았습니다. 반복 운영에도 성과가 유지돼 "
            f"고성과 타겟 중심의 단기 재편성이 가능하며, 이후 2회 이상 연속 하락이 발생할 때 "
            f"최근 ○일간 미편성 후 재운영하는 기준을 적용하는 것이 적절합니다."
        )

    # 등락 반복
    return (
        f"• {_with_topic(short)} 금주 {len(vals)}회 편성의 회차별 주문금액은 {seq_txt}으로 편차가 확인됐습니다. "
        f"단순 반복 횟수보다 각 회차의 성별·연령·SEG·가격 조건을 함께 비교해 고성과 조건을 선별한 뒤 재편성하는 것이 적절합니다."
    )



def _compact_weekly_business_tone(sentence: str) -> str:
    """주간실적 시사점을 간결한 실무 보고체로 정리."""
    s = str(sentence or "").strip()
    replacements = [
        ("기록했습니다.", "기록"),
        ("확인됐습니다.", "확인"),
        ("확인되었습니다.", "확인"),
        ("확인됩니다.", "확인"),
        ("차지했습니다.", "차지"),
        ("높았습니다.", "우수"),
        ("낮았습니다.", "열위"),
        ("유지됐습니다.", "유지"),
        ("유지되었습니다.", "유지"),
        ("개선됐습니다.", "개선"),
        ("개선되었습니다.", "개선"),
        ("필요가 있습니다.", "필요"),
        ("필요합니다.", "필요"),
        ("적절합니다.", "검토 필요"),
        ("검토하는 것이 좋습니다.", "검토 필요"),
        ("검토하는 것이 적절합니다.", "검토 필요"),
        ("활용할 수 있습니다.", "활용 가능"),
        ("병행할 수 있습니다.", "병행 가능"),
        ("재편성이 가능합니다.", "재편성 가능"),
        ("검토할 수 있습니다.", "검토 가능"),
        ("것이 좋습니다.", "검토 필요"),
        ("것이 필요합니다.", "필요"),
    ]
    for old, new in replacements:
        s = s.replace(old, new)
    s = re.sub(r"(했습니다|됩니다|있습니다)\.$", "", s).rstrip(".").strip()
    return s


def _weekly_new_repeat_stats(pw: pd.DataFrame, products_all: pd.DataFrame, week_start) -> dict | None:
    """
    금주 상품을 과거 동일상품 이력 유무 기준으로
    신규·유사신규 후보군 vs 재편성으로 구분해 성과 비교.
    """
    if pw.empty or "상품명" not in pw.columns:
        return None

    wk = pw.groupby("상품명", as_index=False)["주문금액"].sum().copy()
    if wk.empty:
        return None

    hist = products_all.copy()
    hist["_date2"] = pd.to_datetime(hist["_date"], errors="coerce")
    week_start = pd.to_datetime(week_start, errors="coerce")

    rows = []
    for _, r in wk.iterrows():
        pname = str(r["상품명"])
        prior = hist[
            (hist["상품명"].astype(str) == pname)
            & hist["_date2"].notna()
            & (hist["_date2"] < week_start)
        ]
        kind = "재편성" if not prior.empty else "신규·유사신규"
        rows.append({"구분": kind, "주문금액": float(r["주문금액"])})

    d = pd.DataFrame(rows)
    out = {}
    for kind, g in d.groupby("구분"):
        out[kind] = {
            "상품수": int(len(g)),
            "평균매출": float(g["주문금액"].mean()) if len(g) else 0,
            "삼백이상률": float((g["주문금액"] >= 3_000_000).mean()) if len(g) else 0,
            "오백이상률": float((g["주문금액"] >= 5_000_000).mean()) if len(g) else 0,
        }
    return out or None


def _weekly_success_formula_sentence(
    core_count: int,
    core_share: float,
    amount_delta: float | None,
    spm_delta: float | None,
) -> str | None:
    """금주 성과 구조를 과도한 인과 단정 없이 요약."""
    if core_count <= 0 or core_share <= 0:
        return None
    if amount_delta is not None and spm_delta is not None and amount_delta > 0 and spm_delta > 0:
        return (
            f"• 금주 성과 핵심 요인 : 500만원 이상 핵심 상품 {core_count}개가 매출의 {core_share:.1f}%를 차지한 가운데 "
            f"주문금액 {amount_delta*100:+.1f}%, SPM {spm_delta*100:+.1f}% 개선 > "
            f"단순 발송량 확대보다 검증 상품 × 적합 타겟 중심의 편성 효율 개선이 성과 상승과 함께 확인"
        )
    return (
        f"• 금주 성과 구조 : 500만원 이상 핵심 상품 {core_count}개가 매출의 {core_share:.1f}% 차지 > "
        f"핵심 상품 의존도와 신규 매출원 확보 수준을 함께 관리할 필요"
    )



def _weekly_apply_selective_price_condition(sentence: str) -> str:
    """
    차주 운영 제안에 가격·구성 성공 조건을 선택적으로 반영.
    실제 과거 혜택가 근거가 문장에 있을 때만 가격대를 범위화하며,
    근거가 없으면 임의 생성하지 않음.
    """
    s = str(sentence or "")

    if any(k in s for k in ["에어써큘", "서큘", "선풍기"]):
        m = re.search(r"당시 혜택가\s*([0-9,]+)원", s)
        if m:
            price = int(m.group(1).replace(",", ""))
            if 30000 <= price < 40000:
                s = re.sub(r"당시 혜택가\s*[0-9,]+원\.?\s*", "", s)
                s = s.replace("당시와 유사한 가격 조건을 확보한", "3만원대 가격과")

    if any(k in s for k in ["우양산", "양산"]):
        m = re.search(r"당시 혜택가\s*([0-9,]+)원", s)
        if m:
            price = int(m.group(1).replace(",", ""))
            if 8000 <= price < 15000:
                s = re.sub(r"당시 혜택가\s*[0-9,]+원\.?\s*", "", s)
                s = s.replace("당시와 유사한 가격 조건을 확보한", "1만원 내외 가격대의")
                s = s.replace("1만원 내외 가격대의 1만원 내외 가격대의", "1만원 내외 가격대의")

    return s

def build_weekly_analysis(week, year, pw, sw, products_all, sends_all) -> str:
    send_col = first_col(sw, ["발송 성공 건수", "총 발송 건수"])
    week_start = pd.to_datetime(pw["_date"], errors="coerce").min() if not pw.empty else pd.NaT
    amount_delta = None
    spm_delta = None
    click_col = first_col(sw, ["클릭 수(uniq)", "클릭 수"])

    send_count = float(sw[send_col].sum())
    click_count = float(sw[click_col].sum())
    order_count = float(sw["주문건수"].sum())
    qty = float(sw["주문수량"].sum())
    amount = float(sw["주문금액"].sum())
    ctr = click_count / send_count if send_count else 0
    cvr = order_count / click_count if click_count else 0
    aov = amount / order_count if order_count else 0
    spm = amount / send_count if send_count else 0

    # 전주 비교
    all_weeks = sends_all[sends_all["_year"] == year].groupby("주차")["_date"].min().sort_values()
    week_names = [str(x) for x in all_weeks.index]
    prev_sw = pd.DataFrame()
    prev_pw = pd.DataFrame()
    if week in week_names and week_names.index(week) > 0:
        prev_week = week_names[week_names.index(week)-1]
        prev_sw = sends_all[(sends_all["_year"] == year) & (sends_all["주차"].astype(str) == prev_week)].copy()
        prev_pw = products_all[(products_all["_year"] == year) & (products_all["주차"].astype(str) == prev_week)].copy()

    if not prev_sw.empty:
        psend = float(prev_sw[send_col].sum())
        pclick = float(prev_sw[click_col].sum())
        porders = float(prev_sw["주문건수"].sum())
        pqty = float(prev_sw["주문수량"].sum())
        pamount = float(prev_sw["주문금액"].sum())
        pctr = pclick / psend if psend else 0
        pcvr = porders / pclick if pclick else 0
        paov = pamount / porders if porders else 0
        pspm = pamount / psend if psend else 0
        # 주간 요약은 KPI 3행(•) + 핵심 해석 2행(:)으로 고정
        order_delta = (order_count - porders) / abs(porders) if porders else 0
        amount_delta = (amount - pamount) / abs(pamount) if pamount else 0
        spm_delta = (spm - pspm) / abs(pspm) if pspm else 0
        send_delta = (send_count - psend) / abs(psend) if psend else 0

        summary = [
            f"• 발송횟수 {len(sw):,}회({_weekly_plain_delta(len(sw),len(prev_sw))}) / 상품수 {len(pw):,}건({_weekly_plain_delta(len(pw),len(prev_pw))}) / 발송건수 {int(send_count):,}건({_weekly_plain_delta(send_count,psend)}) 운영",
            f"• 주문건수 {int(order_count):,}건({_weekly_plain_delta(order_count,porders)}) / 주문수량 {int(qty):,}건({_weekly_plain_delta(qty,pqty)}) / 주문금액 {compact_money(amount)}({_weekly_plain_delta(amount,pamount)}) 기록",
            f"• CTR {ctr*100:.1f}%({_weekly_plain_delta(ctr,pctr,True)}) / CVR {cvr*100:.1f}%({_weekly_plain_delta(cvr,pcvr,True)}) / 객단가 {int(aov):,}원({_weekly_plain_delta(aov,paov)}) / SPM {spm:.1f}({_weekly_plain_delta(spm,pspm)}) 기록",
        ]

        # 규모 축소/확대와 성과 변화를 함께 해석해 단순 수치 반복을 피함
        if (len(sw) < len(prev_sw) or len(pw) < len(prev_pw) or send_count < psend) and amount > pamount and ctr >= pctr and cvr >= pcvr and spm > pspm:
            summary.append(": 발송횟수·상품수 감소에도 주문건수·주문금액 및 CTR·CVR·SPM이 모두 개선되며 전주 대비 높은 발송 효율 기록")
        elif amount > pamount and spm > pspm:
            summary.append(": 주문금액과 SPM이 함께 개선되며 전주 대비 매출 및 발송 효율 상승")
        elif amount < pamount and spm < pspm:
            summary.append(": 주문금액과 SPM이 함께 감소해 전주 대비 매출 및 발송 효율 둔화")
        else:
            summary.append(": 발송 규모와 주요 성과지표의 증감이 혼재해 상품·타겟별 기여도 추가 확인 필요")

        if order_delta > 0 and spm_delta > 0:
            summary.append(f": 특히 주문건수 {order_delta*100:+.1f}%, SPM {spm_delta*100:+.1f}%로 발송 규모 {'축소' if send_delta < 0 else '변화'} 대비 구매전환 및 매출 효율 크게 개선")
        elif amount_delta > 0:
            summary.append(f": 주문금액 {amount_delta*100:+.1f}% 증가를 기록해 매출 성장 기여 상품과 타겟 중심의 재현 조건 확인 필요")
        elif amount_delta < 0:
            summary.append(f": 주문금액 {amount_delta*100:+.1f}% 감소해 저성과 상품·타겟 및 반복 편성 영향 점검 필요")
    else:
        summary = [
            f"• 발송횟수 {len(sw):,}회 / 상품수 {len(pw):,}건 / 발송건수 {int(send_count):,}건 운영",
            f"• 주문건수 {int(order_count):,}건 / 주문수량 {int(qty):,}건 / 주문금액 {compact_money(amount)} 기록",
            f"• CTR {ctr*100:.1f}% / CVR {cvr*100:.1f}% / 객단가 {int(aov):,}원 / SPM {spm:.1f} 기록",
            ": 전주 비교 데이터가 없어 금주 실적을 기준으로 상품·타겟·편성 효율 확인",
        ]

    week_end = pd.to_datetime(pw["_date"], errors="coerce").max()

    # 상품별 주간 집계
    rank = pw.groupby("상품명",as_index=False).agg(
        주문금액=("주문금액","sum"),
        운영횟수=("상품명","size")
    ).sort_values("주문금액",ascending=False)
    core = rank[rank["주문금액"]>=5_000_000]
    poor = rank[rank["주문금액"]<1_000_000]

    product_points = []

    # 핵심상품 집중도: 실제 비중을 계산한 경우에만
    if not core.empty and amount > 0:
        core_sum = float(core["주문금액"].sum())
        core_share = core_sum / amount * 100
        core_names = ", ".join(_short_weekly_product_name(x) for x in core.head(4)["상품명"].astype(str))
        if core_share >= 50:
            product_points.append(
                f"• 핵심 상품 매출 집중 : 500만원 이상 핵심 상품 {len(core)}개가 전체 주문금액의 {core_share:.1f}% 차지하며 금주 매출 성장 견인 > "
                f"{core_names} 등 검증 상품은 안정적으로 재편성하되 신규·유사신규 고성과 후보를 병행 발굴해 핵심 상품군 확대 필요"
            )
        else:
            product_points.append(
                f"• 핵심 상품 매출 기여 : 500만원 이상 핵심 상품 {len(core)}개가 주문금액의 {core_share:.1f}% 차지 > "
                f"{core_names} 등 검증 상품은 유지하되 특정 상품 의존도와 신규 매출원 확보 수준 지속 점검 필요"
            )

    # 저성과 상품 비중: 편성 비효율을 핵심상품 집중도와 함께 확인
    if not rank.empty:
        poor_count = int(len(poor))
        poor_share = poor_count / len(rank) * 100 if len(rank) else 0
        if poor_count > 0:
            product_points.append(
                f"• 저성과 상품 효율 점검 : 금주 고유 상품 {len(rank)}개 중 100만원 미만 {poor_count}개로 {poor_share:.1f}% 차지 > "
                f"반복 저성과 상품은 재편성 우선순위를 낮추고 과거 300만원 이상 검증 상품 또는 신규·유사신규 후보로 교체 필요"
            )

    # 최고매출 상품: 실제 반복횟수/회당 성과 근거 반영
    if not rank.empty:
        r = rank.iloc[0]
        pname = str(r["상품명"])
        wh = pw[pw["상품명"].astype(str) == pname].copy()
        wk_amounts = pd.to_numeric(wh["주문금액"], errors="coerce").fillna(0)
        ge3 = int((wk_amounts >= 3_000_000).sum())
        count = len(wh)
        short = _short_weekly_product_name(pname)
        if count >= 2:
            product_points.append(
                f"• {_with_topic(short)} 금주 {count}회 편성 중 {ge3}회 300만원 이상을 기록하고 누적 {compact_money(float(r['주문금액']))}으로 최고 매출을 기록했습니다. 회차별 주문금액과 타겟별 성과를 실제 비교해 반복 운영 지속 여부를 판정하는 것이 적절합니다."
            )
        else:
            product_points.append(
                f"• {_with_topic(short)} 금주 {compact_money(float(r['주문금액']))}으로 최고 매출을 기록했습니다. 동일 타겟 1회 추가 검증 후 유사 성과가 유지되면 운영 확대를 검토할 수 있습니다."
            )

    # 식품/건강 카테고리는 전체가 잘돼도 개별상품 부진 가능성 분리
    food_rows = pw[pw["대카"].astype(str).str.contains("식품|건강", na=False)].copy()
    if not food_rows.empty and amount > 0:
        food_amount = float(food_rows["주문금액"].sum())
        food_share = food_amount / amount * 100
        under = _weekly_food_underperformers(pw)
        if under:
            product_points.append(
                f"• 식품/건강은 전체 주문금액의 {food_share:.1f}%로 높은 비중을 차지했으나 {', '.join(under)} 등은 100만원 미만에 그쳤습니다. 카테고리 자체보다 상품 대중성·구성·가격 경쟁력에 따른 편차가 큰 만큼 식품 비중을 줄이기보다 과거 MMS 고성과 이력이 있는 검증 상품 중심으로 교체 편성이 필요합니다."
            )

    # 부진 상품 중 가격 근거를 제시할 수 있는 대표 1개
    if not poor.empty:
        price_col = first_col(pw, ["멤버십 혜택가", "행사가", "판매가", "혜택가"])
        if price_col:
            for _, pr in poor.head(6).iterrows():
                pname = str(pr["상품명"])
                sub = pw[pw["상품명"].astype(str) == pname].copy()
                if sub.empty:
                    continue
                sale = pd.to_numeric(sub[price_col], errors="coerce").dropna()
                if sale.empty:
                    continue
                unit_phrase = _unit_price_phrase(pname, float(sale.iloc[-1]))
                hist = _weekly_product_history_stats(pname, products_all, week_end)
                if unit_phrase and hist and hist["count"] >= 2:
                    price_text = ""
                    if hist["sale_min"] is not None and hist["sale_max"] is not None:
                        if abs(hist["sale_max"] - hist["sale_min"]) / max(hist["sale_min"],1) <= 0.1:
                            price_text = f" 과거에도 유사한 가격대({hist['sale_min']:,.0f}~{hist['sale_max']:,.0f}원)로 운영됐습니다."
                    hp = _latest_and_high_perf_price(pname, products_all)
                    hp_text = ""
                    if hp and hp["high_perf_avg_price"]:
                        hp_text = f" 과거 500만원 이상 고성과 운영 당시 평균 혜택가는 {hp['high_perf_avg_price']:,.0f}원이었습니다."
                    product_points.append(
                        f"• {_with_topic(_short_weekly_product_name(pname))} 금번 {unit_phrase} 수준임에도 100만원 미만을 기록했습니다.{price_text}{hp_text} 과거 유사 가격 조건에서도 반복적으로 100만원 미만이 확인된 경우에 한해 가격보다 MMS 메인 상품 적합도 이슈로 판단하고 편성 우선순위를 조정하는 것이 적절합니다."
                    )
                    break

    # 성별·연령·SEG × 반복 운영 자동 판정
    # 프로모션은 이 판정에서 제외하고 실제 운영횟수·매출·SPM·고성과율만 사용
    if not rank.empty:
        target_sentence_added = 0
        for pname in rank.head(10)["상품명"].astype(str):
            ta = _product_target_strength_analysis(pname, products_all, week_end)
            sentence = _target_strength_sentence(pname, ta)
            if sentence:
                product_points.append(sentence)
                target_sentence_added += 1
            if target_sentence_added >= 2:
                break

    # 프로모션 효과 분리: 일반기간에서도 성과가 유지되는지 확인
    if not rank.empty:
        for pname in rank.head(8)["상품명"].astype(str):
            ps = _promotion_performance_stats(pname, products_all)
            if not ps or ps["promo_n"] < 1 or ps["normal_n"] < 1:
                continue
            pa, na = ps["promo_avg"], ps["normal_avg"]
            if na > 0 and pa / na >= 1.8:
                product_points.append(
                    f"• {_with_topic(_short_weekly_product_name(pname))} 프로모션 기간 평균 {compact_money(pa)} 대비 일반기간 평균 {compact_money(na)}으로 성과 차이가 커 프로모션 의존도가 높은 상품입니다. 일반기간에는 동일 수준의 재편성을 지양하고 프로모션 연계 운영을 우선 검토하는 것이 적절합니다."
                )
                break
            elif na >= 3_000_000:
                product_points.append(
                    f"• {_with_topic(_short_weekly_product_name(pname))} 프로모션 기간 평균 {compact_money(pa)}, 일반기간 평균 {compact_money(na)}으로 일반 운영에서도 안정적인 성과가 확인됩니다. 프로모션 여부와 무관하게 고성과 타겟 중심의 재편성 후보로 활용할 수 있습니다."
                )
                break

    # 편성 운영 시사점
    seg = grouped_send_table(sw,["성별","연령"])
    weekday = grouped_send_table(sw,["요일"])
    time_df = grouped_send_table(sw,["시간대"])
    big_cat = pw.groupby("대카",as_index=False)["주문금액"].sum().sort_values("주문금액",ascending=False)

    op = []
    if not seg.empty:
        s = seg.loc[seg["SPM"].idxmax()]
        g = str(s["성별"]); a = clean_identifier_value(s["연령"])
        top_products = _dominant_target_products(pw, g, a, 3)
        if top_products:
            op.append(
                f"• {g}{a}은 {', '.join(top_products)} 등 고성과 상품이 함께 편성된 가운데 SPM {s['SPM']:.1f}를 기록했습니다. 타겟 자체가 우수하다고 단정하기보다 해당 타겟에서 반복적으로 성과가 확인된 상품군과 SEG를 우선 확인하고, 동일 조건에서 재현 여부를 검증한 뒤 유사 상품 재편성과 미발송 SEG 확대 TEST에 활용하는 것이 좋습니다."
            )

    # 신규·유사신규 vs 재편성 성과 비교
    nr_stats = _weekly_new_repeat_stats(pw, products_all, week_start)
    if nr_stats and "재편성" in nr_stats and "신규·유사신규" in nr_stats:
        rep = nr_stats["재편성"]
        newc = nr_stats["신규·유사신규"]
        if rep["상품수"] > 0 and newc["상품수"] > 0:
            rep_avg = rep["평균매출"]
            new_avg = newc["평균매출"]
            rep_hit = rep["삼백이상률"]
            new_hit = newc["삼백이상률"]

            if rep_avg > new_avg and new_hit > rep_hit:
                verdict = "재편성은 평균매출, 신규·유사신규는 300만원 이상 성공률에서 각각 강점 확인"
                action = "검증 상품 재편성과 신규 후보 TEST 병행 필요"
            elif rep_avg >= new_avg * 1.25 and rep_hit >= new_hit:
                verdict = "재편성이 평균매출과 300만원 이상 성공률 모두 우위"
                action = "핵심 재편성 중심 운영하되 신규·유사신규 후보의 가격·구성 경쟁력 강화 필요"
            elif new_avg >= rep_avg * 1.25 and new_hit >= rep_hit:
                verdict = "신규·유사신규가 평균매출과 300만원 이상 성공률 모두 우위"
                action = "신규 소싱을 유지하고 고성과 신규 상품은 빠르게 재편성 후보로 전환 필요"
            else:
                verdict = "유형별 성과 우위가 혼재"
                action = "평균매출과 300만원 이상 성공률을 함께 기준으로 재편성·신규 TEST 비중 조정 필요"

            op.append(
                f"• 신규·유사신규 vs 재편성 : 신규·유사신규 {newc['상품수']}개 평균 {compact_money(new_avg)}·300만원 이상 비중 {new_hit*100:.1f}%, "
                f"재편성 {rep['상품수']}개 평균 {compact_money(rep_avg)}·300만원 이상 비중 {rep_hit*100:.1f}% 기록 > {verdict}, {action}"
            )

    # 최근 4주 반복성 실제 계산: 3주 이상 동일 우위일 때만 강한 시사점 생성
    pattern4 = _recent_4week_time_pattern(week, year, sends_all)
    if pattern4 and pattern4.get("time"):
        tname, tcnt, ttotal = pattern4["time"]
        if tcnt >= 3:
            op.append(
                f"• 최근 {ttotal}주 중 {tcnt}주에서 {tname} 시간대가 SPM 최고를 기록해 시간대 우위가 반복 확인됐습니다. 해당 시간대에 편성된 고성과 상품·타겟 조합을 기준으로 핵심 상품 우선 배치 TEST를 확대하는 것이 적절합니다."
            )
    if pattern4 and pattern4.get("day"):
        dname, dcnt, dtotal = pattern4["day"]
        if dcnt >= 3:
            op.append(
                f"• 요일별 편성 조건 검증 : 최근 {dtotal}주 중 {dcnt}주에서 {dname}요일 SPM 최고 기록 > "
                f"요일 자체 효과로 단정하기보다 해당 요일의 핵심 상품 비중·타겟·SEG·발송모수를 함께 비교해 공통 고성과 조건 확인 필요"
            )

    if not big_cat.empty and big_cat["주문금액"].sum()>0:
        total = big_cat["주문금액"].sum()
        cats = ", ".join(f"{r['대카']} {r['주문금액']/total*100:.1f}%" for _,r in big_cat.head(3).iterrows())
        op.append(
            f"• 대카테고리 매출은 {cats} 순으로 구성됐습니다. 카테고리 비중만으로 편성 우선순위를 정하기보다 카테고리 내 과거 300만원·500만원 이상 달성 횟수와 가격 경쟁력을 함께 비교해 검증 상품 중심으로 편성을 정교화할 필요가 있습니다."
        )



    # 반복 운영 상품: generic 문장을 실제 회차별 판정 문장으로 교체
    _repeat_replaced = []
    _repeat_done = set()
    for _s in product_points:
        _matched = False
        for _pname in pd.Series(pw.get("상품명", [])).dropna().astype(str).unique() if "상품명" in pw.columns else []:
            _short = _safe_product_label(_pname)
            if _short and _short in str(_s):
                _rows = pw[pw["상품명"].astype(str) == str(_pname)] if "상품명" in pw.columns else pd.DataFrame()
                if len(_rows) >= 2 and ("회차별" in str(_s) or "편성 중" in str(_s)):
                    _actual = _repeat_operation_sentence(_pname, pw)
                    if _actual and _pname not in _repeat_done:
                        _repeat_replaced.append(_actual)
                        _repeat_done.add(_pname)
                        _matched = True
                        break
        if not _matched:
            _repeat_replaced.append(_s)
    product_points = _repeat_replaced

    # 상품 운영 시사점 중복 제거: 동일 문장/동일 반복판정 중복 방지
    _pp_seen = set()
    _pp_dedup = []
    for _s in product_points:
        _key = re.sub(r"\s+", " ", str(_s)).strip()
        if _key and _key not in _pp_seen:
            _pp_seen.add(_key)
            _pp_dedup.append(_s)
    product_points = _pp_dedup

    # 차주 운영 제안: 개수 제한 없이 실제 근거가 있는 제안만 우선순위 순으로 노출
    nxt = []
    ranked_actions = _next_week_action_candidates(pw, products_all, week_end)

    # 유형별 중복을 제한하되 전체 개수는 제한하지 않음.
    # 즉시 재편성은 상품별 최대 2건, 나머지는 유형별 1건 우선.
    used_kinds = {}
    seen_sentences = set()

    for score, kind, sentence in ranked_actions:
        limit = 2 if kind == "즉시 재편성" else 1
        if used_kinds.get(kind, 0) >= limit:
            continue
        clean_sentence = sentence.strip()
        if not clean_sentence or clean_sentence in seen_sentences:
            continue
        nxt.append("• " + clean_sentence)
        seen_sentences.add(clean_sentence)
        used_kinds[kind] = used_kinds.get(kind, 0) + 1

    # 시즌성 실제 상품 근거는 항상 별도 축으로 노출.
    seasonal_evidence = _seasonal_action_sentence(products_all, week_end)
    if seasonal_evidence:
        clean = seasonal_evidence.strip()
        if clean and clean not in seen_sentences:
            nxt.append(clean)
            seen_sentences.add(clean)
    else:
        # 실제 과거 시즌 고성과 상품 근거가 전혀 없을 때만 편성 횟수형 fallback 사용
        season_gap = _season_gap_action(pw, products_all, week_end)
        if season_gap:
            clean = ("• " + season_gap).strip()
            if clean not in seen_sentences:
                nxt.append(clean)
                seen_sentences.add(clean)

    # 최신 NAVER 외부 트렌드는 내부 MMS 근거까지 교차검증된 경우 별도 노출.
    # 다른 제안이 많아도 개수 제한 때문에 잘리지 않음.
    latest_trend_action = _latest_trend_action_sentence(products_all, week_end)
    if latest_trend_action:
        clean = latest_trend_action.strip()
        if clean and clean not in seen_sentences:
            nxt.append(clean)
            seen_sentences.add(clean)

    # 근거형 제안이 거의 없을 때만 일반 원칙 보완
    if len(nxt) < 3:
        poor_now = pw.groupby("상품명", as_index=False)["주문금액"].sum()
        poor_now = poor_now[poor_now["주문금액"] < 1_000_000]
        if not poor_now.empty:
            fallback_sentence = (
                "• 금주 100만원 미만 상품은 과거 운영 이력·개당/장당 실구매가·타겟별 성과에서 "
                "개선 근거가 확인되는 경우에만 재TEST하고, 동일 조건 반복 부진 상품은 검증된 대체 상품으로 "
                "교체하는 것이 적절합니다."
            )
            if fallback_sentence not in seen_sentences:
                nxt.append(fallback_sentence)
                seen_sentences.add(fallback_sentence)

    # 모든 주차 공통 최종 중복 제거:
    # 동일 product master가 여러 추천 규칙에 걸리면 근거를 하나로 병합해 1회만 노출
    nxt = _merge_same_product_recommendations(nxt, products_all)



    nxt = [_weekly_apply_selective_price_condition(x) for x in nxt]

    return "\n".join([
        "■ 주간 실적 요약",*summary,"",
        "■ 상품 운영 시사점",*(product_points[:5] or ["• 금주 상품 성과 기준 재편성 우선순위 점검 필요"]),"",
        "■ 편성 운영 시사점",*(op[:4] or ["• 타겟·요일·시간대·카테고리별 효율을 원인 상품과 함께 비교해 편성 우선순위 조정 필요"]),"",
        "■ 차주 운영 제안",*nxt
    ])


APP_DIR = Path(__file__).resolve().parent
IMAGE_DIR = APP_DIR / "images"
MESSAGE_DIR = APP_DIR / "messages"


def daily_asset_key(date_value, time_value) -> str:
    dt = pd.to_datetime(date_value, errors="coerce")
    if pd.isna(dt):
        return ""

    time_text = str(time_value).strip()
    if time_text.startswith("10") or "오전" in time_text or time_text in ["1", "01"]:
        slot = "01"
    elif time_text.startswith("16") or time_text.startswith("17") or "오후" in time_text or time_text in ["2", "02"]:
        slot = "02"
    else:
        parsed = pd.to_datetime(time_text, errors="coerce")
        if pd.notna(parsed):
            slot = "01" if parsed.hour < 12 else "02"
        else:
            return ""
    return f"{dt:%Y%m%d}_{slot}"


def find_daily_image(asset_key: str, campaign_name: str = ""):
    if not IMAGE_DIR.exists():
        return None

    valid_suffixes = [".jpg", ".jpeg", ".png", ".webp"]
    files = [p for p in IMAGE_DIR.iterdir() if p.is_file() and p.suffix.lower() in valid_suffixes]

    # 캠페인명과 동일한 파일명을 우선 사용
    campaign_name = str(campaign_name).strip()
    if campaign_name:
        exact = [p for p in files if p.stem == campaign_name]
        if exact:
            return sorted(exact, key=lambda p: p.name)[0]

    # 기존 날짜_01/02 방식도 계속 지원
    if asset_key:
        matches = [p for p in files if p.name.startswith(asset_key)]
        if matches:
            return sorted(matches, key=lambda p: p.name)[0]

    return None


def clean_mms_message(value) -> str:
    """앞뒤 큰따옴표만 제거하고 내부 줄바꿈은 그대로 유지합니다."""
    if value is None or pd.isna(value):
        return ""

    text_value = str(value)
    stripped = text_value.strip()

    if len(stripped) >= 2 and stripped.startswith('"') and stripped.endswith('"'):
        stripped = stripped[1:-1]

    return stripped.strip("\r\n")


def extract_mms_message(
    matched: pd.DataFrame,
    send_row: pd.Series,
    messages_df: pd.DataFrame | None = None,
) -> str:
    """문구 시트의 캠페인명을 우선 매칭하고 기존 RAW 컬럼은 보조로 사용합니다."""
    campaign_name = str(send_row.get("캠페인명", "")).strip()

    if messages_df is not None and not messages_df.empty and campaign_name:
        matched_message = messages_df[
            messages_df["캠페인명"].astype(str).str.strip().eq(campaign_name)
        ]
        if not matched_message.empty:
            cleaned = clean_mms_message(matched_message.iloc[-1]["MMS문구"])
            if cleaned:
                return cleaned

    candidate_cols = ["MMS문구", "MMS 문구", "발송문구", "문구"]

    for col in candidate_cols:
        if col in matched.columns:
            for value in matched[col].tolist():
                cleaned = clean_mms_message(value)
                if cleaned:
                    return cleaned

    for col in candidate_cols:
        if col in send_row.index:
            cleaned = clean_mms_message(send_row.get(col))
            if cleaned:
                return cleaned

    return ""


def format_discount_percent(x):
    if pd.isna(x) or str(x).strip() in ["", "nan", "None"]:
        return ""
    try:
        value = float(str(x).replace("%", "").replace(",", "").strip())
        if abs(value) <= 1:
            value *= 100
        return f"{value:.0f}%"
    except (TypeError, ValueError):
        return str(x)


def format_integer_price(x):
    if pd.isna(x) or str(x).strip() in ["", "nan", "None"]:
        return ""
    try:
        return f"{float(str(x).replace(',', '')):,.0f}"
    except (TypeError, ValueError):
        return str(x)



def floor_discount_rate(normal_price, sale_price):
    """엑셀 =ROUNDDOWN(1-행사가/정상가, 2)와 동일하게 계산합니다."""
    try:
        normal = float(str(normal_price).replace(",", "").strip())
        sale = float(str(sale_price).replace(",", "").strip())
        if normal <= 0:
            return pd.NA
        return math.floor((1 - sale / normal) * 100) / 100
    except (TypeError, ValueError):
        return pd.NA


def parse_slot_day(material_name: str) -> str:
    """7/21(1), 7/21(2)처럼 같은 날짜 슬롯을 묶기 위한 날짜 키입니다."""
    value = str(material_name).strip()
    match = re.search(r"(\d{1,2})\s*/\s*(\d{1,2})", value)
    if match:
        return f"{int(match.group(1)):02d}-{int(match.group(2)):02d}"
    return value.split("(")[0].strip()


def parse_target_text(target_text: str) -> dict:
    text_value = str(target_text).replace(" ", "")
    gender = "여성" if "여성" in text_value else ("남성" if "남성" in text_value else "")
    age = "5060" if "5060" in text_value else ("3040" if "3040" in text_value else "")
    seg_match = re.search(r"SEG\s*([123])", text_value, flags=re.I)
    seg = seg_match.group(1) if seg_match else ""
    return {"성별": gender, "연령": age, "SEG": seg}


def is_candidate_gender_compatible(candidate: pd.Series, target_text: str) -> bool:
    """상품명 기준 성별 전용 상품을 타겟 후보에서 제외합니다."""
    target = parse_target_text(target_text)
    target_gender = target.get("성별", "")
    name = str(candidate.get("상품명", "")).lower().replace(" ", "")

    male_only_keywords = [
        "면도기", "전기면도기", "남성그루밍", "그루밍풀세트", "코털제거기",
        "남성드로즈", "남성언더웨어", "남성용", "남자용", "남성팬티",
    ]
    female_only_keywords = [
        "여성언더웨어", "여성용", "여자용", "여성팬티", "여성브라", "브라팬티",
    ]

    if target_gender == "여성" and any(keyword in name for keyword in male_only_keywords):
        return False
    if target_gender == "남성" and any(keyword in name for keyword in female_only_keywords):
        return False
    return True


def match_candidate_history(candidate: pd.Series, history: pd.DataFrame) -> pd.DataFrame:
    """쇼라코드 → 알파코드 → 상품명 순으로 과거 이력을 찾습니다."""
    for key in ["쇼라코드", "알파코드"]:
        value = clean_identifier_value(candidate.get(key, ""))
        if value and key in history.columns:
            matched = history[history[key].map(clean_identifier_value).eq(value)]
            if not matched.empty:
                return matched.sort_values("_date")

    name = str(candidate.get("상품명", "")).strip()
    if name and "상품명" in history.columns:
        return history[history["상품명"].astype(str).str.strip().eq(name)].sort_values("_date")

    return history.iloc[0:0].copy()


def reward_period_label(row: pd.Series) -> str:
    """RAW에 보답 여부 컬럼이 있다면 표시하고, 없으면 일반으로 처리합니다."""
    for col in ["보답", "보답프로그램", "운영구분", "프로그램"]:
        if col in row.index:
            value = str(row.get(col, "")).strip()
            if value and value not in ["-", "nan", "None"]:
                return "보답" if "보답" in value else value
    return "일반"


def candidate_slot_metrics(candidate: pd.Series, target_text: str, history: pd.DataFrame) -> dict:
    hist = match_candidate_history(candidate, history)
    target = parse_target_text(target_text)

    if hist.empty:
        return {
            "이력여부": "신규",
            "추천매출": 0.0,
            "동일타겟평균": 0.0,
            "성연령평균": 0.0,
            "전체평균": 0.0,
            "최고매출": 0.0,
            "최근발송일": None,
            "직전행사가": None,
            "가격증감": None,
            "가격증감률": None,
            "이력": hist,
            "근거": "과거 동일 상품 발송 이력이 없어 신규 TEST 상품으로 후순위 배치",
        }

    h = hist.copy()
    same_exact = h.copy()
    if target["성별"] and "성별" in h.columns:
        same_exact = same_exact[same_exact["성별"].astype(str).str.strip().eq(target["성별"])]
    if target["연령"] and "연령" in h.columns:
        same_exact = same_exact[same_exact["연령"].map(clean_identifier_value).eq(target["연령"])]
    if target["SEG"] and "SEG" in h.columns:
        same_exact = same_exact[same_exact["SEG"].map(clean_identifier_value).eq(target["SEG"])]

    same_demo = h.copy()
    if target["성별"] and "성별" in h.columns:
        same_demo = same_demo[same_demo["성별"].astype(str).str.strip().eq(target["성별"])]
    if target["연령"] and "연령" in h.columns:
        same_demo = same_demo[same_demo["연령"].map(clean_identifier_value).eq(target["연령"])]

    exact_avg = float(same_exact["주문금액"].mean()) if not same_exact.empty else 0.0
    demo_avg = float(same_demo["주문금액"].mean()) if not same_demo.empty else 0.0
    overall_avg = float(h["주문금액"].mean()) if not h.empty else 0.0

    # 매출 우선: 동일 타겟 평균을 가장 강하게, 없으면 성별·연령, 전체 평균 순으로 사용
    if exact_avg > 0:
        expected = exact_avg
        base_reason = f"동일 타겟 과거 평균매출 {compact_money(exact_avg)}"
    elif demo_avg > 0:
        expected = demo_avg * 0.95
        base_reason = f"동일 성별·연령 과거 평균매출 {compact_money(demo_avg)}"
    else:
        expected = overall_avg * 0.85
        base_reason = f"전체 과거 평균매출 {compact_money(overall_avg)}"

    latest = h.sort_values("_date").iloc[-1]
    previous_price = float(latest.get("멤버십혜택가", 0)) if pd.notna(latest.get("멤버십혜택가", pd.NA)) else 0
    current_price = float(candidate.get("행사가", 0))
    price_change = current_price - previous_price if previous_price > 0 else None
    price_rate = price_change / previous_price if previous_price > 0 else None

    # 가격은 매출 다음의 보조 기준으로만 약하게 반영
    if price_rate is not None:
        if price_rate <= -0.05:
            expected *= 1.05
        elif price_rate >= 0.10:
            expected *= 0.90
        elif price_rate >= 0.05:
            expected *= 0.95

    best_row = h.loc[h["주문금액"].idxmax()]
    best_target = target_label(best_row)

    price_reason = ""
    if price_change is not None:
        if price_change > 0:
            price_reason = f"직전 대비 {format_integer_price(price_change)}원 인상"
        elif price_change < 0:
            price_reason = f"직전 대비 {format_integer_price(abs(price_change))}원 인하"
        else:
            price_reason = "직전 발송가와 동일"

    return {
        "이력여부": "이력 있음",
        "추천매출": expected,
        "동일타겟평균": exact_avg,
        "성연령평균": demo_avg,
        "전체평균": overall_avg,
        "최고매출": float(h["주문금액"].max()),
        "최고타겟": best_target,
        "최근발송일": latest["_date"],
        "직전행사가": previous_price,
        "가격증감": price_change,
        "가격증감률": price_rate,
        "이력": h,
        "근거": " · ".join(x for x in [base_reason, price_reason] if x),
    }


def build_schedule_recommendations(
    slots: pd.DataFrame,
    candidates: pd.DataFrame,
    history: pd.DataFrame,
    cooldown_days: int,
    max_weekly_count: int,
) -> tuple[pd.DataFrame, dict]:
    """입력 후보 안에서 매출 우선으로 슬롯별 상품을 자동 배치합니다."""
    result_rows = []
    detail_map = {}
    weekly_counts = {}
    day_products = {}
    planned_dates = pd.to_datetime(slots.get("발송일", pd.Series(dtype=object)), errors="coerce")
    plan_reference = planned_dates.min() if planned_dates.notna().any() else history["_date"].max()

    for slot_idx, slot in slots.iterrows():
        target = str(slot.get("타겟", "")).strip()
        product_count = int(float(slot.get("상품수", 0) or 0))
        if not target or product_count <= 0:
            continue

        slot_date = pd.to_datetime(slot.get("발송일"), errors="coerce")
        day_key = slot_date.strftime("%Y-%m-%d") if pd.notna(slot_date) else f"slot-{slot_idx}"
        day_products.setdefault(day_key, set())

        # 같은 슬롯 안에서는 동일 상품이 여러 행으로 입력되어도 1개 후보로만 사용합니다.
        # 상품 식별 우선순위: 쇼라코드 → 알파코드 → 상품명
        ranked_by_product = {}
        for cand_idx, candidate in candidates.iterrows():
            # 성별 전용 상품은 타겟 부적합 시 점수 계산 전에 후보군에서 제외합니다.
            if not is_candidate_gender_compatible(candidate, target):
                continue

            product_key = (
                clean_identifier_value(candidate.get("쇼라코드", ""))
                or clean_identifier_value(candidate.get("알파코드", ""))
                or str(candidate.get("상품명", "")).strip()
            )
            if not product_key:
                continue
            if product_key in day_products[day_key]:
                continue
            if weekly_counts.get(product_key, 0) >= max_weekly_count:
                continue

            metrics = candidate_slot_metrics(candidate, target, history)
            latest_date = metrics.get("최근발송일")
            if latest_date is not None and pd.notna(latest_date) and pd.notna(plan_reference):
                elapsed = (pd.Timestamp(plan_reference) - pd.Timestamp(latest_date)).days
                if elapsed < cooldown_days:
                    continue

            candidate_rank = (
                float(metrics.get("추천매출", 0) or 0),
                float(candidate.get("할인율계산값", 0) or 0),
            )
            previous = ranked_by_product.get(product_key)
            if previous is None or candidate_rank > previous[0]:
                ranked_by_product[product_key] = (
                    candidate_rank,
                    cand_idx,
                    product_key,
                    metrics,
                )

        ranked = [
            (rank_key[0], cand_idx, product_key, metrics)
            for rank_key, cand_idx, product_key, metrics in ranked_by_product.values()
        ]
        ranked.sort(
            key=lambda x: (
                x[0],
                float(candidates.loc[x[1], "할인율계산값"] or 0),
            ),
            reverse=True,
        )
        selected = ranked[:product_count]

        for order_no, (_, cand_idx, product_key, metrics) in enumerate(selected, start=1):
            candidate = candidates.loc[cand_idx]
            row = {
                "발송일": slot_date.strftime("%Y-%m-%d") if pd.notna(slot_date) else "",
                "시간대": str(slot.get("시간대", "")).strip(),
                "타겟": target,
                "전시순서": order_no,
                "알파코드": clean_identifier_value(candidate.get("알파코드", "")),
                "쇼라코드": clean_identifier_value(candidate.get("쇼라코드", "")),
                "상품명": str(candidate.get("상품명", "")).strip(),
                "정상가": float(candidate.get("정상가", 0)),
                "행사가": float(candidate.get("행사가", 0)),
                "할인율": candidate.get("할인율", ""),
                "예상매출": metrics["추천매출"],
            }
            result_rows.append(row)
            detail_map[(row["발송일"], row["시간대"], target, row["알파코드"], row["쇼라코드"], row["상품명"])] = metrics
            weekly_counts[product_key] = weekly_counts.get(product_key, 0) + 1
            day_products[day_key].add(product_key)

    return pd.DataFrame(result_rows), detail_map


def schedule_history_table(hist: pd.DataFrame, current_price: float) -> pd.DataFrame:
    if hist.empty:
        return pd.DataFrame(columns=[
            "발송일", "타겟", "소재", "멤버십혜택가", "현재가 대비", "주문금액", "프로모션"
        ])

    view = hist.sort_values("_date", ascending=False).copy()
    view["_date"] = pd.to_datetime(view.get("_date"), errors="coerce")
    view["발송일"] = view["_date"].dt.strftime("%Y-%m-%d").fillna("-")
    view["타겟"] = view.apply(target_label, axis=1)

    promotions = st.session_state.get(
        "promotions",
        pd.DataFrame(columns=["프로모션명", "_start_date", "_end_date", "스킴"]),
    )
    view["프로모션"] = view["_date"].map(
        lambda value: promotion_name_for_date(value, promotions)
    )

    view["현재가 대비"] = current_price - view["멤버십혜택가"]
    cols = ["발송일", "타겟", "소재", "멤버십혜택가", "현재가 대비", "주문금액", "프로모션"]
    return view[[c for c in cols if c in view.columns]].head(20)


def schedule_target_summary(hist: pd.DataFrame) -> pd.DataFrame:
    if hist.empty:
        return pd.DataFrame(columns=[
            "타겟", "운영횟수", "평균매출", "최고매출", "최근발송일", "프로모션"
        ])

    h = hist.copy()
    h["타겟"] = h.apply(target_label, axis=1)
    h["_date"] = pd.to_datetime(h.get("_date"), errors="coerce")

    summary = h.groupby("타겟", as_index=False).agg(
        운영횟수=("상품명", "size"),
        평균매출=("주문금액", "mean"),
        최고매출=("주문금액", "max"),
        최근발송일=("_date", "max"),
    ).sort_values(["평균매출", "최고매출"], ascending=False)

    promotions = st.session_state.get(
        "promotions",
        pd.DataFrame(columns=["프로모션명", "_start_date", "_end_date", "스킴"]),
    )
    summary["프로모션"] = summary["최근발송일"].map(
        lambda value: promotion_name_for_date(value, promotions)
    )
    summary["최근발송일"] = summary["최근발송일"].dt.strftime("%Y-%m-%d").fillna("-")

    return summary[[
        "타겟", "운영횟수", "평균매출", "최고매출", "최근발송일", "프로모션"
    ]]


# ─────────────────────────────────────────────────────────────────────────────
# 데이터 연결
# ─────────────────────────────────────────────────────────────────────────────
DEFAULT_GOOGLE_SHEET_URL = (
    "https://docs.google.com/spreadsheets/d/"
    "1I8sAfs8kfMAFThHa_o-aeb2GLWLbFtxf3FxBhA8q-tQ/edit?gid=0#gid=0"
)

if "products" not in st.session_state:
    st.session_state.products = None
    st.session_state.sends = None
    st.session_state.source_name = None
    st.session_state.synced_at = None
    st.session_state.google_sync_error = None
    st.session_state.lowest = pd.DataFrame()
    st.session_state.messages = pd.DataFrame(columns=["캠페인명", "MMS문구"])
    st.session_state.promotions = pd.DataFrame(columns=["프로모션명", "_start_date", "_end_date", "스킴"])
    st.session_state.auto_sync_attempted = False

st.sidebar.markdown(
    """
    <div style="padding: 0.15rem 0 1rem 0;">
        <div style="font-size: 24px; font-weight: 800; letter-spacing: -0.5px; line-height: 1.3;">
            MMS AI 대시보드
        </div>
        <div style="font-size: 13px; color: #6b7280; line-height: 1.55; margin-top: 5px;">
            MMS 일일·주간·월간 실적 및<br>상품 인사이트 통합 분석
        </div>
    </div>
    """,
    unsafe_allow_html=True,
)

source = st.sidebar.radio(
    "데이터 연결",
    ["구글시트 자동연동", "엑셀 업로드"],
    index=0,
)

if source == "구글시트 자동연동":
    sheet_url = st.sidebar.text_input(
        "🔗 구글시트 링크",
        value=DEFAULT_GOOGLE_SHEET_URL,
    )

    # 앱 최초 실행 시 자동 동기화
    if not st.session_state.auto_sync_attempted:
        st.session_state.auto_sync_attempted = True
        try:
            with st.spinner("구글시트 최신 데이터를 자동으로 불러오는 중입니다."):
                sync_google_sheet(sheet_url)
        except Exception as exc:
            st.session_state.google_sync_error = str(exc)

    if st.sidebar.button("🔄 지금 새로고침", use_container_width=True):
        try:
            with st.spinner("구글시트 최신 데이터를 다시 불러오는 중입니다."):
                sync_google_sheet(sheet_url, force=True)
            st.sidebar.success("최신 데이터로 갱신했습니다.")
            st.rerun()
        except Exception as exc:
            st.session_state.google_sync_error = str(exc)
            st.sidebar.error(f"새로고침 실패: {exc}")

    st.sidebar.caption("자동 갱신 주기: 5분")
    if st.session_state.google_sync_error:
        st.sidebar.warning(
            "구글시트 자동연동에 실패했습니다. "
            "공유 권한과 탭 이름(상품, 소재)을 확인하거나 엑셀 업로드를 이용해주세요."
        )
        with st.sidebar.expander("오류 상세"):
            st.code(st.session_state.google_sync_error)

else:
    uploaded = st.sidebar.file_uploader("📁 MMS 파일 업로드", type=["xlsx", "xlsm"])
    if uploaded is not None:
        try:
            products, sends, lowest, messages, promotions = load_excel_bytes(uploaded.getvalue())
            st.session_state.products = products
            st.session_state.sends = sends
            st.session_state.lowest = lowest
            st.session_state.messages = messages
            st.session_state.promotions = promotions
            st.session_state.source_name = uploaded.name
            st.session_state.synced_at = datetime.now()
            st.session_state.google_sync_error = None
        except Exception as exc:
            st.sidebar.error(f"파일을 불러오지 못했습니다: {exc}")

if st.session_state.products is None or st.session_state.sends is None:
    st.info(
        "구글시트 자동연동을 확인 중입니다. 연결되지 않으면 "
        "왼쪽에서 엑셀 업로드로 전환해주세요."
    )
    st.stop()

products = st.session_state.products
sends = st.session_state.sends
lowest = st.session_state.get("lowest", pd.DataFrame())
messages = st.session_state.get("messages", pd.DataFrame(columns=["캠페인명", "MMS문구"]))
promotions = st.session_state.get("promotions", pd.DataFrame(columns=["프로모션명", "_start_date", "_end_date", "스킴"]))

_menu_options = ["홈", "일일실적", "주간실적", "상품구분", "타겟분석", "편성 프로그램"]
_menu_slug_to_name = {
    "home": "홈",
    "daily": "일일실적",
    "weekly": "주간실적",
    "product": "상품구분",
    "target": "타겟분석",
    "planning": "편성 프로그램",
}
_menu_name_to_slug = {v: k for k, v in _menu_slug_to_name.items()}

_query_menu = _get_query_param("menu").strip().lower()
_query_menu_name = _menu_slug_to_name.get(_query_menu, "홈")
_menu_default_index = _menu_options.index(_query_menu_name) if _query_menu_name in _menu_options else 0

menu = st.sidebar.radio(
    "메뉴",
    _menu_options,
    index=_menu_default_index,
)

# 현재 메뉴를 URL에 반영하여 주소를 그대로 공유하면 해당 메뉴로 직접 진입.
try:
    st.query_params["menu"] = _menu_name_to_slug.get(menu, "home")
except Exception:
    pass

# 메뉴별로 사용하지 않는 세부 query param은 정리
try:
    if menu != "주간실적":
        for _qp_key in ["year", "week"]:
            if _qp_key in st.query_params:
                del st.query_params[_qp_key]
    if menu != "일일실적" and "date" in st.query_params:
        del st.query_params["date"]
except Exception:
    pass



# ─────────────────────────────────────────────────────────────────────────────
# 홈
# ─────────────────────────────────────────────────────────────────────────────
if menu == "홈":
    st.markdown('<div class="section-title">홈 · 기간별 실적</div>', unsafe_allow_html=True)
    st.caption("🔗 현재 브라우저 주소를 그대로 공유하면 홈 화면으로 바로 연결됩니다.")

    monthly_all = aggregate_send(sends, "Monthly")
    weekly_all = aggregate_send(sends, "Weekly")
    daily_all = aggregate_send(sends, "Daily")

    # 월간 기간 필터
    c1, c2, c3 = st.columns([1.2, 1.2, 1.2])
    with c1:
        month_option = st.selectbox("월간 조회 기간", ["전체", "최근 3개월", "최근 6개월", "최근 12개월", "직접 선택"])
    month_labels = monthly_all["_label"].astype(str).tolist()
    start_month = end_month = None
    if month_option == "직접 선택" and month_labels:
        with c2:
            start_month = st.selectbox("시작월", month_labels, index=0)
        with c3:
            end_month = st.selectbox("종료월", month_labels, index=len(month_labels)-1)

    monthly = filter_monthly_period(monthly_all, month_option, start_month, end_month)

    weekly = weekly_all.copy()

    # KPI
    latest_df = monthly if not monthly.empty else monthly_all
    if not latest_df.empty:
        latest = latest_df.iloc[-1]
        cards = [
            ("주문금액", f"{fmt_num(latest['주문금액'])}원", delta_for_latest(latest_df, "주문금액")),
            ("발송건수", fmt_num(latest["발송건수"]), delta_for_latest(latest_df, "발송건수")),
            ("CTR", fmt_pct(latest["반응율(Uniq CTR)"]), delta_for_latest(latest_df, "반응율(Uniq CTR)", pp=True)),
            ("클릭 CVR", fmt_pct(latest["클릭 CVR"]), delta_for_latest(latest_df, "클릭 CVR", pp=True)),
            ("SPM", f"{latest['발송대비매출(SPM)']:.1f}", delta_for_latest(latest_df, "발송대비매출(SPM)")),
            ("발송당매출", f"{fmt_num(latest['발송당매출(발송횟수)'])}원", delta_for_latest(latest_df, "발송당매출(발송횟수)")),
        ]
        cols = st.columns(6)
        for col, (label, value, delta) in zip(cols, cards):
            col.markdown(
                f'<div class="metric-card"><div class="metric-label">{label}</div>'
                f'<div class="metric-value">{value}</div>'
                f'<div class="metric-delta">직전 대비 {delta}</div></div>',
                unsafe_allow_html=True,
            )

    # 월간 그래프와 표
    st.markdown('<div class="section-title">월별 SPM / 발송대비매출</div>', unsafe_allow_html=True)
    st.plotly_chart(
        trend_chart(monthly, "월별 SPM / 발송대비매출", "#fdbb00"),
        use_container_width=True,
        config={"displayModeBar": False},
    )
    st.dataframe(
        clean_identifier_columns(format_home_table_with_summary(monthly, "Monthly")),
        use_container_width=True,
        hide_index=True,
        height=400,
    )

    # 주간 그래프와 표 - 독립 조회 기간
    st.markdown('<div class="section-title">주간 SPM / 발송대비매출</div>', unsafe_allow_html=True)
    st.caption("주간 조회 기간")
    week_labels = weekly_all["_label"].astype(str).tolist()
    if week_labels:
        week_start_col, week_end_col = st.columns(2)
        with week_start_col:
            start_week = st.selectbox("시작 주차", week_labels, index=0, key="home_week_start")
        with week_end_col:
            end_week = st.selectbox("종료 주차", week_labels, index=len(week_labels)-1, key="home_week_end")
        weekly = filter_weekly_period(weekly_all, start_week, end_week)
    else:
        weekly = weekly_all.copy()
    if weekly.empty:
        st.info("선택한 기간의 주간 데이터가 없습니다.")
    else:
        st.plotly_chart(
            trend_chart(weekly, "주간 SPM / 발송대비매출", "#70ad47"),
            use_container_width=True,
            config={"displayModeBar": False},
        )
        st.dataframe(
            clean_identifier_columns(format_home_table_with_summary(weekly, "Weekly")),
            use_container_width=True,
            hide_index=True,
            height=520,
        )

    # 일간 표: 실제 날짜 필터
    st.markdown('<div class="section-title">Daily</div>', unsafe_allow_html=True)
    st.caption("일간 조회 기간")
    daily_start_col, daily_end_col = st.columns(2)
    with daily_start_col:
        daily_start_date = st.date_input(
            "시작일",
            value=sends["_date"].min().date(),
            min_value=sends["_date"].min().date(),
            max_value=sends["_date"].max().date(),
            format="YYYY/MM/DD",
            key="home_daily_start",
        )
    with daily_end_col:
        daily_end_date = st.date_input(
            "종료일",
            value=sends["_date"].max().date(),
            min_value=sends["_date"].min().date(),
            max_value=sends["_date"].max().date(),
            format="YYYY/MM/DD",
            key="home_daily_end",
        )
    if daily_start_date > daily_end_date:
        daily_start_date, daily_end_date = daily_end_date, daily_start_date
    daily_source = sends[
        (sends["_date"].dt.date >= daily_start_date)
        & (sends["_date"].dt.date <= daily_end_date)
    ].copy()
    daily = aggregate_send(daily_source, "Daily")
    if daily.empty:
        st.info("선택한 기간의 일간 데이터가 없습니다.")
    else:
        st.dataframe(
            clean_identifier_columns(format_home_table_with_summary(daily, "Daily")),
            use_container_width=True,
            hide_index=True,
            height=480,
        )


# ─────────────────────────────────────────────────────────────────────────────
# 일일실적
# ─────────────────────────────────────────────────────────────────────────────
elif menu == "일일실적":
    st.markdown('<div class="section-title">📊 일일실적 분석</div>', unsafe_allow_html=True)
    dates = sorted(products["_date"].dt.date.unique(), reverse=True)

    _query_date_raw = _get_query_param("date").strip()
    try:
        _query_date = pd.to_datetime(_query_date_raw, errors="raise").date() if _query_date_raw else None
    except Exception:
        _query_date = None

    _daily_default_index = dates.index(_query_date) if _query_date in dates else 0
    selected_date = st.selectbox(
        "📅 날짜 선택",
        dates,
        index=_daily_default_index,
        key="daily_selected_date",
    )

    try:
        st.query_params["menu"] = "daily"
        st.query_params["date"] = selected_date.isoformat()
    except Exception:
        pass

    st.caption("🔗 현재 브라우저 주소를 그대로 공유하면 이 날짜의 일일실적으로 바로 연결됩니다.")

    pday = products[products["_date"].dt.date == selected_date].copy()
    sday = sends[sends["_date"].dt.date == selected_date].copy()
    pday = merge_lowest_price(pday, lowest)

    if pday.empty and sday.empty:
        st.info("선택한 날짜의 데이터가 없습니다.")
        st.stop()

    total_amount = pday["주문금액"].sum()
    total_orders = pday["주문건수"].sum()
    total_qty = pday["주문수량"].sum()
    send_col = first_col(sday, ["발송 성공 건수", "총 발송 건수"])
    click_col = first_col(sday, ["클릭 수(uniq)", "클릭 수"])
    send_success = sday[send_col].sum() if send_col else 0
    clicks = sday[click_col].sum() if click_col else 0

    cards = [
        ("주문건수", f"{int(total_orders):,}건"),
        ("주문수량", f"{int(total_qty):,}개"),
        ("주문금액", f"{int(total_amount):,}원"),
        ("CTR", f"{(clicks/send_success*100 if send_success else 0):.1f}%"),
        ("CVR", f"{(total_orders/clicks*100 if clicks else 0):.1f}%"),
        ("SPM", f"{(total_amount/send_success if send_success else 0):.1f}"),
    ]
    metric_cols = st.columns(len(cards))
    for c, (label, value) in zip(metric_cols, cards):
        c.markdown(
            f'<div class="metric-card"><div class="metric-label">{label}</div>'
            f'<div class="metric-value">{value}</div></div>',
            unsafe_allow_html=True,
        )

    # 오전/오후 또는 소재 단위로 분리
    # 운영 이슈는 선택 날짜의 오전·오후 전체 상품을 묶어 상단에서 한 번만 관리
    st.markdown('<div class="subsection-title">운영 이슈</div>', unsafe_allow_html=True)
    issue_products = pday.copy()
    issue_products = issue_products.sort_values([c for c in ["시간대", "전시순서", "상품명"] if c in issue_products.columns])
    issue_rows = list(issue_products.iterrows())
    if issue_rows:
        issue_options = list(range(len(issue_rows)))
        def issue_option_label(option_idx: int) -> str:
            _, option_row = issue_rows[option_idx]
            return str(option_row.get("상품명", "상품명 없음"))
        issue_cols = st.columns([2.4, 1.2, 3.2, 0.8, 0.8], gap="small")
        with issue_cols[0]:
            selected_issue_idx = st.selectbox("상품 선택", issue_options, format_func=issue_option_label, key=f"daily_issue_product_top_{selected_date}")
        selected_issue_row = issue_rows[selected_issue_idx][1]
        saved_issue = get_saved_issue(selected_issue_row)
        saved_type = (saved_issue.get("유형") or ["선택 안 함"])[0]
        type_options = ["선택 안 함", "판매중단", "가격오류", "기타"]
        with issue_cols[1]:
            issue_type = st.selectbox("이슈 유형", type_options, index=type_options.index(saved_type) if saved_type in type_options else 0, key=f"daily_issue_type_{selected_date}_{selected_issue_idx}")
        with issue_cols[2]:
            issue_memo = st.text_input("상세 메모", value=saved_issue.get("메모", ""), placeholder="예: 11시 30분 판매중단 발생", key=f"daily_issue_memo_{selected_date}_{selected_issue_idx}")
        with issue_cols[3]:
            st.write("")
            st.write("")
            if st.button("저장", use_container_width=True, key=f"daily_issue_save_{selected_date}_{selected_issue_idx}"):
                save_operation_issue(selected_issue_row, issue_type, issue_memo)
                st.success("저장했습니다.")
                st.rerun()
        with issue_cols[4]:
            st.write("")
            st.write("")
            if st.button("삭제", use_container_width=True, key=f"daily_issue_delete_{selected_date}_{selected_issue_idx}"):
                delete_operation_issue(selected_issue_row)
                st.success("삭제했습니다.")
                st.rerun()
    else:
        st.caption("등록할 상품이 없습니다.")

    if "시간대" in sday.columns:
        sday["_sort_time"] = pd.to_datetime(sday["시간대"].astype(str), errors="coerce")
        sday = sday.sort_values(["_sort_time", "소재"] if "소재" in sday.columns else ["_sort_time"])

    for idx, send_row in sday.iterrows():
        time_value = str(send_row.get("시간대", ""))
        material = str(send_row.get("소재", ""))
        part_title = f"{time_value} · {material}" if material else time_value
        st.markdown(f'<div class="section-title">🕒 {part_title}</div>', unsafe_allow_html=True)

        matched = pday.copy()
        if "시간대" in matched.columns and time_value:
            matched = matched[matched["시간대"].astype(str) == time_value]
        if "소재" in matched.columns and material:
            material_match = matched[matched["소재"].astype(str) == material]
            if not material_match.empty:
                matched = material_match

        if matched.empty:
            st.info("해당 발송 건과 연결된 상품 실적이 없습니다.")
            continue

        asset_key = daily_asset_key(selected_date, time_value)
        campaign_name = str(send_row.get("캠페인명", "")).strip()
        image_path = find_daily_image(asset_key, campaign_name)
        message_text = extract_mms_message(matched, send_row, messages)

        st.markdown('<div class="subsection-title">발송 소재</div>', unsafe_allow_html=True)

        import base64
        import html

        if image_path is not None:
            mime_map = {
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".png": "image/png",
                ".webp": "image/webp",
            }
            mime = mime_map.get(image_path.suffix.lower(), "image/jpeg")
            encoded = base64.b64encode(image_path.read_bytes()).decode("utf-8")
            image_body = (
                f'<img src="data:{mime};base64,{encoded}" '
                f'alt="{html.escape(image_path.name)}">'
            )
        else:
            image_body = (
                '<div class="asset-empty">images 폴더에<br>'
                + html.escape(f"{asset_key}_ 로 시작하는 이미지가 없습니다.")
                + '</div>'
            )

        if message_text:
            clean_message_text = str(message_text).lstrip()
            message_body = html.escape(clean_message_text).replace("\n", "<br>")
        else:
            message_body = (
                "로우데이터의 MMS문구 컬럼에 문구를 입력하면 "
                "이곳에 자동으로 표시됩니다."
            )

        # 같은 HTML Grid 안에서 렌더링해 좌우 카드가 항상 동일한 세로 높이를 사용합니다.
        # 이미지는 object-fit: contain으로 원본 비율/내용을 변경하지 않습니다.
        asset_pair_html = f"""
        <div class="daily-asset-pair">
            <div class="asset-card asset-image-card">{image_body}</div>
            <div class="asset-card asset-message-card">{message_body}</div>
        </div>
        """
        st.markdown(asset_pair_html, unsafe_allow_html=True)

        # 발송 통계: 요청 컬럼만 표시
        send_count = float(send_row.get(send_col, 0)) if send_col else 0
        click_count = float(send_row.get(click_col, 0)) if click_col else 0
        orders = float(send_row.get("주문건수", 0))
        amount = float(send_row.get("주문금액", 0))
        send_view = pd.DataFrame([{
            "성별": send_row.get("성별", ""),
            "연령": send_row.get("연령", ""),
            "SEG": send_row.get("SEG", ""),
            "소재": material,
            "URL": send_row.get("URL", ""),
            "발송건수": fmt_num(send_count),
            "클릭수": fmt_num(click_count),
            "CTR": fmt_pct(click_count / send_count if send_count else 0),
            "CVR": fmt_pct(orders / click_count if click_count else 0),
            "객단가": fmt_num(amount / orders if orders else 0),
            "SPM": f"{(amount/send_count if send_count else 0):.1f}",
        }])
        st.markdown('<div class="subsection-title">발송 통계</div>', unsafe_allow_html=True)
        st.dataframe(
            clean_identifier_columns(send_view),
            use_container_width=True,
            hide_index=True,
        )

        matched["상품등급"] = matched["주문금액"].apply(product_grade)
        matched = add_history_columns(matched, products)
        sort_cols = [c for c in ["전시순서", "상품명"] if c in matched.columns]
        if sort_cols:
            matched = matched.sort_values(sort_cols)

        display_cols = [
            "전시순서", "MD", "알파코드", "쇼라코드", "상품명",
            "정상가", "멤버십혜택가", "할인율", "추가노출",
            "주문건수", "주문수량", "주문금액",
            "재편성", "발송일 최저가", "최저가 확보",
            "최고매출", "최고일자", "최고타겟"
        ]
        product_view = matched[[c for c in display_cols if c in matched.columns]].copy()

        if "할인율" in product_view.columns:
            product_view["할인율"] = product_view["할인율"].map(format_discount_percent)
        for price_col in ["정상가", "멤버십혜택가", "주문금액", "최고매출", "발송일 최저가"]:
            if price_col in product_view.columns:
                product_view[price_col] = product_view[price_col].map(format_integer_price)

        if "최저가 확보" in product_view.columns:
            product_view = product_view.rename(columns={"최저가 확보": "최저가 여부"})

        # 합계행
        total_row = {c: "" for c in product_view.columns}
        first_display = product_view.columns[0]
        total_row[first_display] = "합계"
        for c in ["주문건수", "주문수량", "주문금액"]:
            if c in product_view.columns:
                total_row[c] = matched[c].sum()
        product_view = pd.concat([product_view, pd.DataFrame([total_row])], ignore_index=True)
        for c in ["주문건수", "주문수량", "주문금액"]:
            if c in product_view.columns:
                product_view[c] = product_view[c].map(
                    lambda x: format_integer_price(x) if str(x).strip() not in ["", "nan", "None"] else ""
                )

        st.markdown('<div class="subsection-title">상품 실적</div>', unsafe_allow_html=True)
        st.dataframe(
            clean_identifier_columns(product_view),
            use_container_width=True,
            hide_index=True,
            height=280,
        )

        st.markdown('<div class="subsection-title">상품 인사이트</div>', unsafe_allow_html=True)
        st.caption("상품별 영역은 기본 접힘 상태이며, 클릭하면 주요 인사이트와 최근 발송 이력을 확인할 수 있습니다.")

        for _, product_row in matched.iterrows():
            saved_issue = get_saved_issue(product_row)
            report = generate_insight_report(product_row, products, saved_issue)
            product_name = report["상품명"] or "상품명 없음"
            amount_text = compact_money(float(product_row.get("주문금액", 0) or 0))
            grade_text = report["상품등급"]
            target_text = target_label(product_row) or "타겟 정보 없음"

            with st.expander(
                f"{product_name} · {target_text} · {amount_text} · {grade_text}",
                expanded=False,
            ):
                rows_html = []
                for item in report["인사이트"]:
                    insight_text = format_daily_insight_item(item)
                    rows_html.append(f"<div class='insight-row'>{html.escape(insight_text)}</div>")
                st.markdown("<div class='compact-insight'>" + "".join(rows_html) + "</div>", unsafe_allow_html=True)

                if report["위험요인"]:
                    st.caption("주의: " + " · ".join(report["위험요인"]))

                st.markdown("**최근 발송 이력**")
                if report["발송이력"].empty:
                    st.caption("동일 상품의 발송 이력이 없습니다.")
                else:
                    st.dataframe(
                        clean_identifier_columns(report["발송이력"]),
                        use_container_width=True,
                        hide_index=True,
                        height=min(210, 38 * (len(report["발송이력"]) + 1)),
                    )

    all_insights = [make_insight(r, products) for _, r in pday.iterrows()]
    ppt = build_ppt(
        f"{selected_date} MMS 일일실적",
        all_insights,
        pday[[c for c in ["상품명", "주문건수", "주문수량", "주문금액"] if c in pday.columns]],
    )
    st.download_button(
        "📥 PPT 다운로드",
        ppt,
        file_name=f"{selected_date}_MMS_일일실적.pptx",
    )


# ─────────────────────────────────────────────────────────────────────────────
# 주간실적
# ─────────────────────────────────────────────────────────────────────────────
elif menu == "주간실적":
    st.markdown('<div class="section-title">📈 주간실적 분석</div>', unsafe_allow_html=True)

    # 연도 → 주차 순서로 선택하여 같은 주차명이 연도별로 섞이지 않도록 함
    # 공유 URL 예: ?menu=weekly&year=2026&week=0713주차
    weekly_years = sorted(products["_year"].dropna().astype(int).unique(), reverse=True)

    _query_year_raw = _get_query_param("year")
    try:
        _query_year = int(_query_year_raw)
    except (TypeError, ValueError):
        _query_year = None

    _year_default_index = (
        weekly_years.index(_query_year)
        if _query_year in weekly_years
        else 0
    )
    selected_year = st.selectbox(
        "연도 선택",
        weekly_years,
        index=_year_default_index,
        key="weekly_selected_year",
    )

    year_products = products[products["_year"] == selected_year].copy()
    year_sends = sends[sends["_year"] == selected_year].copy()
    week_order = (
        year_products.groupby("주차")["_date"].min().sort_values()
        if not year_products.empty else pd.Series(dtype="datetime64[ns]")
    )
    weeks = [str(x) for x in week_order.index]
    if not weeks:
        st.info("선택한 연도의 주차 데이터가 없습니다.")
        st.stop()

    _query_week = _get_query_param("week").strip()
    _week_default_index = (
        weeks.index(_query_week)
        if _query_week in weeks
        else len(weeks) - 1
    )
    week = st.selectbox(
        "주차 선택",
        weeks,
        index=_week_default_index,
        key="weekly_selected_week",
    )

    # 현재 선택값을 주소창 URL에 즉시 반영.
    # 이 주소를 그대로 복사해 공유하면 동일 연도·주차로 바로 진입합니다.
    _set_weekly_deeplink(selected_year, week)
    st.caption("🔗 현재 브라우저 주소를 그대로 공유하면 이 주차 화면으로 바로 연결됩니다.")

    pw = year_products[year_products["주차"].astype(str) == week].copy()
    sw = year_sends[year_sends["주차"].astype(str) == week].copy()

    if pw.empty or sw.empty:
        st.info("선택한 연도·주차의 상품 또는 소재 데이터가 없습니다.")
        st.stop()

    send_col = first_col(sw, ["발송 성공 건수", "총 발송 건수"])
    click_col = first_col(sw, ["클릭 수(uniq)", "클릭 수"])
    send_count = sw[send_col].sum()
    click_count = sw[click_col].sum()
    order_count = sw["주문건수"].sum()
    amount = sw["주문금액"].sum()
    ctr = click_count / send_count if send_count else 0
    cvr = order_count / click_count if click_count else 0
    aov = amount / order_count if order_count else 0
    spm = amount / send_count if send_count else 0

    # 전주 대비 카드 증감
    year_week_names = [str(x) for x in (
        year_sends.groupby("주차")["_date"].min().sort_values().index
    )]
    prev_sw = pd.DataFrame()
    if week in year_week_names and year_week_names.index(week) > 0:
        prev_week = year_week_names[year_week_names.index(week) - 1]
        prev_sw = year_sends[year_sends["주차"].astype(str) == prev_week]

    def prev_metric(column, default=0):
        return float(prev_sw[column].sum()) if not prev_sw.empty and column in prev_sw.columns else default

    prev_send = prev_metric(send_col)
    prev_click = prev_metric(click_col)
    prev_orders = prev_metric("주문건수")
    prev_amount = prev_metric("주문금액")
    prev_ctr = prev_click / prev_send if prev_send else 0
    prev_cvr = prev_orders / prev_click if prev_click else 0
    prev_aov = prev_amount / prev_orders if prev_orders else 0
    prev_spm = prev_amount / prev_send if prev_send else 0

    cards = [
        ("발송횟수", f"{len(sw):,}회", weekly_delta(len(sw), len(prev_sw)) if not prev_sw.empty else "-"),
        ("상품수", f"{len(pw):,}건", "-"),
        ("발송건수", f"{int(send_count):,}", weekly_delta(send_count, prev_send) if prev_sw is not None and not prev_sw.empty else "-"),
        ("주문건수", f"{int(order_count):,}건", weekly_delta(order_count, prev_orders) if not prev_sw.empty else "-"),
        ("주문금액", f"{int(amount):,}원", weekly_delta(amount, prev_amount) if not prev_sw.empty else "-"),
        ("CTR", f"{ctr*100:.1f}%", weekly_delta(ctr, prev_ctr, pp=True) if not prev_sw.empty else "-"),
        ("CVR", f"{cvr*100:.1f}%", weekly_delta(cvr, prev_cvr, pp=True) if not prev_sw.empty else "-"),
        ("객단가", f"{int(aov):,}원", weekly_delta(aov, prev_aov) if not prev_sw.empty else "-"),
        ("SPM", f"{spm:.1f}", weekly_delta(spm, prev_spm) if not prev_sw.empty else "-"),
    ]
    card_cols = st.columns(5)
    for i, (label, value, delta) in enumerate(cards):
        with card_cols[i % 5]:
            st.markdown(
                f'<div class="metric-card"><div class="metric-label">{label}</div>'
                f'<div class="metric-value">{value}</div>'
                f'<div class="metric-delta">전주 대비 {delta}</div></div>',
                unsafe_allow_html=True,
            )
        if i == 4:
            card_cols = st.columns(4)

    chart_left, chart_right = st.columns(2)
    with chart_left:
        st.plotly_chart(
            weekly_product_chart(sw),
            use_container_width=True,
            config={"displayModeBar": False},
        )
    with chart_right:
        st.plotly_chart(
            weekly_send_chart(sw),
            use_container_width=True,
            config={"displayModeBar": False},
        )

    # 대·중카테고리 편성 및 주문 비중
    cat_left, cat_right = st.columns(2)

    with cat_left:
        big_table = category_summary_table(pw, "대카", week, selected_year)
        st.plotly_chart(
            category_pie_chart(big_table, "대카테고리 주문비중"),
            use_container_width=True,
            config={"displayModeBar": False},
        )
        _weekly_table_title("대카테고리 편성 및 주문 비중")
        _big_table_display = clean_identifier_columns(weekly_display_format(big_table))
        st.dataframe(
            _style_weekly_category_total(_big_table_display),
            use_container_width=True,
            hide_index=True,
            height=430,
        )

    with cat_right:
        mid_table = category_summary_table(pw, "중카", week, selected_year)
        st.plotly_chart(
            category_pie_chart(mid_table, "중카테고리 주문비중"),
            use_container_width=True,
            config={"displayModeBar": False},
        )
        _weekly_table_title("중카테고리 편성 및 주문 비중")
        _mid_table_display = clean_identifier_columns(weekly_display_format(mid_table))
        st.dataframe(
            _style_weekly_category_total(_mid_table_display),
            use_container_width=True,
            hide_index=True,
            height=560,
        )

    tabs = st.tabs([
        "주간실적 분석", "상품 실적", "소재 실적",
        "SEG 실적", "요일 실적", "시간대 실적", "MMS 상품"
    ])

    with tabs[0]:
        report = build_weekly_analysis(
            week, selected_year, pw, sw, products, sends
        )
        # 4개 의사결정 섹션: 제목은 굵게, 내용은 불필요한 빈 줄 없이 한 줄씩 표시
        report_lines = [line.strip() for line in report.splitlines() if line.strip()]
        report_html = []
        for line in report_lines:
            if line.startswith("■ "):
                report_html.append(
                    f'<div style="font-weight:700; margin-top:14px; margin-bottom:4px;">{line}</div>'
                )
            else:
                report_html.append(
                    f'<div style="margin:0 0 3px 0; line-height:1.55;">{line}</div>'
                )
        st.markdown(
            '<div class="insight-box">' + "".join(report_html) + '</div>',
            unsafe_allow_html=True,
        )

        st.markdown("### 상세 데이터 보기")

        # MD 의사결정용 상세
        try:
            _week_end = pd.to_datetime(pw["_date"], errors="coerce").max()
            _md_rec_df, _md_src_df = _md_recommendation_tables(products, pw, _week_end)

            with st.expander("▶ 금주 재편성 추천 상품", expanded=False):
                if _md_rec_df.empty:
                    st.caption("근거 기준을 충족한 재편성 추천 상품이 없습니다.")
                else:
                    st.dataframe(_md_rec_df, use_container_width=True, hide_index=True)

            with st.expander("▶ 신규·유사신규 소싱 제안", expanded=False):
                if _md_src_df.empty:
                    st.caption("전년·과거 동시즌 고성과 근거를 충족한 소싱 제안이 없습니다.")
                else:
                    st.dataframe(_md_src_df, use_container_width=True, hide_index=True)
        except Exception as _md_exc:
            st.caption(f"MD 상세 분석을 불러오지 못했습니다: {type(_md_exc).__name__}")
        detail_sections = [
            "MMS 상품 실적",
            "MMS 발송 통계",
            "카테고리 분석",
            "SEG 분석",
            "요일·시간대 분석",
            "상품별 상세 인사이트",
            "최저가 미확보 상품",
        ]
        detail_report = build_weekly_detail_analysis(
            week, selected_year, pw, sw, products, sends
        )
        # 상세 데이터 보기에서도 화면용 상품명 축약을 동일 적용
        detail_product_names = sorted(
            pw["상품명"].dropna().astype(str).unique().tolist(),
            key=len,
            reverse=True,
        )
        for original_name in detail_product_names:
            short_name = _short_weekly_product_name(original_name)
            if short_name and short_name != original_name:
                detail_report = detail_report.replace(original_name, short_name)
        # 기존 상세 분석의 각 섹션을 개별 한 줄 expander로 분리
        detail_map = {}
        current_title = None
        current_lines = []
        title_map = {
            "■ MMS 상품 실적": "MMS 상품 실적",
            "■ MMS 발송 통계": "MMS 발송 통계",
            "■ 카테고리 분석": "카테고리 분석",
            "■ SEG 분석": "SEG 분석",
            "■ 요일·시간대 분석": "요일·시간대 분석",
            "■ 상품 인사이트": "상품별 상세 인사이트",
            "■ 최저가 미확보 상품": "최저가 미확보 상품",
        }
        for raw_line in detail_report.splitlines():
            line = raw_line.strip()
            if line in title_map:
                if current_title is not None:
                    detail_map[current_title] = "\n".join(current_lines).strip()
                current_title = title_map[line]
                current_lines = []
            elif current_title is not None and line:
                current_lines.append(line)
        if current_title is not None:
            detail_map[current_title] = "\n".join(current_lines).strip()

        for section_name in detail_sections:
            with st.expander(f"▶ {section_name}", expanded=False):
                section_text = detail_map.get(section_name, "해당 기간 데이터가 없습니다.")
                section_lines = [x.strip() for x in section_text.splitlines() if x.strip()]
                section_html = "".join(
                    f'<div style="margin:0 0 3px 0; line-height:1.55;">{x}</div>'
                    for x in section_lines
                )
                st.markdown(
                    f'<div class="insight-box">{section_html}</div>',
                    unsafe_allow_html=True,
                )

    with tabs[1]:
        total_amount = pw["주문금액"].sum()
        pw["주문비중"] = pw["주문금액"] / total_amount if total_amount else 0
        sort_cols = [c for c in ["_date", "시간대", "전시순서"] if c in pw.columns]
        product_sorted = pw.sort_values(sort_cols) if sort_cols else pw
        cols = [
            "일자", "요일", "시간대", "성별", "연령", "소재",
            "전시순서", "추가노출", "상품명", "멤버십혜택가",
            "주문건수", "주문수량", "주문금액", "주문비중"
        ]
        view = product_sorted[[c for c in cols if c in product_sorted.columns]].copy()
        total = {c: "" for c in view.columns}
        total[view.columns[0]] = "총합계"
        for c in ["주문건수", "주문수량", "주문금액"]:
            if c in view.columns:
                total[c] = product_sorted[c].sum()
        if "주문비중" in view.columns:
            total["주문비중"] = 1.0
        view = pd.concat([view, pd.DataFrame([total])], ignore_index=True)
        raw_amounts = (
            list(pd.to_numeric(view["주문금액"], errors="coerce"))
            if "주문금액" in view.columns else []
        )
        formatted_view = clean_identifier_columns(weekly_display_format(view))
        styled_view = formatted_view.style.apply(
            lambda _: style_weekly_product_rows(formatted_view, raw_amounts),
            axis=None,
        )
        st.dataframe(
            styled_view,
            use_container_width=True,
            hide_index=True,
            height=680,
        )

    with tabs[2]:
        material = pd.DataFrame({
            "연도": selected_year,
            "주차": sw["주차"],
            "일자": sw.get("일자", sw["_date"].dt.strftime("%m%d")),
            "요일": sw.get("요일", ""),
            "시간대": sw.get("시간대", ""),
            "소재": sw.get("소재", ""),
            "성별": sw.get("성별", ""),
            "연령": sw.get("연령", ""),
            "상품수": sw.get("상품수", 0),
            "발송성공건수": sw[send_col],
            "클릭수(uniq)": sw[click_col],
            "CTR(uniq)": safe_div(sw[click_col], sw[send_col]),
            "CVR(클릭>구매)": safe_div(sw["주문건수"], sw[click_col]),
            "객단가": safe_div(sw["주문금액"], sw["주문건수"]),
            "SPM": safe_div(sw["주문금액"], sw[send_col]),
            "주문건수": sw["주문건수"],
            "주문수량": sw["주문수량"],
            "주문금액": sw["주문금액"],
        }).fillna(0)
        st.dataframe(
            clean_identifier_columns(weekly_display_format(material)),
            use_container_width=True, hide_index=True, height=570
        )

    with tabs[3]:
        seg = grouped_send_table(sw, ["성별", "연령"])
        seg.insert(0, "주차", week)
        seg.insert(0, "연도", selected_year)
        cols = ["연도", "주차", "성별", "연령", "발송횟수", "CTR(uniq)", "CVR(클릭>구매)", "객단가", "SPM", "주문금액", "발송당매출(발송횟수)"]
        st.dataframe(
            clean_identifier_columns(weekly_display_format(seg[cols])),
            use_container_width=True, hide_index=True
        )

    with tabs[4]:
        weekday = grouped_send_table(sw, ["요일"])
        weekday.insert(0, "주차", week)
        weekday.insert(0, "연도", selected_year)
        cols = ["연도", "주차", "요일", "발송횟수", "CTR(uniq)", "CVR(클릭>구매)", "객단가", "SPM", "주문금액", "발송당매출(발송횟수)"]
        st.dataframe(
            clean_identifier_columns(weekly_display_format(weekday[cols])),
            use_container_width=True, hide_index=True
        )

    with tabs[5]:
        time_df = grouped_send_table(sw, ["시간대"])
        time_df.insert(0, "주차", week)
        time_df.insert(0, "연도", selected_year)
        cols = ["연도", "주차", "시간대", "발송횟수", "CTR(uniq)", "CVR(클릭>구매)", "객단가", "SPM", "주문금액", "발송당매출(발송횟수)"]
        st.dataframe(
            clean_identifier_columns(weekly_display_format(time_df[cols])),
            use_container_width=True, hide_index=True
        )

    with tabs[6]:
        rank = pw.groupby(["쇼라코드", "상품명"], as_index=False).agg(
            발송횟수=("상품명", "size"),
            주문건수=("주문건수", "sum"),
            주문수량=("주문수량", "sum"),
            주문금액=("주문금액", "sum"),
        ).sort_values("주문금액", ascending=False)
        rank.insert(0, "주차", week)
        rank.insert(0, "연도", selected_year)
        st.dataframe(
            clean_identifier_columns(weekly_display_format(rank)),
            use_container_width=True, hide_index=True, height=680
        )

    weekly_summary = [
        f"{selected_year}년 {week} 주문금액 {fmt_num(pw['주문금액'].sum())}원",
        f"발송횟수 {len(sw)}회 / 상품수 {len(pw)}건",
        f"핵심 상품 {sum(pw['주문금액'].apply(product_grade) == '핵심 상품')}개",
    ]
    ppt = build_ppt(
        f"{selected_year}년 {week} MMS 주간실적",
        weekly_summary,
        rank[[c for c in ["상품명", "발송횟수", "주문건수", "주문수량", "주문금액"] if c in rank.columns]],
    )
    st.download_button(
        "📥 주간실적 PPT 다운로드",
        ppt,
        file_name=f"{selected_year}_{week}_MMS_주간실적.pptx",
    )


# ─────────────────────────────────────────────────────────────────────────────
# 상품구분
# ─────────────────────────────────────────────────────────────────────────────
elif menu == "상품구분":
    st.markdown('<div class="section-title">상품구분</div>', unsafe_allow_html=True)
    st.caption("🔗 현재 브라우저 주소를 그대로 공유하면 상품구분 화면으로 바로 연결됩니다.")

    st.markdown('<div class="subsection-title">상품 등급 기준</div>', unsafe_allow_html=True)
    st.caption("평균 주문금액 기준")
    grade_rule_df = pd.DataFrame({
        "주문금액": [
            "100만원 미만",
            "100만원 이상 ~ 200만원 미만",
            "200만원 이상 ~ 300만원 미만",
            "300만원 이상 ~ 500만원 미만",
            "500만원 이상",
        ],
        "등급": ["🔴 부진 상품", "🟠 관찰 상품", "🟡 안정 상품", "🟢 우수 상품", "🔵 핵심 상품"],
    })
    st.dataframe(grade_rule_df, use_container_width=True, hide_index=True, height=212)

    filter_col1, filter_col2 = st.columns([1.5, 1])
    with filter_col1:
        date_range = st.date_input(
            "기간 선택",
            [products["_date"].min().date(), products["_date"].max().date()],
            key="product_group_date_range",
        )
    with filter_col2:
        grade_filter = st.multiselect(
            "상품등급 선택",
            GRADE_ORDER,
            default=GRADE_ORDER,
            key="product_group_grade_filter",
        )

    if len(date_range) == 2:
        start, end = date_range
        filt = products[
            (products["_date"].dt.date >= start)
            & (products["_date"].dt.date <= end)
        ].copy()
    else:
        filt = products.copy()

    search_col1, search_col2 = st.columns(2)
    with search_col1:
        product_number_search = st.text_input(
            "상품번호 검색",
            placeholder="쇼라코드 또는 알파코드",
            key="product_group_number_search",
        )
    with search_col2:
        product_name_search = st.text_input(
            "상품명 검색",
            placeholder="상품명 일부 입력",
            key="product_group_name_search",
        )

    group_keys = [c for c in ["쇼라코드", "알파코드", "상품명"] if c in filt.columns]
    grouped = filt.groupby(group_keys, dropna=False, as_index=False).agg(
        운영횟수=("상품명", "size"),
        최고실적=("주문금액", "max"),
        최저실적=("주문금액", "min"),
        평균실적=("주문금액", "mean"),
    )
    grouped["등급"] = grouped["평균실적"].apply(product_grade)

    kpi_cols = st.columns(6)
    kpi_values = [
        ("상품수", len(grouped)),
        ("🔵 핵심", int((grouped["등급"] == "핵심 상품").sum())),
        ("🟢 우수", int((grouped["등급"] == "우수 상품").sum())),
        ("🟡 안정", int((grouped["등급"] == "안정 상품").sum())),
        ("🟠 관찰", int((grouped["등급"] == "관찰 상품").sum())),
        ("🔴 부진", int((grouped["등급"] == "부진 상품").sum())),
    ]
    for col, (label, value) in zip(kpi_cols, kpi_values):
        with col:
            st.metric(label, f"{value:,}개")

    case_rows = []
    for _, group_row in grouped.iterrows():
        mask = pd.Series(True, index=filt.index)
        for key in group_keys:
            value = group_row.get(key)
            if pd.isna(value):
                mask &= filt[key].isna()
            else:
                mask &= filt[key].astype(str).eq(str(value))
        hist = filt[mask].sort_values("_date")
        cases = []
        for _, hist_row in hist.iterrows():
            cases.extend(classify_cases(hist_row, products))
        case_rows.append(", ".join(dict.fromkeys(cases)))
    grouped["사례"] = case_rows

    result = grouped[grouped["등급"].isin(grade_filter)].copy()

    if product_number_search.strip():
        number_query = product_number_search.strip()
        number_mask = pd.Series(False, index=result.index)
        for code_col in ["쇼라코드", "알파코드"]:
            if code_col in result.columns:
                number_mask |= result[code_col].astype(str).str.contains(
                    number_query, case=False, na=False, regex=False
                )
        result = result[number_mask]

    if product_name_search.strip():
        result = result[result["상품명"].astype(str).str.contains(
            product_name_search.strip(), case=False, na=False, regex=False
        )]

    display_cols = [
        c for c in [
            "쇼라코드", "알파코드", "상품명", "운영횟수",
            "최고실적", "최저실적", "평균실적", "등급", "사례",
        ] if c in result.columns
    ]
    display_df = result[display_cols].copy().reset_index(drop=True)

    for money_col in ["최고실적", "최저실적", "평균실적"]:
        if money_col in display_df.columns:
            display_df[money_col] = display_df[money_col].map(format_integer_price)

    st.caption(f"조회 상품 {len(display_df):,}개 · 행을 선택하면 아래에 발송 이력이 표시됩니다.")
    selection_event = st.dataframe(
        display_df,
        use_container_width=True,
        hide_index=True,
        height=500,
        on_select="rerun",
        selection_mode="single-row",
        key="product_group_summary_table",
    )

    selected_rows = getattr(getattr(selection_event, "selection", None), "rows", [])
    if selected_rows:
        selected_pos = selected_rows[0]
        selected_record = result.reset_index(drop=True).iloc[selected_pos]

        history_mask = pd.Series(True, index=products.index)
        for key in group_keys:
            value = selected_record.get(key)
            if pd.isna(value):
                history_mask &= products[key].isna()
            else:
                history_mask &= products[key].astype(str).eq(str(value))
        history = products[history_mask].sort_values("_date", ascending=False).copy()

        st.markdown(
            f'<div class="subsection-title">발송 이력 · {selected_record.get("상품명", "")}</div>',
            unsafe_allow_html=True,
        )

        history["발송일"] = history["_date"].dt.strftime("%Y-%m-%d")
        history["타겟"] = history.apply(target_label, axis=1)
        history["프로모션"] = history.apply(promotion_label, axis=1)

        history_cols = [
            c for c in [
                "발송일", "타겟", "캠페인명", "소재", "정상가", "멤버십혜택가",
                "할인율", "주문건수", "주문수량", "주문금액", "발송일 최저가", "프로모션",
            ] if c in history.columns
        ]
        history_view = history[history_cols].copy()

        for price_col in ["정상가", "멤버십혜택가", "주문금액", "발송일 최저가"]:
            if price_col in history_view.columns:
                history_view[price_col] = history_view[price_col].map(format_integer_price)
        if "할인율" in history_view.columns:
            history_view["할인율"] = history_view["할인율"].map(
                lambda x: f"{float(x) * 100:.0f}%" if float(x) <= 1 else f"{float(x):.0f}%"
            )

        st.dataframe(
            history_view,
            use_container_width=True,
            hide_index=True,
            height=min(420, 42 + len(history_view) * 35),
        )


elif menu == "타겟분석":
    st.caption("🔗 현재 브라우저 주소를 그대로 공유하면 타겟분석 화면으로 바로 연결됩니다.")
    st.markdown('<div class="section-title">타겟분석</div>', unsafe_allow_html=True)

    target_date_range = st.date_input(
        "기간 선택",
        [sends["_date"].min().date(), sends["_date"].max().date()],
        key="target_analysis_date_range",
    )

    if len(target_date_range) == 2:
        target_start, target_end = target_date_range
        target_sends = sends[
            (sends["_date"].dt.date >= target_start)
            & (sends["_date"].dt.date <= target_end)
        ].copy()
        target_products = products[
            (products["_date"].dt.date >= target_start)
            & (products["_date"].dt.date <= target_end)
        ].copy()
    else:
        target_sends = sends.copy()
        target_products = products.copy()

    if target_sends.empty:
        st.info("선택한 기간에 해당하는 발송 데이터가 없습니다.")
    else:
        def target_analysis_raw(data: pd.DataFrame, group_keys: list[str]) -> pd.DataFrame:
            view = grouped_send_table(data, group_keys).copy()
            view = view.rename(columns={
                "CTR(uniq)": "CTR",
                "CVR(클릭>구매)": "CVR",
                "발송당매출(발송횟수)": "발송당매출",
            })
            if "SEG" in view.columns:
                view["SEG"] = view["SEG"].map(clean_identifier_value)
            if "객단가" not in view.columns:
                view["객단가"] = view.apply(
                    lambda r: float(r.get("주문금액", 0)) / float(r.get("주문건수", 0))
                    if float(r.get("주문건수", 0)) else 0,
                    axis=1,
                )
            view = view.sort_values("SPM", ascending=False).reset_index(drop=True)
            view.insert(0, "순위", range(1, len(view) + 1))
            return view

        def format_target_view(view: pd.DataFrame, group_keys: list[str]) -> pd.DataFrame:
            out = view.copy()
            column_order = ["순위"] + group_keys + [
                "발송횟수", "발송건수", "클릭수", "주문건수", "주문금액",
                "CTR", "CVR", "SPM", "객단가", "발송당매출",
            ]
            out = out[[c for c in column_order if c in out.columns]].copy()
            for col in ["발송횟수", "발송건수", "클릭수", "주문건수", "주문금액", "객단가", "발송당매출"]:
                if col in out.columns:
                    out[col] = out[col].map(fmt_num)
            for col in ["CTR", "CVR"]:
                if col in out.columns:
                    out[col] = out[col].map(fmt_pct)
            if "SPM" in out.columns:
                out["SPM"] = out["SPM"].map(lambda x: f"{float(x):.1f}")
            return out

        gender_age_raw = target_analysis_raw(target_sends, ["성별", "연령"])
        gender_age_seg_raw = target_analysis_raw(target_sends, ["성별", "연령", "SEG"])

        st.markdown('<div class="subsection-title">성별·연령별 SPM</div>', unsafe_allow_html=True)
        chart_data = gender_age_raw.copy()
        chart_data["타겟"] = chart_data.apply(
            lambda r: f"{str(r.get('성별', '')).strip()} {clean_identifier_value(r.get('연령', ''))}".strip(),
            axis=1,
        )
        fig_target_spm = go.Figure(
            go.Bar(
                x=chart_data["타겟"],
                y=chart_data["SPM"],
                text=chart_data["SPM"].map(lambda x: f"{float(x):.1f}"),
                textposition="outside",
                cliponaxis=False,
            )
        )
        fig_target_spm.update_layout(
            height=420,
            margin=dict(l=45, r=25, t=35, b=60),
            xaxis_title="",
            yaxis_title="SPM",
            showlegend=False,
            plot_bgcolor="#ffffff",
            yaxis=dict(gridcolor="#e5e7eb"),
        )
        st.plotly_chart(fig_target_spm, use_container_width=True, key="target_spm_chart")

        st.markdown('<div class="subsection-title">성별·연령별 실적</div>', unsafe_allow_html=True)
        gender_age_event = st.dataframe(
            format_target_view(gender_age_raw, ["성별", "연령"]),
            use_container_width=True,
            hide_index=True,
            height=min(520, 42 + len(gender_age_raw) * 35),
            on_select="rerun",
            selection_mode="single-row",
            key="target_gender_age_table",
        )

        st.markdown('<div class="subsection-title">성별·연령·SEG별 실적</div>', unsafe_allow_html=True)
        gender_age_seg_event = st.dataframe(
            format_target_view(gender_age_seg_raw, ["성별", "연령", "SEG"]),
            use_container_width=True,
            hide_index=True,
            height=min(620, 42 + len(gender_age_seg_raw) * 35),
            on_select="rerun",
            selection_mode="single-row",
            key="target_gender_age_seg_table",
        )

        selected_target = None
        selected_seg = None
        seg_rows = list(getattr(getattr(gender_age_seg_event, "selection", None), "rows", []) or [])
        age_rows = list(getattr(getattr(gender_age_event, "selection", None), "rows", []) or [])
        if seg_rows:
            selected_target = gender_age_seg_raw.iloc[int(seg_rows[0])]
            selected_seg = clean_identifier_value(selected_target.get("SEG", ""))
        elif age_rows:
            selected_target = gender_age_raw.iloc[int(age_rows[0])]

        if selected_target is not None:
            selected_gender = str(selected_target.get("성별", "")).strip()
            selected_age = clean_identifier_value(selected_target.get("연령", ""))
            history = target_products.copy()
            if "성별" in history.columns:
                history = history[history["성별"].astype(str).str.strip().eq(selected_gender)]
            if "연령" in history.columns:
                history = history[history["연령"].map(clean_identifier_value).eq(selected_age)]
            if selected_seg and "SEG" in history.columns:
                history = history[history["SEG"].map(clean_identifier_value).eq(selected_seg)]

            if history.empty:
                st.info("선택한 타겟의 상품 발송 이력이 없습니다.")
            else:
                send_count_col = first_col(target_sends, ["발송 성공 건수", "총 발송 건수"])
                campaign_col_send = first_col(target_sends, ["캠페인명", "캠페인", "소재"])
                campaign_col_product = first_col(history, ["캠페인명", "캠페인", "소재"])

                send_lookup = target_sends.copy()
                send_lookup["_date_key"] = send_lookup["_date"].dt.strftime("%Y-%m-%d")
                if send_count_col:
                    send_lookup["_send_count"] = pd.to_numeric(send_lookup[send_count_col], errors="coerce").fillna(0)
                else:
                    send_lookup["_send_count"] = 0
                if campaign_col_send:
                    send_lookup["_campaign_key"] = send_lookup[campaign_col_send].fillna("").astype(str).str.strip()
                    campaign_map = send_lookup.groupby(["_date_key", "_campaign_key"])["_send_count"].sum().to_dict()
                else:
                    campaign_map = {}
                date_map = send_lookup.groupby("_date_key")["_send_count"].sum().to_dict()

                history["발송일"] = history["_date"].dt.strftime("%Y-%m-%d")
                history["SEG"] = history.get("SEG", "").map(clean_identifier_value) if "SEG" in history.columns else ""
                history["행사가"] = pd.to_numeric(history.get("멤버십혜택가", 0), errors="coerce").fillna(0)
                history["_campaign_key"] = history[campaign_col_product].fillna("").astype(str).str.strip() if campaign_col_product else ""
                history["_send_count"] = history.apply(
                    lambda r: campaign_map.get((r["발송일"], r["_campaign_key"]), date_map.get(r["발송일"], 0)),
                    axis=1,
                )
                history["SPM"] = history.apply(
                    lambda r: float(r.get("주문금액", 0)) / float(r.get("_send_count", 0))
                    if float(r.get("_send_count", 0)) else 0,
                    axis=1,
                )
                history["할인율"] = history.apply(
                    lambda r: floor_discount_rate(float(r.get("정상가", 0) or 0), float(r.get("행사가", 0) or 0)),
                    axis=1,
                )

                target_name = f"{selected_gender} {selected_age}" + (f" SEG{selected_seg}" if selected_seg else "")
                st.markdown(f'<div class="subsection-title">{target_name} 발송 이력</div>', unsafe_allow_html=True)
                history_cols = [
                    "발송일", "SEG", "쇼라코드", "알파코드", "상품명",
                    "정상가", "행사가", "할인율", "주문금액", "SPM",
                ]
                history_view = history[[c for c in history_cols if c in history.columns]].sort_values("발송일", ascending=False).copy()
                for col in ["정상가", "행사가", "주문금액"]:
                    if col in history_view.columns:
                        history_view[col] = history_view[col].map(format_integer_price)
                if "할인율" in history_view.columns:
                    history_view["할인율"] = history_view["할인율"].map(lambda x: f"{float(x) * 100:.0f}%")
                if "SPM" in history_view.columns:
                    history_view["SPM"] = history_view["SPM"].map(lambda x: f"{float(x):.1f}")
                st.dataframe(
                    history_view,
                    use_container_width=True,
                    hide_index=True,
                    height=min(620, 42 + len(history_view) * 35),
                )

elif menu == "편성 프로그램":
    st.caption("🔗 현재 브라우저 주소를 그대로 공유하면 편성 프로그램 화면으로 바로 연결됩니다.")
    st.markdown('<div class="section-title">🗓️ 편성 프로그램</div>', unsafe_allow_html=True)
    st.caption(
        "발송 슬롯과 소재를 각각 입력한 뒤, 주력 상품 안에서 과거 주문금액을 우선으로 자동 편성합니다. "
        "같은 날짜의 동일 상품 중복 편성은 자동으로 제외됩니다."
    )

    tab_input, tab_result, tab_history = st.tabs([
        "① 편성 조건 입력",
        "② 자동 편성 결과",
        "③ 발송 이력 구분",
    ])

    time_options = [f"{hour:02d}:{minute:02d}" for hour in range(24) for minute in (0, 30)]
    target_options = [
        f"{age} {gender} SEG{seg}"
        for age in ["3040", "5060"]
        for gender in ["남성", "여성"]
        for seg in [1, 2, 3]
    ]

    with tab_input:
        st.markdown('<div class="subsection-title">발송 슬롯 입력</div>', unsafe_allow_html=True)
        st.caption("발송일은 달력에서 선택하고, 시간대는 00:00~23:30 범위에서 30분 단위로 설정하세요.")

        today_value = pd.Timestamp.today().normalize()
        default_slots = pd.DataFrame([
            {"발송일": today_value, "시간대": "11:30", "타겟": "5060 여성 SEG1", "상품수": 4},
            {"발송일": today_value, "시간대": "16:00", "타겟": "3040 남성 SEG1", "상품수": 4},
        ])
        if "schedule_slots" not in st.session_state:
            st.session_state.schedule_slots = default_slots

        edited_slots = st.data_editor(
            st.session_state.schedule_slots,
            num_rows="dynamic",
            use_container_width=True,
            hide_index=True,
            column_config={
                "발송일": st.column_config.DateColumn("발송일", format="YYYY-MM-DD", required=True),
                "시간대": st.column_config.SelectboxColumn("시간대", options=time_options, required=True),
                "타겟": st.column_config.SelectboxColumn("타겟", options=target_options, required=True),
                "상품수": st.column_config.NumberColumn("상품수", min_value=1, max_value=20, step=1, required=True),
            },
            key="schedule_slots_editor_v2",
        )
        st.session_state.schedule_slots = edited_slots

        st.markdown('<div class="subsection-title">주력 상품 입력</div>', unsafe_allow_html=True)
        st.caption("알파코드·쇼라코드·상품명·정상가·행사가를 입력하세요. 할인율은 자동 편성 결과에서 자동 계산됩니다.")

        upload_col, paste_col = st.columns([1, 1])
        uploaded_candidates = None
        with upload_col:
            schedule_file = st.file_uploader(
                "주력 상품 엑셀/CSV 업로드",
                type=["xlsx", "xls", "csv"],
                key="schedule_candidate_upload",
            )
            if schedule_file is not None:
                try:
                    if schedule_file.name.lower().endswith(".csv"):
                        uploaded_candidates = pd.read_csv(schedule_file)
                    else:
                        uploaded_candidates = pd.read_excel(schedule_file)
                except Exception as exc:
                    st.error(f"상품 파일을 읽지 못했습니다: {exc}")
        with paste_col:
            st.caption("파일이 없으면 아래 표에 직접 붙여넣어도 됩니다.")

        required_candidate_cols = ["알파코드", "쇼라코드", "상품명", "정상가", "행사가"]
        if uploaded_candidates is not None:
            for col in required_candidate_cols:
                if col not in uploaded_candidates.columns:
                    uploaded_candidates[col] = ""
            st.session_state.schedule_candidates = uploaded_candidates[required_candidate_cols].copy()

        if "schedule_candidates" not in st.session_state:
            st.session_state.schedule_candidates = pd.DataFrame(
                [{c: "" for c in required_candidate_cols} for _ in range(8)]
            )

        edited_candidates = st.data_editor(
            st.session_state.schedule_candidates,
            num_rows="dynamic",
            use_container_width=True,
            hide_index=True,
            column_config={
                "알파코드": st.column_config.TextColumn("알파코드"),
                "쇼라코드": st.column_config.TextColumn("쇼라코드"),
                "상품명": st.column_config.TextColumn("상품명", width="large"),
                "정상가": st.column_config.NumberColumn("정상가", min_value=0, step=100),
                "행사가": st.column_config.NumberColumn("행사가", min_value=0, step=100),
            },
            key="schedule_candidates_editor",
        )
        st.session_state.schedule_candidates = edited_candidates

        candidates_calc = edited_candidates.copy()
        for col in ["정상가", "행사가"]:
            candidates_calc[col] = num(candidates_calc[col])
        candidates_calc["할인율계산값"] = candidates_calc.apply(
            lambda r: floor_discount_rate(r["정상가"], r["행사가"]), axis=1
        )
        candidates_calc["할인율"] = candidates_calc["할인율계산값"].map(
            lambda x: f"{float(x):.0%}" if pd.notna(x) else "-"
        )
        st.session_state.schedule_candidates_calc = candidates_calc

        c1, c2, c3 = st.columns(3)
        with c1:
            cooldown_days = st.number_input(
                "발송 후 재편성 제한일",
                min_value=0,
                max_value=60,
                value=int(st.session_state.get("schedule_cooldown", 0)),
                step=1,
            )
        with c2:
            max_weekly_count = st.number_input(
                "상품별 주간 최대 편성횟수",
                min_value=1,
                max_value=10,
                value=int(st.session_state.get("schedule_max_weekly", 3)),
                step=1,
            )
        with c3:
            st.metric("편성 슬롯 수", f"{len(edited_slots.dropna(how='all'))}회")

        st.session_state.schedule_cooldown = cooldown_days
        st.session_state.schedule_max_weekly = max_weekly_count

        if st.button("🤖 매출 우선 자동 편성 실행", type="primary", use_container_width=True):
            clean_slots = edited_slots.copy().reset_index(drop=True)
            clean_slots["발송일"] = pd.to_datetime(clean_slots["발송일"], errors="coerce")
            clean_slots = clean_slots[
                clean_slots["발송일"].notna()
                & clean_slots["시간대"].astype(str).str.strip().ne("")
                & clean_slots["타겟"].astype(str).str.strip().ne("")
            ].copy()

            clean_candidates = candidates_calc.copy()
            clean_candidates = clean_candidates[
                clean_candidates["상품명"].astype(str).str.strip().ne("")
                & clean_candidates["정상가"].gt(0)
                & clean_candidates["행사가"].gt(0)
            ].copy()

            if clean_slots.empty:
                st.error("발송일·시간대·타겟이 입력된 슬롯을 1개 이상 등록해주세요.")
            elif clean_candidates.empty:
                st.error("주력 상품을 1개 이상 입력해주세요.")
            else:
                result, detail_map = build_schedule_recommendations(
                    clean_slots,
                    clean_candidates,
                    products,
                    int(cooldown_days),
                    int(max_weekly_count),
                )
                st.session_state.schedule_result = result
                st.session_state.schedule_detail_map = detail_map
                st.success("자동 편성이 완료되었습니다. '② 자동 편성 결과' 탭에서 확인해주세요.")

    with tab_result:
        result = st.session_state.get("schedule_result", pd.DataFrame())
        detail_map = st.session_state.get("schedule_detail_map", {})

        if result.empty:
            st.info("먼저 '① 편성 조건 입력'에서 자동 편성을 실행해주세요.")
        else:
            st.markdown('<div class="subsection-title">전체 편성안</div>', unsafe_allow_html=True)
            result_cols = [
                "발송일", "시간대", "타겟", "전시순서",
                "알파코드", "쇼라코드", "상품명", "정상가", "행사가", "할인율",
            ]
            copy_view = result[[c for c in result_cols if c in result.columns]].copy()
            for col in ["정상가", "행사가"]:
                if col in copy_view.columns:
                    copy_view[col] = copy_view[col].map(format_integer_price)
            st.dataframe(copy_view, use_container_width=True, hide_index=True)

            csv_bytes = copy_view.to_csv(index=False, encoding="utf-8-sig").encode("utf-8-sig")
            excel_buffer = io.BytesIO()
            with pd.ExcelWriter(excel_buffer, engine="openpyxl") as writer:
                copy_view.to_excel(writer, index=False, sheet_name="편성안")
            d1, d2 = st.columns(2)
            with d1:
                st.download_button(
                    "📥 편성안 CSV 다운로드", data=csv_bytes,
                    file_name="MMS_자동편성안.csv", mime="text/csv", use_container_width=True,
                )
            with d2:
                st.download_button(
                    "📥 편성안 엑셀 다운로드", data=excel_buffer.getvalue(),
                    file_name="MMS_자동편성안.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True,
                )

            st.markdown('<div class="subsection-title">슬롯별 추천 결과 및 근거</div>', unsafe_allow_html=True)
            group_cols = [c for c in ["발송일", "시간대", "타겟"] if c in result.columns]
            for group_key, group in result.groupby(group_cols, sort=False):
                if not isinstance(group_key, tuple):
                    group_key = (group_key,)
                label = " · ".join(str(x) for x in group_key)
                st.markdown(f"### {label}")
                main_view = group[["전시순서", "알파코드", "쇼라코드", "상품명", "정상가", "행사가", "할인율"]].copy()
                for col in ["정상가", "행사가"]:
                    main_view[col] = main_view[col].map(format_integer_price)
                st.dataframe(main_view, use_container_width=True, hide_index=True)

                for _, row in group.iterrows():
                    detail_key = (
                        str(row["발송일"]), str(row["시간대"]), str(row["타겟"]),
                        str(row["알파코드"]), str(row["쇼라코드"]), str(row["상품명"]),
                    )
                    metrics = detail_map.get(detail_key, {})
                    with st.expander(f"{row['상품명']} · 추천 근거 및 발송 이력"):
                        if not metrics:
                            st.info("상세 이력이 없습니다.")
                            continue
                        if metrics.get("이력여부") == "신규":
                            st.warning("과거 동일 상품 이력이 없는 신규 TEST 후보입니다.")
                            st.write(metrics.get("근거", ""))
                            continue
                        m1, m2, m3, m4 = st.columns(4)
                        m1.metric("예상 기준매출", compact_money(metrics.get("추천매출", 0)))
                        m2.metric("역대 최고매출", compact_money(metrics.get("최고매출", 0)))
                        m3.metric("최고타겟", metrics.get("최고타겟", "-"))
                        recent = metrics.get("최근발송일")
                        m4.metric("최근 발송일", pd.Timestamp(recent).strftime("%Y-%m-%d") if pd.notna(recent) else "-")
                        st.markdown("**편성 근거**")
                        st.write(f"- {metrics.get('근거', '')}\n- 과거 주문금액 우선 정렬\n- 동일 날짜 중복 및 주간 최대횟수 조건 반영")
                        st.markdown("**타겟별 성과**")
                        target_summary = schedule_target_summary(metrics.get("이력", pd.DataFrame()))
                        if not target_summary.empty:
                            target_display = target_summary.copy()
                            for c in ["평균매출", "최고매출"]:
                                target_display[c] = target_display[c].map(format_integer_price)
                            st.dataframe(target_display, use_container_width=True, hide_index=True)
                        st.markdown("**발송 이력**")
                        history_view = schedule_history_table(metrics.get("이력", pd.DataFrame()), float(row["행사가"]))
                        if not history_view.empty:
                            for c in ["멤버십혜택가", "현재가 대비", "주문금액"]:
                                if c in history_view.columns:
                                    history_view[c] = history_view[c].map(format_integer_price)
                            st.dataframe(history_view, use_container_width=True, hide_index=True)
                st.divider()

    with tab_history:
        st.markdown('<div class="subsection-title">입력 상품 발송 이력 구분</div>', unsafe_allow_html=True)
        candidates_calc = st.session_state.get("schedule_candidates_calc", pd.DataFrame())
        if candidates_calc.empty:
            st.info("먼저 '① 편성 조건 입력'에서 주력 상품을 입력해주세요.")
        else:
            valid_candidates = candidates_calc[
                candidates_calc["상품명"].astype(str).str.strip().ne("")
            ].copy()
            history_rows = []
            for _, candidate in valid_candidates.iterrows():
                hist = match_candidate_history(candidate, products)
                history_rows.append({
                    "알파코드": clean_identifier_value(candidate.get("알파코드", "")),
                    "쇼라코드": clean_identifier_value(candidate.get("쇼라코드", "")),
                    "상품명": str(candidate.get("상품명", "")).strip(),
                    "정상가": float(candidate.get("정상가", 0) or 0),
                    "행사가": float(candidate.get("행사가", 0) or 0),
                    "할인율": candidate.get("할인율", "-"),
                    "발송이력": "있음" if not hist.empty else "없음",
                    "운영횟수": int(len(hist)),
                    "최근발송일": hist["_date"].max().strftime("%Y-%m-%d") if not hist.empty else "-",
                    "평균주문금액": float(hist["주문금액"].mean()) if not hist.empty else 0,
                })
            history_columns = [
                "알파코드", "쇼라코드", "상품명", "정상가", "행사가", "할인율",
                "발송이력", "운영횟수", "최근발송일", "평균주문금액",
            ]
            history_status = pd.DataFrame(history_rows, columns=history_columns)
            hist_tab, new_tab = st.tabs(["발송 이력 있는 상품", "발송 이력 없는 상품"])
            for target_tab, status in [(hist_tab, "있음"), (new_tab, "없음")]:
                with target_tab:
                    if history_status.empty:
                        view = pd.DataFrame(columns=history_columns)
                    else:
                        view = history_status[history_status["발송이력"].fillna("").eq(status)].copy()
                    if view.empty:
                        st.info(f"발송 이력 {status} 상품이 없습니다.")
                    else:
                        for c in ["정상가", "행사가", "평균주문금액"]:
                            view[c] = view[c].map(format_integer_price)
                        st.dataframe(view, use_container_width=True, hide_index=True)
